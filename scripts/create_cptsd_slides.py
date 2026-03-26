from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# === カラーパレット ===
NAVY      = RGBColor(0x1A, 0x2F, 0x5E)   # 濃紺（タイトル背景）
TEAL      = RGBColor(0x00, 0x7A, 0x87)   # ティール（アクセント）
ACCENT    = RGBColor(0xE8, 0x4A, 0x5F)   # 赤（強調）
LIGHT_BG  = RGBColor(0xF4, 0xF7, 0xFC)   # 薄青（コンテンツ背景）
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x1C, 0x1C, 0x2E)
MID_GRAY  = RGBColor(0x6B, 0x7A, 0x99)
YELLOW    = RGBColor(0xFF, 0xC1, 0x07)
GREEN     = RGBColor(0x28, 0xA7, 0x45)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── ユーティリティ ────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_color, border_color=None, border_pt=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=18, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_para(tf, text, font_size=16, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, space_before=0, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def add_textbox_multi(slide, lines, x, y, w, h, base_size=16):
    """lines: list of (text, size, bold, color, align)"""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for (text, size, bold, color, align) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txBox

def header_bar(slide, title, subtitle=""):
    """濃紺のヘッダーバー"""
    add_rect(slide, 0, 0, W, Inches(1.35), NAVY)
    add_text(slide, title, Inches(0.4), Inches(0.1), W - Inches(0.8), Inches(0.75),
             font_size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.82), W - Inches(0.8), Inches(0.45),
                 font_size=14, color=RGBColor(0xA8, 0xC4, 0xE8), align=PP_ALIGN.LEFT)
    # ティールのアクセントライン
    add_rect(slide, 0, Inches(1.35), W, Pt(4), TEAL)

def content_bg(slide):
    add_rect(slide, 0, 0, W, H, LIGHT_BG)

def bottom_bar(slide, note=""):
    add_rect(slide, 0, H - Inches(0.4), W, Inches(0.4), NAVY)
    if note:
        add_text(slide, note, Inches(0.3), H - Inches(0.38), W - Inches(0.6), Inches(0.35),
                 font_size=9, color=RGBColor(0xA8, 0xC4, 0xE8), align=PP_ALIGN.LEFT)

def bullet_box(slide, items, x, y, w, h, title="", bg=WHITE,
               bullet="●", font_size=15, title_size=14):
    shape = add_rect(slide, x, y, w, h, bg, TEAL, 1.5)
    if title:
        add_text(slide, title, x + Inches(0.15), y + Inches(0.1),
                 w - Inches(0.3), Inches(0.35),
                 font_size=title_size, bold=True, color=TEAL)
    txBox = slide.shapes.add_textbox(
        x + Inches(0.15), y + (Inches(0.45) if title else Inches(0.15)),
        w - Inches(0.3), h - (Inches(0.6) if title else Inches(0.3)))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = DARK_TEXT

def pmid_tag(slide, text, x, y):
    add_rect(slide, x, y, Inches(3.8), Inches(0.32), TEAL)
    add_text(slide, text, x + Inches(0.1), y + Inches(0.03),
             Inches(3.6), Inches(0.28), font_size=10, color=WHITE, bold=True)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド0：タイトルスライド
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, NAVY)
# 装飾ライン
add_rect(slide, 0, Inches(2.8), W, Pt(5), TEAL)
add_rect(slide, 0, Inches(2.85), W, Pt(2), YELLOW)
# タイトル
add_text(slide, "複雑性PTSD（C-PTSD）の診断と治療", Inches(0.8), Inches(1.2),
         W - Inches(1.6), Inches(1.1), font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "内科・総合診療医が知っておくべき実践知識", Inches(0.8), Inches(2.3),
         W - Inches(1.6), Inches(0.6), font_size=20, color=RGBColor(0xA8, 0xC4, 0xE8),
         align=PP_ALIGN.CENTER)
# サブ情報
add_text(slide, "30分レクチャー  ｜  対象：研修医〜若手医師", Inches(0.8), Inches(3.2),
         W - Inches(1.6), Inches(0.5), font_size=16, color=RGBColor(0x80, 0xA8, 0xD0),
         align=PP_ALIGN.CENTER)
add_text(slide, "おもろまちメディカルセンター  副院長・内科/呼吸器内科", Inches(0.8), Inches(3.75),
         W - Inches(1.6), Inches(0.4), font_size=13, color=RGBColor(0x80, 0xA8, 0xD0),
         align=PP_ALIGN.CENTER)
add_text(slide, "2026-03-20", Inches(0.8), Inches(4.2),
         W - Inches(1.6), Inches(0.4), font_size=13,
         color=RGBColor(0x80, 0xA8, 0xD0), align=PP_ALIGN.CENTER)
# アイコン的な装飾テキスト
add_text(slide, "EMDR  ｜  DBT-PTSD  ｜  TF-CBT  ｜  IFS  ｜  薬物療法",
         Inches(0.8), Inches(5.5), W - Inches(1.6), Inches(0.5),
         font_size=13, color=TEAL, align=PP_ALIGN.CENTER, italic=True)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド1：症例提示
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_bg(slide)
header_bar(slide, "症例提示", "オープニング｜1.5分")
bottom_bar(slide, "Slide 1 / 17")

# 症例ボックス
add_rect(slide, Inches(0.35), Inches(1.55), Inches(7.8), Inches(2.9), WHITE, NAVY, 2)
add_rect(slide, Inches(0.35), Inches(1.55), Inches(7.8), Inches(0.42), NAVY)
add_text(slide, "【症例】40代女性", Inches(0.5), Inches(1.57), Inches(5), Inches(0.38),
         font_size=16, bold=True, color=WHITE)
items_case = [
    "主訴：慢性腰痛・不眠・「原因不明の疲労」",
    "整形外科 → 消化器内科 → 心療内科 と転々",
    "現在 6つ目の外来受診",
    "カルテ記載：「感情的になりやすい」「信頼関係構築困難」「アドヒアランス不良」",
]
txBox = slide.shapes.add_textbox(Inches(0.55), Inches(2.05), Inches(7.4), Inches(2.3))
tf = txBox.text_frame; tf.word_wrap = True
first = True
for it in items_case:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_before = Pt(4)
    run = p.add_run(); run.text = f"▶  {it}"
    run.font.size = Pt(15); run.font.color.rgb = DARK_TEXT

# 右側：問いかけ
add_rect(slide, Inches(8.4), Inches(1.55), Inches(4.55), Inches(2.9), WHITE, TEAL, 2)
add_text(slide, "この患者が抱えているのは\n何でしょうか？", Inches(8.55), Inches(1.75),
         Inches(4.2), Inches(1.1), font_size=17, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
add_text(slide, "うつ病？\n慢性疼痛症？\nパーソナリティ障害？",
         Inches(8.55), Inches(2.85), Inches(4.2), Inches(1.1),
         font_size=15, color=MID_GRAY, align=PP_ALIGN.CENTER)

# 答えボックス
add_rect(slide, Inches(0.35), Inches(4.65), W - Inches(0.7), Inches(0.95), ACCENT)
add_text(slide, "実は DV歴10年以上。これらは複雑性PTSD（C-PTSD）のサインです。",
         Inches(0.55), Inches(4.72), W - Inches(1.1), Inches(0.8),
         font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド2：C-PTSDとは何か
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_bg(slide)
header_bar(slide, "C-PTSDとは何か", "オープニング｜1.5分")
bottom_bar(slide, "Slide 2 / 17")

# 左：定義
bullet_box(slide,
    ["2018年 ICD-11 に独立疾患単位として正式収載",
     "DSM-5 にはまだ独立した診断名なし",
     "原因：反復・長期・逃れにくいトラウマ",
     "　幼少期虐待 / DV / 人身売買 / 拷問 / 難民経験"],
    Inches(0.35), Inches(1.55), Inches(5.8), Inches(2.5),
    title="定義・背景", bg=WHITE)

# 右：通常PTSDとの違い矢印
add_rect(slide, Inches(6.4), Inches(1.55), Inches(6.55), Inches(2.5), WHITE, TEAL, 1.5)
add_text(slide, "通常 PTSD", Inches(6.55), Inches(1.65), Inches(2.8), Inches(0.4),
         font_size=14, bold=True, color=TEAL)
add_text(slide, "再体験\n回避\n過覚醒",
         Inches(6.55), Inches(2.05), Inches(2.8), Inches(0.9),
         font_size=13, color=DARK_TEXT)
add_text(slide, "＋", Inches(9.3), Inches(2.3), Inches(0.6), Inches(0.5),
         font_size=32, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_rect(slide, Inches(9.85), Inches(1.65), Inches(2.85), Inches(2.25), TEAL)
add_text(slide, "DSO（自己組織化の障害）", Inches(9.9), Inches(1.7), Inches(2.7), Inches(0.45),
         font_size=13, bold=True, color=WHITE)
add_text(slide, "① 情動調節困難\n② 否定的自己概念\n③ 対人関係障害",
         Inches(9.9), Inches(2.15), Inches(2.7), Inches(1.0),
         font_size=13, bold=True, color=YELLOW)

# 下：一言
add_rect(slide, Inches(0.35), Inches(4.25), W - Inches(0.7), Inches(0.7), NAVY)
add_text(slide, "「通常PTSDとは症状の構造が異なる」——この理解が治療選択のすべての出発点",
         Inches(0.55), Inches(4.3), W - Inches(1.1), Inches(0.6),
         font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド3：ICD-11 DSO
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_bg(slide)
header_bar(slide, "ICD-11診断基準：DSO（自己組織化の障害）", "第1章：C-PTSDを理解する｜2分")
bottom_bar(slide, "Slide 3 / 17")

# 左列：通常PTSD
add_rect(slide, Inches(0.35), Inches(1.55), Inches(4.0), Inches(3.6), WHITE, MID_GRAY, 1.5)
add_text(slide, "通常 PTSD（共通症状）", Inches(0.5), Inches(1.65),
         Inches(3.7), Inches(0.4), font_size=14, bold=True, color=MID_GRAY)
for i, t in enumerate(["再体験（フラッシュバック・悪夢）","回避","過覚醒（過度の警戒・易驚愕）"]):
    add_text(slide, f"✓  {t}", Inches(0.5), Inches(2.15 + i*0.7),
             Inches(3.7), Inches(0.6), font_size=14, color=DARK_TEXT)

# 中央：プラス
add_text(slide, "＋", Inches(4.5), Inches(2.8), Inches(0.8), Inches(0.8),
         font_size=40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# 右列：DSO 3項目
dso_titles = ["情動調節困難", "否定的自己概念", "対人関係障害"]
dso_descs = [
    "感情の爆発 ↔ 感情の麻痺\n激しい感情の振れ幅",
    "「自分はダメだ」「穢れている」\n深い自責・羞恥・無力感",
    "人を信じることが極めて困難\n親密さへの恐怖・孤立"
]
dso_colors = [ACCENT, TEAL, NAVY]
for i in range(3):
    add_rect(slide, Inches(5.5), Inches(1.55 + i*1.2), Inches(7.45), Inches(1.1), dso_colors[i])
    add_text(slide, f"DSO ③{i+1}  {dso_titles[i]}", Inches(5.65), Inches(1.6 + i*1.2),
             Inches(5), Inches(0.4), font_size=16, bold=True, color=WHITE)
    add_text(slide, dso_descs[i], Inches(5.65), Inches(2.0 + i*1.2),
             Inches(7.1), Inches(0.55), font_size=13, color=RGBColor(0xFF,0xFF,0xCC))

# 症例への適用
add_rect(slide, Inches(0.35), Inches(5.2), W - Inches(0.7), Inches(0.9), LIGHT_BG, TEAL, 1.5)
add_text(slide, "症例に当てはめると：「感情的になりやすい」＝情動調節困難 ／ 「信頼構築困難」＝対人関係障害 ／ 「アドヒアランス不良」＝否定的自己概念からの無力感",
         Inches(0.55), Inches(5.27), W - Inches(1.1), Inches(0.75),
         font_size=13, color=DARK_TEXT, italic=True)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド4：通常PTSDとの臨床的な違い
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_bg(slide)
header_bar(slide, "通常PTSDとC-PTSDの臨床的な違い", "第1章：C-PTSDを理解する｜2分")
bottom_bar(slide, "Slide 4 / 17")

# テーブルヘッダー
headers = ["比較項目", "通常 PTSD", "複雑性 PTSD"]
col_x = [Inches(0.35), Inches(3.6), Inches(8.5)]
col_w = [Inches(3.1), Inches(4.7), Inches(4.4)]
add_rect(slide, Inches(0.35), Inches(1.55), W - Inches(0.7), Inches(0.45), NAVY)
for j, h in enumerate(headers):
    add_text(slide, h, col_x[j] + Inches(0.1), Inches(1.57),
             col_w[j], Inches(0.4), font_size=14, bold=True, color=WHITE)

rows = [
    ("トラウマの性質", "一回性・限定的", "反復・長期・逃れられない"),
    ("自己感覚", "比較的保たれる", "根本から傷ついている"),
    ("対人関係", "特定のトリガーで悪化", "全般的に障害されている"),
    ("解離症状", "軽度〜中等度", "高頻度・重篤"),
    ("治療反応", "標準PTSD治療で効く", "フェーズ準備なしには効きにくい"),
]
row_colors = [LIGHT_BG, WHITE, LIGHT_BG, WHITE, ACCENT]
text_colors = [DARK_TEXT, DARK_TEXT, DARK_TEXT, DARK_TEXT, WHITE]
for i, (item, normal, complex_) in enumerate(rows):
    y = Inches(2.0 + i * 0.72)
    add_rect(slide, Inches(0.35), y, W - Inches(0.7), Inches(0.68), row_colors[i])
    tc = text_colors[i]
    add_text(slide, item, col_x[0] + Inches(0.1), y + Inches(0.08),
             col_w[0], Inches(0.55), font_size=13, bold=(i==4), color=tc)
    add_text(slide, normal, col_x[1] + Inches(0.1), y + Inches(0.08),
             col_w[1], Inches(0.55), font_size=13, color=tc)
    add_text(slide, complex_, col_x[2] + Inches(0.1), y + Inches(0.08),
             col_w[2], Inches(0.55), font_size=13, bold=(i==4), color=tc)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド5：見落とされやすい場面
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_bg(slide)
header_bar(slide, "C-PTSDが見落とされやすい場面", "第1章：C-PTSDを理解する｜1分")
bottom_bar(slide, "Slide 5 / 17")

scenes = [
    ("内科外来", "慢性疼痛（腰痛・頭痛）・原因不明の身体症状\n機能性身体症状・身体表現性障害の診断がついている"),
    ("精神科外来", "「治療抵抗性うつ病」として複数の抗うつ薬を試されている"),
    ("要注意", "「境界性パーソナリティ障害（BPD）」診断に埋没している\nBPDの多くは幼少期虐待歴を持ち、C-PTSDとの重複が非常に多い"),
]
colors = [TEAL, NAVY, ACCENT]
for i, (label, desc) in enumerate(scenes):
    x = Inches(0.35 + i * 4.35)
    add_rect(slide, x, Inches(1.55), Inches(4.1), Inches(3.5), colors[i])
    add_text(slide, label, x + Inches(0.15), Inches(1.65),
             Inches(3.8), Inches(0.4), font_size=16, bold=True, color=WHITE)
    add_text(slide, desc, x + Inches(0.15), Inches(2.15),
             Inches(3.8), Inches(2.6), font_size=14, color=WHITE)

add_rect(slide, Inches(0.35), Inches(5.2), W - Inches(0.7), Inches(0.85), YELLOW)
add_text(slide, "「問題患者」として切り捨てないこと——反復トラウマの既往を必ず問診で確認する",
         Inches(0.55), Inches(5.27), W - Inches(1.1), Inches(0.7),
         font_size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド6：フェーズ基盤型アプローチ全体像
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_bg(slide)
header_bar(slide, "フェーズ基盤型アプローチの全体像（ISTSS 2018）", "第2章：治療の原則｜2分")
bottom_bar(slide, "Slide 6 / 17  ｜  PMID: 31283056（Bisson JI et al., J Trauma Stress, 2019）")

# ISTSS根拠
add_rect(slide, Inches(0.35), Inches(1.55), Inches(12.6), Inches(0.55), TEAL)
add_text(slide, "根拠：ISTSSガイドライン 2018 ｜ 361件のRCT ＋ 208件のメタ解析に基づく最も包括的な指針（PMID: 31283056）",
         Inches(0.5), Inches(1.58), Inches(12.3), Inches(0.5),
         font_size=13, color=WHITE, bold=True)

# 4フェーズ
phases = [
    ("Phase 1", "安全確立\nSafety", "DV継続中か確認\n社会資源の整備\nシェルター・法的介入"),
    ("Phase 2", "安定化\nStabilization", "グラウンディング\n感情調節スキル\nマインドフルネス"),
    ("Phase 3", "トラウマ処理\nTrauma Processing", "EMDR\nPE / CPT\nDBT-PTSD"),
    ("Phase 4", "社会統合\nReintegration", "対人スキル訓練\n職場復帰\nコミュニティ参加"),
]
ph_colors = [GREEN, TEAL, NAVY, ACCENT]
for i, (num, title, desc) in enumerate(phases):
    x = Inches(0.35 + i * 3.25)
    add_rect(slide, x, Inches(2.25), Inches(3.05), Inches(3.2), ph_colors[i])
    add_text(slide, num, x + Inches(0.1), Inches(2.32),
             Inches(2.8), Inches(0.35), font_size=13, color=WHITE)
    add_text(slide, title, x + Inches(0.1), Inches(2.65),
             Inches(2.8), Inches(0.75), font_size=16, bold=True, color=WHITE)
    add_text(slide, desc, x + Inches(0.1), Inches(3.45),
             Inches(2.8), Inches(1.8), font_size=13, color=RGBColor(0xFF,0xFF,0xCC))
    # 矢印
    if i < 3:
        add_text(slide, "→", Inches(3.2 + i * 3.25), Inches(3.4),
                 Inches(0.4), Inches(0.5), font_size=28, bold=True,
                 color=DARK_TEXT, align=PP_ALIGN.CENTER)

# 警告
add_rect(slide, Inches(0.35), Inches(5.6), W - Inches(0.7), Inches(0.6), LIGHT_BG, ACCENT, 2)
add_text(slide, "⚠  この順番を守ることが C-PTSD 治療の鉄則。Phase 3 を急いではいけない。",
         Inches(0.55), Inches(5.65), W - Inches(1.1), Inches(0.5),
         font_size=16, bold=True, color=ACCENT)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド7：なぜ急いでトラウマ処理をしてはいけないか
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_bg(slide)
header_bar(slide, "なぜ急いでトラウマ処理をしてはいけないか", "第2章：治療の原則｜2分")
bottom_bar(slide, "Slide 7 / 17")

add_text(slide, "感情調節スキルが不十分な状態でトラウマ記憶を掘り起こすと…",
         Inches(0.4), Inches(1.6), W - Inches(0.8), Inches(0.5),
         font_size=17, bold=True, color=NAVY)

consequences = [
    ("解離の悪化", "現実感が失われ、診察中に「消えて」しまう"),
    ("自傷・希死念慮の増悪", "感情の洪水を調節できず、自傷行為が悪化"),
    ("治療関係の崩壊", "セラピストへの不信感・怒りが爆発する"),
    ("治療脱落", "DBT-PTSD vs CPT 比較RCTでCPT群の39%が脱落"),
]
for i, (title, desc) in enumerate(consequences):
    add_rect(slide, Inches(0.35 + (i % 2) * 6.35), Inches(2.15 + (i // 2) * 1.55),
             Inches(6.1), Inches(1.35), ACCENT if i == 3 else WHITE,
             ACCENT, 1.5)
    c = WHITE if i == 3 else ACCENT
    d = WHITE if i == 3 else DARK_TEXT
    add_text(slide, f"❌  {title}", Inches(0.5 + (i % 2) * 6.35),
             Inches(2.22 + (i // 2) * 1.55), Inches(5.8), Inches(0.42),
             font_size=15, bold=True, color=c)
    add_text(slide, desc, Inches(0.5 + (i % 2) * 6.35),
             Inches(2.64 + (i // 2) * 1.55), Inches(5.8), Inches(0.6),
             font_size=13, color=d)

add_rect(slide, Inches(0.35), Inches(5.35), W - Inches(0.7), Inches(0.75), NAVY)
add_text(slide, "「急がば回れ」—— これが C-PTSD 治療の核心",
         Inches(0.55), Inches(5.42), W - Inches(1.1), Inches(0.6),
         font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド8：フェーズ1〜4の実際
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_bg(slide)
header_bar(slide, "フェーズ1〜4の実際", "第2章：治療の原則｜1分")
bottom_bar(slide, "Slide 8 / 17")

details = [
    ("Phase 1  安全確立", TEAL,
     ["DV・ホームレス・自傷の有無確認", "シェルターや法的介入の手配", "社会的資源（行政・NPO）の整備"]),
    ("Phase 2  安定化", NAVY,
     ["グラウンディング（「今ここにいる」感覚の訓練）", "感情日記・感情のラベリング", "マインドフルネス・呼吸法"]),
    ("Phase 3  トラウマ処理", ACCENT,
     ["EMDR（眼球運動による脱感作と再処理）", "PE（持続エクスポージャー）", "CPT / DBT-PTSD"]),
    ("Phase 4  社会統合", GREEN,
     ["対人スキルトレーニング", "職場復帰・就労支援", "コミュニティ参加・ピアサポート"]),
]
for i, (title, color, items) in enumerate(details):
    x = Inches(0.35 + i * 3.25)
    add_rect(slide, x, Inches(1.55), Inches(3.05), Inches(4.5), color)
    add_text(slide, title, x + Inches(0.1), Inches(1.62),
             Inches(2.8), Inches(0.5), font_size=14, bold=True, color=WHITE)
    for j, item in enumerate(items):
        add_text(slide, f"• {item}", x + Inches(0.1), Inches(2.2 + j * 0.9),
                 Inches(2.8), Inches(0.8), font_size=13, color=WHITE)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド9：EMDR
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_bg(slide)
header_bar(slide, "EMDR — 眼球運動による脱感作と再処理", "第3章：各治療法のエビデンス｜2分")
bottom_bar(slide, "Slide 9 / 17  ｜  PMID: 30857567（Karatzias T et al., Psychological Medicine, 2019）")

# 効果量ビッグナンバー
add_rect(slide, Inches(0.35), Inches(1.55), Inches(4.5), Inches(2.8), TEAL)
add_text(slide, "効果量\nHedges' g", Inches(0.5), Inches(1.75),
         Inches(4.2), Inches(0.65), font_size=14, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "−1.26", Inches(0.5), Inches(2.3), Inches(4.2), Inches(1.1),
         font_size=52, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
add_text(slide, "非常に大きな効果量", Inches(0.5), Inches(3.35),
         Inches(4.2), Inches(0.4), font_size=13, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "RCT 51件メタ解析（Karatzias 2019）", Inches(0.5), Inches(3.7),
         Inches(4.2), Inches(0.4), font_size=11, color=RGBColor(0xA0,0xD8,0xE8),
         align=PP_ALIGN.CENTER)

# 右：所見
bullet_box(slide,
    ["否定的自己概念・対人関係にも中等度〜大の効果量（k=4）",
     "青年C-PTSDへの3ヶ月EMDRで全指標が有意改善（Rolling 2024）"],
    Inches(5.1), Inches(1.55), Inches(7.85), Inches(1.55), title="主な知見", bg=WHITE)

# 注意事項
add_rect(slide, Inches(5.1), Inches(3.2), Inches(7.85), Inches(1.15), WHITE, ACCENT, 1.5)
add_text(slide, "⚠  限界・注意点", Inches(5.25), Inches(3.28),
         Inches(5), Inches(0.38), font_size=13, bold=True, color=ACCENT)
add_text(slide, "解離症状が顕著な症例では効果が著しく低下\n幼少期発症トラウマほど効果が下がる傾向\n→ 解離安定化介入を必ず先行させること",
         Inches(5.25), Inches(3.65), Inches(7.5), Inches(0.65),
         font_size=13, color=DARK_TEXT)

add_rect(slide, Inches(0.35), Inches(4.5), W - Inches(0.7), Inches(0.55), NAVY)
add_text(slide, "「解離がある → まず安定化 → その後EMDR」の順番を守る",
         Inches(0.55), Inches(4.55), W - Inches(1.1), Inches(0.45),
         font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド10：TF-CBT/PE/CPT
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_bg(slide)
header_bar(slide, "TF-CBT / PE / CPT — 強力だが脱落率に注意", "第3章：各治療法のエビデンス｜2分")
bottom_bar(slide, "Slide 10 / 17  ｜  PMID: 40737648（Korovynskyi I et al., 2025）/ PMID: 32697288（Bohus M et al., 2020）")

# 3つの効果量
stats = [("PE\n持続エクスポージャー", "d = 1.24", "反応率 68%"),
         ("CPT\n認知処理療法", "d = 1.16", "反応率 64%"),
         ("TF-CBT", "d = 1.08", "反応率 61%")]
for i, (name, es, rate) in enumerate(stats):
    x = Inches(0.35 + i * 3.0)
    add_rect(slide, x, Inches(1.55), Inches(2.8), Inches(2.5), NAVY)
    add_text(slide, name, x + Inches(0.1), Inches(1.65),
             Inches(2.6), Inches(0.6), font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, es, x + Inches(0.1), Inches(2.25), Inches(2.6), Inches(0.7),
             font_size=30, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    add_text(slide, rate, x + Inches(0.1), Inches(2.95), Inches(2.6), Inches(0.4),
             font_size=14, color=RGBColor(0xA0,0xD0,0xFF), align=PP_ALIGN.CENTER)
add_text(slide, "RCT 45件  n = 4,267（Korovynskyi 2025, PMID: 40737648）",
         Inches(0.35), Inches(4.1), Inches(9.0), Inches(0.38),
         font_size=12, color=MID_GRAY, italic=True)

# 脱落率警告（右側）
add_rect(slide, Inches(9.5), Inches(1.55), Inches(3.45), Inches(2.55), ACCENT)
add_text(slide, "脱落率の衝撃", Inches(9.65), Inches(1.65),
         Inches(3.2), Inches(0.4), font_size=14, bold=True, color=WHITE)
add_text(slide, "CPT群\n39%\n脱落", Inches(9.65), Inches(2.05),
         Inches(3.2), Inches(1.2), font_size=28, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
add_text(slide, "DBT-PTSDとの比較RCT", Inches(9.65), Inches(3.25),
         Inches(3.2), Inches(0.55), font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.35), Inches(4.6), W - Inches(0.7), Inches(0.65), LIGHT_BG, ACCENT, 1.5)
add_text(slide, "⚠  感情調節障害が強いC-PTSD患者への単独適用には要注意。「誰に使うか」が最も重要。",
         Inches(0.55), Inches(4.65), W - Inches(1.1), Inches(0.55),
         font_size=14, bold=True, color=ACCENT)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド11：DBT-PTSD
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_bg(slide)
header_bar(slide, "DBT-PTSD — JAMA Psychiatryの衝撃的結果", "第3章：各治療法のエビデンス｜2分")
bottom_bar(slide, "Slide 11 / 17  ｜  PMID: 32697288（Bohus M et al., JAMA Psychiatry, 2020）/ PMID: 39221987（Vonderlin R 2024）")

# 左：対象
add_rect(slide, Inches(0.35), Inches(1.55), Inches(3.8), Inches(3.5), WHITE, TEAL, 1.5)
add_text(slide, "対象患者", Inches(0.5), Inches(1.65),
         Inches(3.5), Inches(0.4), font_size=14, bold=True, color=TEAL)
add_text(slide, "幼少期虐待後 C-PTSD\n＋ BPD基準3つ以上\n（n = 193）",
         Inches(0.5), Inches(2.1), Inches(3.5), Inches(1.2),
         font_size=15, color=DARK_TEXT)
add_text(slide, "DBT-PTSD\n（外来45週 or 入院12週）",
         Inches(0.5), Inches(3.4), Inches(3.5), Inches(0.8),
         font_size=13, bold=True, color=TEAL)
add_text(slide, "vs  CPT（認知処理療法）",
         Inches(0.5), Inches(4.2), Inches(3.5), Inches(0.5),
         font_size=13, color=DARK_TEXT)

# 中央・右：結果比較
metrics = [
    ("寛解率", "58.4 %", "40.7 %", TEAL, ACCENT),
    ("効果量", "d = 1.35", "d = 0.98", TEAL, ACCENT),
    ("脱落率", "25.5 %", "39.0 %", TEAL, ACCENT),
]
for i, (label, dbt, cpt, dc, cc) in enumerate(metrics):
    y = Inches(1.55 + i * 1.25)
    add_rect(slide, Inches(4.4), y, Inches(4.5), Inches(1.1), dc)
    add_text(slide, f"DBT-PTSD  {dbt}", Inches(4.55), y + Inches(0.08),
             Inches(4.2), Inches(0.85), font_size=26, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(9.1), y, Inches(4.0), Inches(1.1), LIGHT_BG)
    add_text(slide, f"CPT  {cpt}", Inches(9.25), y + Inches(0.08),
             Inches(3.7), Inches(0.85), font_size=26, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.CENTER)
    add_text(slide, label, Inches(4.4), y + Inches(1.1), Inches(8.7), Inches(0.3),
             font_size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, "9ヶ月後フォローアップでも効果維持（Vonderlin 2024, PMID: 39221987）",
         Inches(4.4), Inches(5.0), Inches(8.7), Inches(0.35),
         font_size=12, color=MID_GRAY, italic=True, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.35), Inches(5.45), W - Inches(0.7), Inches(0.7), NAVY)
add_text(slide, "幼少期虐待後 C-PTSD ＋ 感情調節障害 → DBT-PTSDが現時点で最強エビデンス",
         Inches(0.55), Inches(5.52), W - Inches(1.1), Inches(0.6),
         font_size=18, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド12：薬物療法
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_bg(slide)
header_bar(slide, "薬物療法 — 補助療法の位置づけ", "第3章：各治療法のエビデンス｜2分")
bottom_bar(slide, "Slide 12 / 17  ｜  PMID: 36162349（Bajor LA et al., Psychiatry Res, 2022）")

add_rect(slide, Inches(0.35), Inches(1.55), W - Inches(0.7), Inches(0.5), ACCENT)
add_text(slide, "C-PTSD特化の薬物療法RCTはほぼ存在しない  ——  心理療法の補助として使用する",
         Inches(0.55), Inches(1.6), W - Inches(1.1), Inches(0.42),
         font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# アルゴリズム
steps = [
    ("STEP 1", "まず睡眠障害を治療する", TEAL),
    ("STEP 2\n悪夢", "プラゾシン（α1遮断薬）\nノルアドレナリン過活動を抑制", GREEN),
    ("STEP 2\n入眠障害", "ヒドロキシジン\nトラゾドン", NAVY),
    ("STEP 3\n残存症状", "SSRI\nセルトラリン・パロキセチン", TEAL),
]
for i, (step, desc, color) in enumerate(steps):
    x = Inches(0.35 + i * 3.25)
    add_rect(slide, x, Inches(2.2), Inches(3.05), Inches(2.3), color)
    add_text(slide, step, x + Inches(0.1), Inches(2.27),
             Inches(2.8), Inches(0.55), font_size=15, bold=True, color=WHITE)
    add_text(slide, desc, x + Inches(0.1), Inches(2.82),
             Inches(2.8), Inches(1.45), font_size=14, color=RGBColor(0xFF,0xFF,0xCC))
    if i < 3:
        add_text(slide, "→", Inches(3.2 + i * 3.25), Inches(3.05),
                 Inches(0.3), Inches(0.4), font_size=24, bold=True,
                 color=DARK_TEXT, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.35), Inches(4.65), W - Inches(0.7), Inches(0.55), LIGHT_BG, MID_GRAY, 1)
add_text(slide, "ケタミン・TMS・サイケデリクスは現在研究段階——標準治療として提供できる段階ではない",
         Inches(0.55), Inches(4.7), W - Inches(1.1), Inches(0.45),
         font_size=13, color=MID_GRAY, italic=True)

add_rect(slide, Inches(0.35), Inches(5.3), W - Inches(0.7), Inches(0.7), NAVY)
add_text(slide, "薬は補助。心理療法ができる環境を整えるための「下地づくり」として使う。",
         Inches(0.55), Inches(5.37), W - Inches(1.1), Inches(0.6),
         font_size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド13：スキーマ療法・IFS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_bg(slide)
header_bar(slide, "スキーマ療法・IFS — 有望だがまだ研究段階", "第3章：各治療法のエビデンス｜2分")
bottom_bar(slide, "Slide 13 / 17  ｜  PMID: 41160015（van Hemert 2025）/ PMID: 38934934（Comeau 2024）")

# 左：スキーマ療法
add_rect(slide, Inches(0.35), Inches(1.55), Inches(6.15), Inches(4.05), WHITE, TEAL, 2)
add_text(slide, "スキーマ療法（Schema Therapy）", Inches(0.5), Inches(1.65),
         Inches(5.8), Inches(0.42), font_size=16, bold=True, color=TEAL)
add_text(slide,
    "• 幼少期に形成された「早期不適応スキーマ」に働きかける\n  （「見捨てられる」「欠陥がある」などの深い信念パターン）\n\n"
    "• ST + イメージリスクリプティング + EMDRの段階的統合で\n  CAPS-5が重篤→閾値未満まで改善（van Hemert 2025）\n\n"
    "• 難民・治療抵抗例への「橋渡し」として有望\n\n"
    "⚠  RCTほぼなし——症例・パイロットレベル",
    Inches(0.5), Inches(2.15), Inches(5.8), Inches(3.2),
    font_size=13, color=DARK_TEXT)

# 右：IFS
add_rect(slide, Inches(6.75), Inches(1.55), Inches(6.2), Inches(4.05), WHITE, NAVY, 2)
add_text(slide, "IFS（内的家族システム療法）", Inches(6.9), Inches(1.65),
         Inches(5.9), Inches(0.42), font_size=16, bold=True, color=NAVY)
add_text(slide,
    "• 心の中の様々な「部分（parts）」に働きかける\n  独特のアプローチ\n\n"
    "• Comeau 2024（n=15）16週パイロット：\n  → 53%が臨床的有意な反応（CAPS-5改善10点以上）\n  → 自己慈悲・感情調節・脱中心化が改善\n\n"
    "• 忍容性良好（12回以上参加73%）\n\n"
    "⚠  対照群なし——RCT不可欠",
    Inches(6.9), Inches(2.15), Inches(5.9), Inches(3.2),
    font_size=13, color=DARK_TEXT)

add_rect(slide, Inches(0.35), Inches(5.7), W - Inches(0.7), Inches(0.5), LIGHT_BG, MID_GRAY, 1)
add_text(slide, "現時点：治療抵抗例や他の療法の準備段階として統合的に用いる可能性はあるが、RCT確立まで標準推奨不可",
         Inches(0.55), Inches(5.75), W - Inches(1.1), Inches(0.4),
         font_size=12, color=MID_GRAY, italic=True)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド14：実践的判断フロー
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_bg(slide)
header_bar(slide, "どれを選ぶか — 実践的判断フロー", "第3章：各治療法のエビデンス｜2分")
bottom_bar(slide, "Slide 14 / 17")

# スタート
add_rect(slide, Inches(5.4), Inches(1.55), Inches(2.6), Inches(0.55), NAVY)
add_text(slide, "C-PTSDと診断 / 疑い", Inches(5.4), Inches(1.57),
         Inches(2.6), Inches(0.5), font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

flow = [
    (Inches(5.4), Inches(2.25), Inches(2.6), Inches(0.6),
     LIGHT_BG, "感情調節障害が強い？\nBPD基準3つ以上？", DARK_TEXT, 12),
]
add_rect(slide, Inches(5.4), Inches(2.25), Inches(2.6), Inches(0.65), LIGHT_BG, TEAL, 1.5)
add_text(slide, "感情調節障害が強い？\nBPD基準3つ以上？",
         Inches(5.5), Inches(2.28), Inches(2.4), Inches(0.58),
         font_size=12, color=DARK_TEXT, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.35), Inches(2.55), Inches(3.5), Inches(0.65), TEAL)
add_text(slide, "YES → DBT-PTSD（第一選択）\n寛解率58%（PMID: 32697288）",
         Inches(0.45), Inches(2.58), Inches(3.3), Inches(0.58),
         font_size=12, bold=True, color=WHITE)
add_text(slide, "← YES", Inches(3.85), Inches(2.7), Inches(1.5), Inches(0.35),
         font_size=13, bold=True, color=TEAL)

add_rect(slide, Inches(5.4), Inches(3.1), Inches(2.6), Inches(0.65), LIGHT_BG, TEAL, 1.5)
add_text(slide, "解離症状が顕著？",
         Inches(5.5), Inches(3.15), Inches(2.4), Inches(0.5),
         font_size=13, color=DARK_TEXT, align=PP_ALIGN.CENTER)
add_text(slide, "↓ NO", Inches(6.3), Inches(2.9), Inches(1.0), Inches(0.2),
         font_size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.35), Inches(3.2), Inches(3.5), Inches(0.65), NAVY)
add_text(slide, "YES → 解離安定化介入\n       → その後 EMDR",
         Inches(0.45), Inches(3.23), Inches(3.3), Inches(0.58),
         font_size=12, bold=True, color=WHITE)
add_text(slide, "← YES", Inches(3.85), Inches(3.38), Inches(1.5), Inches(0.35),
         font_size=13, bold=True, color=NAVY)

add_rect(slide, Inches(5.4), Inches(3.95), Inches(2.6), Inches(0.65), LIGHT_BG, TEAL, 1.5)
add_text(slide, "幼少期発症トラウマ？\n長期的介入が必要？",
         Inches(5.5), Inches(3.98), Inches(2.4), Inches(0.58),
         font_size=12, color=DARK_TEXT, align=PP_ALIGN.CENTER)
add_text(slide, "↓ NO", Inches(6.3), Inches(3.78), Inches(1.0), Inches(0.2),
         font_size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.35), Inches(4.0), Inches(3.5), Inches(0.7), ACCENT)
add_text(slide, "YES → STAIRナラティブ療法\n       ＋ フェーズ基盤型 CPT",
         Inches(0.45), Inches(4.03), Inches(3.3), Inches(0.63),
         font_size=12, bold=True, color=WHITE)
add_text(slide, "← YES", Inches(3.85), Inches(4.2), Inches(1.5), Inches(0.35),
         font_size=13, bold=True, color=ACCENT)

add_text(slide, "↓ NO", Inches(6.3), Inches(4.63), Inches(1.0), Inches(0.2),
         font_size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(5.4), Inches(4.83), Inches(2.6), Inches(0.6), TEAL)
add_text(slide, "EMDR  または  PE / CPT",
         Inches(5.5), Inches(4.86), Inches(2.4), Inches(0.52),
         font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(9.2), Inches(2.25), Inches(3.85), Inches(3.6), LIGHT_BG, GREEN, 2)
add_text(slide, "全員に共通", Inches(9.35), Inches(2.35),
         Inches(3.5), Inches(0.4), font_size=14, bold=True, color=GREEN)
add_text(slide, "睡眠障害があれば\n薬物療法を補助的に\n\n悪夢 → プラゾシン\n入眠障害 → トラゾドン\n残存症状 → SSRI",
         Inches(9.35), Inches(2.8), Inches(3.5), Inches(2.8),
         font_size=13, color=DARK_TEXT)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド15：C-PTSDを疑うべきサイン5つ
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_bg(slide)
header_bar(slide, "C-PTSDを疑うべきサイン 5つ", "第4章：臨床チェックポイント｜1分")
bottom_bar(slide, "Slide 15 / 17")

signs = [
    ("1", "複数の科を転々としている「難治」患者"),
    ("2", "感情の激しい起伏  または  感情の平板化"),
    ("3", "医療者への不信・「問題患者」ラベルがついている"),
    ("4", "慢性疼痛・機能性身体症状が遷延している"),
    ("5", "反復・長期トラウマの既往（問診で必ず確認 ← 最重要）"),
]
colors5 = [TEAL, TEAL, ACCENT, TEAL, NAVY]
for i, (num, text) in enumerate(signs):
    add_rect(slide, Inches(0.35), Inches(1.55 + i * 0.95), Inches(0.7), Inches(0.8), colors5[i])
    add_text(slide, num, Inches(0.35), Inches(1.6 + i * 0.95), Inches(0.7), Inches(0.65),
             font_size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bold_ = i == 4
    add_rect(slide, Inches(1.15), Inches(1.55 + i * 0.95), W - Inches(1.5), Inches(0.8),
             LIGHT_BG if i % 2 == 0 else WHITE)
    add_text(slide, text, Inches(1.3), Inches(1.62 + i * 0.95), W - Inches(1.8), Inches(0.65),
             font_size=16, bold=bold_, color=DARK_TEXT)

add_rect(slide, Inches(0.35), Inches(6.35), W - Inches(0.7), Inches(0.55), YELLOW)
add_text(slide, "No.5 が最重要：「幼少期に大変なことはありましたか？」と意識的に問診する",
         Inches(0.55), Inches(6.4), W - Inches(1.1), Inches(0.45),
         font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド16：紹介のタイミングとやってはいけないこと
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_bg(slide)
header_bar(slide, "紹介のタイミングと「やってはいけないこと」", "第4章：臨床チェックポイント｜2分")
bottom_bar(slide, "Slide 16 / 17")

# 紹介適応
bullet_box(slide,
    ["自傷・希死念慮が出現したとき（緊急）",
     "解離症状が顕著なとき",
     "感情調節障害が強く通常の外来対応が困難なとき"],
    Inches(0.35), Inches(1.55), Inches(6.0), Inches(2.0),
    title="紹介のタイミング", bg=WHITE)

bullet_box(slide,
    ["DBTを提供できる精神科・心療内科クリニック",
     "トラウマ専門外来（大学病院）",
     "EMDR認定セラピストのいる施設"],
    Inches(0.35), Inches(3.65), Inches(6.0), Inches(2.0),
    title="紹介先の選び方", bg=WHITE)

# やってはいけないこと
ng = [
    ("❌  1", "安定化なしにトラウマを詳しく掘り起こす\n「詳しく教えてください」は最初の診察では禁物"),
    ("❌  2", "「もう過去のことだから大丈夫」と言う\n過去は今この瞬間も続いている体験をしている"),
    ("❌  3", "「問題患者」として切り捨て、紹介もしない\nできなければできる施設へ——それが最大の支援"),
]
for i, (num, desc) in enumerate(ng):
    add_rect(slide, Inches(6.55), Inches(1.55 + i * 1.55), Inches(6.4), Inches(1.4), ACCENT)
    add_text(slide, num, Inches(6.7), Inches(1.62 + i * 1.55),
             Inches(5.8), Inches(0.4), font_size=15, bold=True, color=WHITE)
    add_text(slide, desc, Inches(6.7), Inches(2.02 + i * 1.55),
             Inches(5.8), Inches(0.7), font_size=13, color=WHITE)

add_text(slide, "やってはいけないこと", Inches(6.55), Inches(1.35),
         Inches(3.5), Inches(0.35), font_size=14, bold=True, color=ACCENT)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド17：キーメッセージ・クロージング
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, W, H, NAVY)
add_rect(slide, 0, Inches(1.2), W, Pt(4), TEAL)

add_text(slide, "今日の 3 つのキーメッセージ", Inches(0.5), Inches(0.15),
         W - Inches(1.0), Inches(0.85), font_size=26, bold=True, color=WHITE)

keys = [
    ("1", "C-PTSDは通常PTSDとは別物",
     "DSO（情動調節困難・否定的自己概念・対人関係障害）を意識的に拾う"),
    ("2", "フェーズ基盤型が原則",
     "急いでトラウマを掘り起こさない——安全確立→安定化→トラウマ処理の順番を守る"),
    ("3", "感情調節障害＋BPD共存 → DBT-PTSDが最強エビデンス",
     "寛解率 58.4%（JAMA Psychiatry 2020, PMID: 32697288）"),
]
for i, (num, title, desc) in enumerate(keys):
    add_rect(slide, Inches(0.35), Inches(1.45 + i * 1.45), Inches(0.65), Inches(1.2), TEAL)
    add_text(slide, num, Inches(0.35), Inches(1.55 + i * 1.45), Inches(0.65), Inches(1.0),
             font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(1.1), Inches(1.45 + i * 1.45), W - Inches(1.45), Inches(1.2),
             RGBColor(0x0E, 0x1F, 0x44))
    add_text(slide, title, Inches(1.25), Inches(1.52 + i * 1.45),
             W - Inches(1.8), Inches(0.45), font_size=18, bold=True, color=YELLOW)
    add_text(slide, desc, Inches(1.25), Inches(1.97 + i * 1.45),
             W - Inches(1.8), Inches(0.42), font_size=14, color=RGBColor(0xC0, 0xD8, 0xFF))

add_rect(slide, Inches(0.35), Inches(5.85), W - Inches(0.7), Inches(1.0), TEAL)
add_text(slide, "「問題患者」は「見落とされたC-PTSD患者」かもしれない",
         Inches(0.55), Inches(5.9), W - Inches(1.1), Inches(0.85),
         font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# スライド18：参考文献
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
content_bg(slide)
header_bar(slide, "主要参考文献", "補足資料")
bottom_bar(slide, "Slide 18 / 17（補足）")

refs = [
    ("PMID: 31283056", "Bisson JI et al., J Trauma Stress, 2019 — ISTSSガイドライン"),
    ("PMID: 30857567", "Karatzias T et al., Psychological Medicine, 2019 — EMDRメタ解析"),
    ("PMID: 40737648", "Korovynskyi I et al., Georgian Medical News, 2025 — TF-CBT/PE/CPTメタ解析"),
    ("PMID: 32697288", "Bohus M et al., JAMA Psychiatry, 2020 — DBT-PTSD RCT（最重要）"),
    ("PMID: 39221987", "Vonderlin R et al., Eur J Psychotraumatol, 2024 — DBT-PTSD長期フォロー"),
    ("PMID: 41160015", "van Hemert MTC et al., Eur J Psychotraumatol, 2025 — スキーマ療法"),
    ("PMID: 38934934", "Comeau A et al., Psychol Trauma, 2024 — IFSパイロット"),
    ("PMID: 36162349", "Bajor LA et al., Psychiatry Res, 2022 — 薬物療法アルゴリズム"),
    ("PMID: 40234083", "Schaug JP et al., BMJ Mental Health, 2025 — 最新ガイドライン"),
]
txBox = slide.shapes.add_textbox(Inches(0.35), Inches(1.6), W - Inches(0.7), Inches(5.5))
tf = txBox.text_frame; tf.word_wrap = True
first = True
for pmid, desc in refs:
    if first:
        p = tf.paragraphs[0]; first = False
    else:
        p = tf.add_paragraph()
    p.space_before = Pt(5)
    r1 = p.add_run(); r1.text = f"{pmid}  "; r1.font.size = Pt(13); r1.font.bold = True; r1.font.color.rgb = TEAL
    r2 = p.add_run(); r2.text = desc; r2.font.size = Pt(13); r2.font.color.rgb = DARK_TEXT

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 保存
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
out = "/Users/torukubota/ai-management/docs/education/2026-03-20_cptsd_lecture_slides.pptx"
prs.save(out)
print(f"Saved: {out}")
