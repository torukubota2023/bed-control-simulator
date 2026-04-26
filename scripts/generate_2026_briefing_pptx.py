"""2026新基準 医師・看護師向けブリーフィング PPTX 生成スクリプト.

入力:
    docs/admin/2026新基準_医師看護師向けブリーフィング_根拠表とストーリーボード_2026-04-26.md
    の第2部 ストーリーボード（slide 0 〜 15）を構造化して PPTX 化する。

出力:
    docs/admin/2026新基準_医師看護師向けブリーフィング_2026-04-26.pptx

PPTX 化ルール（ブリーフィング第3部より）:
    - 1枚1メッセージ
    - タイトルは結論文（テーマ名ではない）
    - 本文 32pt 以上
    - キーメッセージ 48〜80pt
    - 1枚最大 5 行
    - 詳細はスピーカーノートへ
    - 表は最小限、図・棒グラフ・タイムラインで見せる
    - 赤は「6F必要度II不足」「Day6ジャンプ」「誤解注意」だけに絞る
    - 医療倫理スライドを冒頭 (slide 0) に置く

実行:
    .venv/bin/python scripts/generate_2026_briefing_pptx.py
"""

from __future__ import annotations

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path(__file__).resolve().parent.parent / "docs" / "admin" / "2026新基準_医師看護師向けブリーフィング_2026-04-26.pptx"

# 色トークン（赤は警告のみ、それ以外はニュートラル）
COLOR_TEXT = RGBColor(0x1F, 0x29, 0x37)         # ダークグレー（主要テキスト）
COLOR_KEY = RGBColor(0x1F, 0x29, 0x37)          # キーメッセージ
COLOR_DANGER = RGBColor(0xDC, 0x26, 0x26)       # 赤（警告のみ）
COLOR_ACCENT = RGBColor(0x37, 0x41, 0x51)       # アクセント（ニュートラル）
COLOR_CAPTION = RGBColor(0x6B, 0x72, 0x80)      # キャプション

# スライド定義（0 〜 15）
SLIDES = [
    {
        "n": 0,
        "title": "この資料は、必要な医療を正しく数字に残すためのものです",
        "key": "不必要な処置・虚偽記載・患者選別はしません",
        "body_lines": [
            "適応のある治療・処置・ケアを、A項目 / C項目として正しく記録します。",
            "数字目的の処置追加ではない、という前提で進めます。",
        ],
        "action": "数字目的の処置追加ではない、という前提で全員が話を聞く",
        "note": (
            "冒頭で倫理ラインを固定する。必要度達成のための処置追加ではなく、"
            "医学的適応のある治療・処置・ケアを正しく記録し制度評価と同期する資料だと明確にする。"
        ),
        "key_color": "danger",
    },
    {
        "n": 1,
        "title": "2026 年 6 月 1 日から、6F は「今まで通り」では必要度 II が足りません",
        "key": "6F 必要度 II: 直近 3 ヶ月で月 99 患者日 不足",
        "body_lines": [
            "目標 18%、現状 9.17%、不足 99.4 患者日/月（応需係数 +1.48% 加算後）",
            "12 ヶ月平均は最低ライン、直近 3 ヶ月は安全ラインとして見る。",
        ],
        "action": "「不足 pt」ではなく「月 99 患者日」を共通言語にする",
        "note": (
            "制度説明から入らず、まず危機の大きさを患者日で見せる。"
            "6F は入院数を減らさず、高稼働を維持する前提で分子を増やす必要がある。"
        ),
        "key_color": "danger",
    },
    {
        "n": 2,
        "title": "看護必要度は「重い患者を診ていること」を数字で証明する指標です",
        "key": "やった医療は、記録とコードにつながって初めて数字になります",
        "body_lines": [
            "患者 → 医師指示 → 看護記録 → 医事コード → 必要度 の流れで数字に残る。",
            "忙しい病棟だから自動的に上がる、わけではない。",
        ],
        "action": "医師が「自分の指示・記録が必要度に影響する」と理解する",
        "note": (
            "医師には看護必要度の前提知識がない。制度用語ではなく、"
            "治療実態を数字に残す仕組みとして説明する。"
        ),
    },
    {
        "n": 3,
        "title": "6F は入院数を減らさず、分子の該当患者日を増やします",
        "key": "分母は維持。増やすのは、適応のある該当日です",
        "body_lines": [
            "稼働率維持のため、非該当日を持つ長期入院患者で病床を支える現実があります。",
            "分母を減らす提案ではありません。拾い漏れを減らす運用改善です。",
        ],
        "action": "「入院を絞る」ではなく「拾い漏れを減らす」に合意する",
        "note": (
            "稼働率維持のため、非該当日を持つ長期入院患者で病床を支える現実がある。"
            "よって分母を減らす提案はしない。"
        ),
    },
    {
        "n": 4,
        "title": "6F は月 75〜80 件の入院を維持する病棟です",
        "key": "入院数を減らさず、100 患者日/月を取り戻す",
        "body_lines": [
            "6F 過去 1 年: 入院 930 件/年、平均 77.5 件/月、範囲 70〜85 件/月。",
            "この入院ボリュームを維持したまま、分子を積み増します。",
        ],
        "action": "月 75〜80 件の入院を維持したまま改善する前提を共有する",
        "note": (
            "6F の過去 1 年の入院ボリュームを示す。"
            "これを今後も維持する前提で、必要度の分子を積み増す。"
        ),
    },
    {
        "n": 5,
        "title": "A3 は「点滴ライン 3 本」ではなく「注射薬剤 3 種類以上」です",
        "key": "ラインではなく薬剤数。最大 7 日間。除外薬剤あり。",
        "body_lines": [
            "悪い例: 維持輸液だけで 3 種類カウント。抗生剤単独。",
            "良い例: 抗菌薬 + 輸液 + 電解質補正など医学的に必要な 3 種類同時投与。",
        ],
        "action": "A3 を確認するとき、医師・看護師・医事が薬剤種類と除外薬剤を確認する",
        "note": (
            "前回 PPTX の重大誤りを修正するスライド。"
            "維持輸液だけで 3 種類にしない。抗生剤単独も A3 ではない。"
        ),
        "key_color": "danger",
    },
    {
        "n": 6,
        "title": "肺炎・心不全・COPD 増悪は、A2/A3 の組み合わせで拾います",
        "key": "酸素 + 注射薬 3 種類 = A 合計 2 点の入口",
        "body_lines": [
            "Day1-2 は A7 (救急 2 点) でカバー、Day3 以降は A2/A3 で維持。",
            "A2 は喀痰吸引のみ除外。コード一覧該当行為が必要。",
        ],
        "action": "入院時 30 秒で「この肺炎は A2/A3 になるか」を考える",
        "note": (
            "医師が目の前の患者で当てはめられるようにする。"
            "A2 は喀痰吸引のみ除外、コード一覧該当行為が必要。"
            "A3 は薬剤数と除外薬剤確認。"
        ),
    },
    {
        "n": 7,
        "title": "敗血症・ショックでは A6 の 3 点項目を見落としません",
        "key": "昇圧薬・麻薬注射・抗不整脈薬・抗血栓薬持続点滴は高得点候補",
        "body_lines": [
            "条件: いずれも「注射剤のみ」または「持続点滴」が必須。",
            "内服・貼付・坐剤の麻薬は 3 点項目ではない。",
            "複数 A6 該当でも合算しない（最高得点を採用）。",
        ],
        "action": "薬剤名、投与経路、投与日、目的を当日記録する",
        "note": (
            "A6 は強いが誤用が危険。内服・貼付麻薬は 3 点ではない。"
            "複数 A6 の足し算もしない。"
        ),
    },
    {
        "n": 8,
        "title": "内科処置・検査は C 項目ごとに 2 日・4 日・5 日で数えます",
        "key": "C22 = 2 日、C21 = 4 日、C23 = 5 日",
        "body_lines": [
            "C22 (2日): 経食道心エコー、気管支鏡 等",
            "C21 (4日): CV 挿入、腰椎穿刺、CHDF、ERCP、内視鏡止血 等",
            "C23 (5日): PEG、CART、PTCD、消化管ステント 等",
        ],
        "action": "処置をした日に「C 項目候補」と医事へつなぐ",
        "note": (
            "PT-OT-ST.NET 別紙7別表1 で照合済。"
            "中心静脈注射用カテーテル挿入、腰椎穿刺、持続緩徐式血液濾過、"
            "吸着式血液浄化法、エンドトキシン選択除去用吸着式血液浄化法はいずれも C21。"
            "PEG、CART、PTCD、消化管ステントは C23。"
        ),
    },
    {
        "n": 9,
        "title": "ペイン科は「薬剤名・目的・投与日」を残すだけで改善余地があります",
        "key": "治療実態を A6 候補として照合する",
        "body_lines": [
            "薬剤名 → コード照合 → 投与日記録 → 該当日 のフローを徹底。",
            "ペイン科は在院日数が長く、6F の稼働率を支えています。",
            "点数目的ではなく、適応ある治療の可視化です。",
        ],
        "action": "ペイン科症例は週 1 回、A6 候補を医師・看護師・医事で照合する",
        "note": (
            "ペイン科は在院日数が長く、6F の稼働率を支える一方で、"
            "必要度分子の取りこぼしが起きやすい。点数目的ではなく、適応ある治療の可視化。"
        ),
    },
    {
        "n": 10,
        "title": "A7 は入院初期 2 日間だけです",
        "key": "3 日目から落ちる前提で、次の A/C を探す",
        "body_lines": [
            "Day1-2: A7（救急搬送後 / 緊急入院、2 点）でカバー。",
            "Day3 以降: A2/A3/A4/A5/A6 または C 項目で維持する設計。",
        ],
        "action": "入院 3 日目の朝に、A7 後の受け皿を確認する",
        "note": (
            "救急搬送後入院や緊急入院は 2 日間（旧 5 日間から短縮）。"
            "3 日目以降の該当維持を A/C 項目で設計する必要がある。"
        ),
    },
    {
        "n": 11,
        "title": "救急 15% は、病院全体ではなく病棟別・3 ヶ月で見られます",
        "key": "5F / 6F それぞれ、rolling 3 ヶ月で 15% 以上",
        "body_lines": [
            "予定外入院割合と制度上の救急搬送後割合は違う。",
            "下り搬送は 15% 分子の内訳（独立基準ではない）。",
            "経路記録（救急車 / 他院救急患者連携搬送料）を曖昧にしない。",
        ],
        "action": "救急・下り搬送の経路記録を曖昧にしない",
        "note": (
            "予定外入院割合と制度上の救急搬送後割合は違う。"
            "下り搬送は 15% 分子の内訳。"
        ),
    },
    {
        "n": 12,
        "title": "短手 3 は Day6 に延びた瞬間、LOS 分母が跳ねます",
        "key": "Day5 までは除外、Day6 で +6 日ジャンプ",
        "body_lines": [
            "Day1-5: LOS 分母から除外（短手 3 患者）。",
            "Day6 に延長した瞬間: 入院初日まで遡って全日数を分母にカウント。",
            "救急 15% 分母には短手 3 を含む。",
        ],
        "action": "Day5 短手 3 患者はカンファで明示判断する",
        "note": (
            "短手 3 は便利だが、Day6 延長時に平均在院日数へ不連続な影響が出る。"
            "医学的判断は変えないが、判断点を可視化する。"
        ),
        "key_color": "danger",
    },
    {
        "n": 13,
        "title": "退院集中は、空床と月曜負荷を作ります",
        "key": "金曜・土曜集中 → 週末空床 → 月曜入院集中",
        "body_lines": [
            "過去 1 年で退院集中 49 回。機会損失 年間 812〜1,215 万円規模。",
            "6F は火曜・金曜集中が目立つ。",
            "金額ではなく業務負荷の因果で伝える。",
        ],
        "action": "木曜カンファで退院日を 1 日単位で分散検討する",
        "note": (
            "退院集中は制度の直接項目ではないが、"
            "稼働率・現場負荷・患者受入に影響する。金額だけでなく業務負荷で伝える。"
        ),
    },
    {
        "n": 14,
        "title": "明日から変えるのは、3 つだけです",
        "key": "入院時 30 秒、朝ラウンド即答、木曜 6F ハドル",
        "body_lines": [
            "医師: 入院時 30 秒で A2 以上 / C1 以上 / A7 該当を考える。",
            "看護師: 朝の A 項目評価で疑問あれば医師に当日確認、翌日に持ち越さない。",
            "管理者: 週 1 回、残不足患者日を共有する。",
        ],
        "action": "医師: 入院時 30 秒確認。看護師: A 項目同日確認。管理者: 週 1 回、残不足患者日を共有する",
        "note": (
            "行動人間学的に、覚える項目を増やさない。固定習慣に落とす。"
        ),
    },
    {
        "n": 15,
        "title": "医師と看護師の協力で、月 100 患者日は現実的に届きます",
        "key": "記録 11 + 内科 40 + ペイン 9 + C 項目 40 = 100 患者日/月",
        "body_lines": [
            "記録回収 11、内科 A 項目 40、ペイン科 A6 9、C21系 16、C22系 4、C23系 20。",
            "「無理なお願い」ではなく「どの行動で何患者日が増えるか」を見る。",
            "次週から 6F で「残不足患者日/月」を運用 KPI にする。",
        ],
        "action": "次週から 6F で「残不足患者日/月」を運用 KPI にする",
        "note": (
            "最後は「無理なお願い」ではなく「どの行動で何患者日が増えるか」を見せる。"
            "責める資料ではなく、同じ数字を見る資料にする。"
        ),
    },
]


def add_text_with_style(
    text_frame,
    text: str,
    font_size_pt: int,
    bold: bool = False,
    color: Optional[RGBColor] = None,
    alignment: int = PP_ALIGN.LEFT,
    new_paragraph: bool = True,
) -> None:
    """テキストフレームにスタイル付きで段落を追加."""
    if new_paragraph:
        para = text_frame.add_paragraph()
    else:
        para = text_frame.paragraphs[0]
    para.alignment = alignment
    run = para.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    run.font.name = "Hiragino Kaku Gothic ProN"


def build_slide(prs: Presentation, spec: dict) -> None:
    """1 枚のスライドを構築する."""
    blank_layout = prs.slide_layouts[6]  # 「白紙」レイアウト
    slide = prs.slides.add_slide(blank_layout)

    # スライド全体の領域（標準 16:9）
    # title 領域: y=0.5", height=1.5"
    # key 領域:   y=2.2", height=2.4"
    # body 領域:  y=4.8", height=2.0"
    # action 領域: y=6.9", height=0.6"

    # タイトル
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(12.3), Inches(1.4)
    )
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    add_text_with_style(
        title_tf,
        spec["title"],
        font_size_pt=28,
        bold=True,
        color=COLOR_ACCENT,
        new_paragraph=False,
    )
    add_text_with_style(
        title_tf,
        f"スライド {spec['n']}",
        font_size_pt=12,
        color=COLOR_CAPTION,
    )

    # キーメッセージ（中央大）
    key_color = COLOR_DANGER if spec.get("key_color") == "danger" else COLOR_KEY
    key_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.2), Inches(12.3), Inches(2.4)
    )
    key_tf = key_box.text_frame
    key_tf.word_wrap = True
    add_text_with_style(
        key_tf,
        spec["key"],
        font_size_pt=44,
        bold=True,
        color=key_color,
        alignment=PP_ALIGN.CENTER,
        new_paragraph=False,
    )

    # 補足テキスト（5 行以内）
    body_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.0)
    )
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    for idx, line in enumerate(spec.get("body_lines", [])):
        add_text_with_style(
            body_tf,
            f"・ {line}",
            font_size_pt=20,
            color=COLOR_TEXT,
            new_paragraph=(idx > 0),
        )

    # アクション（フッター）
    action_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.5)
    )
    action_tf = action_box.text_frame
    action_tf.word_wrap = True
    add_text_with_style(
        action_tf,
        f"📌 残したい行動: {spec['action']}",
        font_size_pt=14,
        color=COLOR_CAPTION,
        new_paragraph=False,
    )

    # スピーカーノート
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = (
        f"【スライド {spec['n']}】\n"
        f"タイトル: {spec['title']}\n"
        f"キーメッセージ: {spec['key']}\n\n"
        f"【スピーカーノート】\n{spec['note']}\n\n"
        f"【聴衆に残したい行動】\n{spec['action']}"
    )


def main() -> int:
    prs = Presentation()
    # 16:9 ワイドスクリーン
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for spec in SLIDES:
        build_slide(prs, spec)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"✅ PPTX 生成完了: {OUTPUT_PATH}")
    print(f"   スライド数: {len(SLIDES)} (slide 0 〜 {SLIDES[-1]['n']})")
    print(f"   サイズ: {OUTPUT_PATH.stat().st_size:,} bytes")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
