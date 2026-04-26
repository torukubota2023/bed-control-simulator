"""2026新基準ブリーフィングスライド（医師・看護師合同向け、15枚）.

副院長判断 (2026-04-26): 2026-06-01 適用の地域包括医療病棟 新基準
（看護必要度 19%/18% + 救急搬送後15% + 平均在院日数20日）への
当院適合状況を、医師・看護師合同で共有し、明日から動ける運用変更を
合意するための 20分プレゼン資料。

設計方針:
- 危機を直視 → 達成可能な共同作戦 → 一致団結 の段階型トーン
- 個別医師名・医師コードは出さない（病棟別・診療科別・匿名分布のみ）
- アクション提案は「初期案」、現場修正提案を歓迎する明文化
- 制度基準は変えられないが、運用方法は変えられる、を明確化
- 退院集中は経営金額より現場負荷の言葉で表現
- スピーカーノート付き（20分プレゼン想定）

使い方:
    .venv/bin/python scripts/generate_2026_briefing_slides.py
    → docs/admin/2026新基準_医師看護師向けブリーフィング_2026-04-26.pptx を出力
"""
from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# テーマカラー
# ---------------------------------------------------------------------------

COLOR_PRIMARY = RGBColor(0x37, 0x41, 0x51)
COLOR_ACCENT = RGBColor(0x25, 0x63, 0xEB)
COLOR_SUCCESS = RGBColor(0x10, 0xB9, 0x81)
COLOR_WARNING = RGBColor(0xF5, 0x9E, 0x0B)
COLOR_DANGER = RGBColor(0xDC, 0x26, 0x26)
COLOR_TEXT = RGBColor(0x1F, 0x29, 0x37)
COLOR_SUBTEXT = RGBColor(0x6B, 0x72, 0x80)
COLOR_BG_LIGHT = RGBColor(0xF9, 0xFA, 0xFB)
COLOR_BG_DANGER = RGBColor(0xFE, 0xE2, 0xE2)
COLOR_BG_WARNING = RGBColor(0xFE, 0xF3, 0xC7)
COLOR_BG_SUCCESS = RGBColor(0xD1, 0xFA, 0xE5)
COLOR_BG_INFO = RGBColor(0xDB, 0xEA, 0xFE)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TODAY = date(2026, 4, 26)
TRANSITION_END = date(2026, 5, 31)
DAYS_LEFT = (TRANSITION_END - TODAY).days


# ---------------------------------------------------------------------------
# ヘルパー
# ---------------------------------------------------------------------------

def add_title(slide, text, color=COLOR_PRIMARY, size=28):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = color


def add_body(slide, left, top, width, height, text, size=16, color=COLOR_TEXT,
             bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False


def add_metric_card(slide, left, top, width, label, value, note="",
                    accent=COLOR_ACCENT, value_size=24):
    h = 1.4
    add_rect(slide, left, top, width, h, COLOR_BG_LIGHT, line_color=accent)
    add_body(slide, left + 0.15, top + 0.08, width - 0.3, 0.35,
             label, size=11, color=COLOR_SUBTEXT)
    add_body(slide, left + 0.15, top + 0.42, width - 0.3, 0.7,
             value, size=value_size, color=accent, bold=True)
    if note:
        add_body(slide, left + 0.15, top + 1.05, width - 0.3, 0.32,
                 note, size=10, color=COLOR_SUBTEXT)


def add_footer(slide, text):
    add_body(slide, 0.5, 7.05, 12.3, 0.3, text, size=9, color=COLOR_SUBTEXT)


def add_notes(slide, text):
    """スピーカーノートを追加."""
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def add_signal_dot(slide, left, top, size, color):
    """信号機の丸印."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top),
                                    Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = COLOR_PRIMARY
    shape.shadow.inherit = False


# ---------------------------------------------------------------------------
# プレゼンテーション本体
# ---------------------------------------------------------------------------

def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # =====================================================================
    # SLIDE 1: タイトル
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 13.333, 7.5, COLOR_PRIMARY)
    add_rect(s, 0, 5.6, 13.333, 1.9, COLOR_ACCENT)
    add_body(s, 0.5, 1.5, 12.3, 1.0,
             "2026年6月1日 新基準",
             size=44, color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_body(s, 0.5, 2.7, 12.3, 1.4,
             "地域包括医療病棟を守るために、\nいま医師・看護師に伝えたいこと",
             size=36, color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_body(s, 0.5, 4.6, 12.3, 0.6,
             "— 高齢救急の地域受け皿を、私たち全員で続けていく —",
             size=20, color=RGBColor(0xD1, 0xD5, 0xDB), align=PP_ALIGN.CENTER)
    add_body(s, 0.5, 5.8, 12.3, 0.5,
             "医師・看護師合同ブリーフィング", size=20, color=COLOR_WHITE,
             bold=True, align=PP_ALIGN.CENTER)
    add_body(s, 0.5, 6.4, 12.3, 0.5,
             f"おもろまちメディカルセンター  副院長  久保田 徹",
             size=14, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_body(s, 0.5, 6.85, 12.3, 0.4,
             f"2026-04-26  経過措置終了まで あと {DAYS_LEFT} 日",
             size=14, color=RGBColor(0xE5, 0xE7, 0xEB), align=PP_ALIGN.CENTER)
    add_notes(s,
        "（30秒）今日は20分のブリーフィングです。\n"
        "2026年6月1日から、当院の地域包括医療病棟は新しい施設基準で運営されます。"
        "経過措置（困難月の除外）は5月31日で終了し、それ以降は例外なく本則適用です。\n"
        "5月31日まであと35日。今日のブリーフィングは『責める資料』ではなく、"
        "『地域の高齢救急の受け皿である当院の病棟を、医師・看護師全員で守るための共同作戦の確認』です。"
    )

    # =====================================================================
    # SLIDE 2: 今日の30秒結論
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_title(s, "📌 今日の結論：3つの基準＋退院集中、このままだと未達。でも一緒なら達成できる。")

    # 4枚カード横並び
    add_metric_card(s, 0.5, 1.5, 3.0,
                    "❶ 看護必要度", "5F あと一歩 / 6F 危機",
                    note="6F は新基準 19%・18% 両方で未達",
                    accent=COLOR_DANGER, value_size=14)
    add_metric_card(s, 3.7, 1.5, 3.0,
                    "❷ 救急搬送後 15%", "rolling 3ヶ月 で判定",
                    note="救急車搬送＋他院からの下り搬送 が分子",
                    accent=COLOR_WARNING, value_size=14)
    add_metric_card(s, 6.9, 1.5, 3.0,
                    "❸ 平均在院日数 20日", "rolling 90日 / 短手3 階段関数",
                    note="Day5/Day6 の境界に注意",
                    accent=COLOR_WARNING, value_size=14)
    add_metric_card(s, 10.1, 1.5, 2.7,
                    "＋ 退院集中", "現場負荷 + 機会損失",
                    note="金+土集中＋月曜入院ピーク",
                    accent=COLOR_ACCENT, value_size=14)

    # 共同作戦
    add_rect(s, 0.5, 3.4, 12.3, 1.6, COLOR_BG_INFO, line_color=COLOR_ACCENT)
    add_body(s, 0.8, 3.55, 11.7, 0.5,
             "🤝 共同作戦の合意ポイント",
             size=18, color=COLOR_ACCENT, bold=True)
    add_body(s, 0.8, 4.1, 11.7, 0.9,
             "・制度基準は変えられない（19% / 18% / 15% / 20日 は固定）\n"
             "・しかし、運用方法は私たちで変えられる ← 今日のお願い\n"
             "・本日の提案は「初期案」、現場の修正提案を歓迎します",
             size=15, color=COLOR_TEXT)

    # 5本柱
    add_body(s, 0.5, 5.3, 12.3, 0.5,
             "今日のお話 — 5つの柱",
             size=16, color=COLOR_PRIMARY, bold=True)
    add_body(s, 0.7, 5.85, 12.0, 1.1,
             "① 制度全体像と当院の現状（信号機で見る5F・6F）\n"
             "② 危機の中身（看護必要度・救急15%・LOS・退院集中の4テーマ）\n"
             "③ このまま6/1を迎えると何が起きるか\n"
             "④ 医師・看護師それぞれに明日からお願いしたいこと\n"
             "⑤ 一緒なら達成できる — 月次振り返りで進捗を共有",
             size=14, color=COLOR_TEXT)
    add_footer(s, "Slide 2 / 15  •  2026-04-26 ブリーフィング")
    add_notes(s,
        "（90秒）まず結論からお話しします。\n"
        "当院の地域包括医療病棟には、6月1日以降クリアし続けないといけない4つのハードルがあります。"
        "看護必要度、救急搬送後15%、平均在院日数、そして退院集中問題。\n"
        "結論としては、いずれもこのままだと未達か未達ギリギリです。"
        "しかし制度基準そのものは動かせなくても、私たちの運用方法は変えられる。\n"
        "今日の提案はあくまで『初期案』です。"
        "現場の医師・看護師の皆さんからの修正提案を歓迎します。一緒に作戦を立て直しましょう。"
    )

    # =====================================================================
    # SLIDE 3: 制度全体像
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_title(s, "🗺️ 制度全体像：4つのハードルが互いに繋がっている")

    add_body(s, 0.5, 1.3, 12.3, 0.5,
             "それぞれが独立した別問題ではなく、運用が連動する4つの基準です。",
             size=14, color=COLOR_SUBTEXT)

    # 4要素を矢印で結ぶ
    cards = [
        ("看護必要度", "新基準\nⅠ 19% / Ⅱ 18%", "rolling 3ヶ月\n病棟別判定", COLOR_DANGER, 0.5),
        ("救急搬送後\n15%", "救急車搬送\n＋下り搬送", "rolling 3ヶ月\n短手3を分母含む", COLOR_WARNING, 3.7),
        ("平均在院日数\n20日", "rolling 90日\n短手3 階段関数", "Day5/Day6\n境界に注意", COLOR_WARNING, 6.9),
        ("退院集中", "金+土に集中\n月曜入院ピーク", "現場負荷増 +\n機会損失", COLOR_ACCENT, 10.1),
    ]
    for label, val, note, color, x in cards:
        add_rect(s, x, 2.0, 3.0, 2.4, COLOR_BG_LIGHT, line_color=color)
        add_body(s, x + 0.15, 2.1, 2.7, 0.7, label, size=15, color=color, bold=True)
        add_body(s, x + 0.15, 2.85, 2.7, 0.85, val, size=12, color=COLOR_TEXT, bold=True)
        add_body(s, x + 0.15, 3.75, 2.7, 0.6, note, size=10, color=COLOR_SUBTEXT)

    # 連動の説明
    add_rect(s, 0.5, 4.7, 12.3, 1.1, COLOR_BG_INFO, line_color=COLOR_ACCENT)
    add_body(s, 0.8, 4.85, 11.7, 0.4,
             "🔗 4つは互いに引っ張り合う", size=14, color=COLOR_ACCENT, bold=True)
    add_body(s, 0.8, 5.25, 11.7, 0.55,
             "在院日数を短くすれば → 退院増 → 退院集中 → 月曜入院ピーク → 現場負荷／"
             "稼働率を上げるため在院日数を延ばすと → 必要度の分母が膨らみ該当割合↓",
             size=11, color=COLOR_TEXT)

    # キーメッセージ
    add_body(s, 0.5, 6.0, 12.3, 0.4,
             "💡 だから「単一指標で頑張る」のではなく、4つを同時に成り立たせる運用が必要",
             size=15, color=COLOR_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_body(s, 0.5, 6.5, 12.3, 0.4,
             "鍵: 入院初期に必要度をしっかり拾い、Day8 以降は退院支援を加速する",
             size=14, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    add_footer(s, "Slide 3 / 15")
    add_notes(s,
        "（90秒）この4つのハードルは独立した別問題ではありません。\n"
        "例えば、稼働率を上げようとして在院日数を延ばすと、看護必要度の分母が膨らんで該当割合が下がります。"
        "逆に、在院日数を短くすれば退院数が増えて、退院集中が起きやすくなる。\n"
        "つまり、単一の指標だけを追いかけて頑張ると、別の指標が悪化する構造です。"
        "だから今日のお話は『4つを同時に成り立たせる運用ルール』が鍵になります。\n"
        "後ほどお話しする協力依頼は、すべてこの『同時成立』を意識した提案です。"
    )

    # =====================================================================
    # SLIDE 4: 当院の信号機
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_title(s, "🚦 当院の現状（5F / 6F、12ヶ月平均ベース）")

    add_body(s, 0.5, 1.3, 12.3, 0.5,
             "信号機の色 = 6/1 適用後に達成できるかの見込み（応需係数 +1.48pt 加算後）",
             size=13, color=COLOR_SUBTEXT)

    # ヘッダー
    add_body(s, 0.5, 1.95, 1.5, 0.4, "病棟", size=14, color=COLOR_TEXT, bold=True)
    add_body(s, 2.1, 1.95, 2.3, 0.4, "看護必要度Ⅰ (19%)", size=14, color=COLOR_TEXT, bold=True)
    add_body(s, 4.5, 1.95, 2.3, 0.4, "看護必要度Ⅱ (18%)", size=14, color=COLOR_TEXT, bold=True)
    add_body(s, 6.9, 1.95, 2.3, 0.4, "救急搬送後 15%", size=14, color=COLOR_TEXT, bold=True)
    add_body(s, 9.3, 1.95, 1.8, 0.4, "平均在院日数20日", size=14, color=COLOR_TEXT, bold=True)
    add_body(s, 11.2, 1.95, 1.8, 0.4, "退院集中対策", size=14, color=COLOR_TEXT, bold=True)

    rows = [
        # ward, I_status, I_text, II_status, II_text, ER, LOS, DC
        ("5F", COLOR_SUCCESS, "19.70%\n達成見込み",
         COLOR_WARNING, "17.85%\nあと0.15pt",
         COLOR_WARNING, "rolling精査中",
         COLOR_SUCCESS, "16.5日達成",
         COLOR_WARNING, "枠超過22日"),
        ("6F", COLOR_DANGER, "17.54%\n−1.46pt 不足",
         COLOR_DANGER, "14.62%\n−3.38pt 不足",
         COLOR_WARNING, "rolling精査中",
         COLOR_SUCCESS, "15.6日達成",
         COLOR_DANGER, "枠超過27日"),
    ]
    y = 2.5
    for ward, c1, t1, c2, t2, c3, lt3, c4, lt4, c5, lt5 in rows:
        add_rect(s, 0.5, y, 12.3, 1.4, COLOR_BG_LIGHT)
        add_body(s, 0.7, y + 0.4, 1.2, 0.6, ward, size=22, color=COLOR_PRIMARY, bold=True)
        # ⅠⅡ ER LOS DC
        for i, (color, txt) in enumerate(
                [(c1, t1), (c2, t2), (c3, lt3), (c4, lt4), (c5, lt5)]):
            x = 2.1 + i * 2.4 if i < 3 else (9.3 if i == 3 else 11.2)
            w = 2.3 if i < 3 else 1.8
            add_signal_dot(s, x + 0.05, y + 0.15, 0.35, color)
            add_body(s, x + 0.5, y + 0.15, w - 0.5, 0.35, "", size=10, color=COLOR_TEXT)
            add_body(s, x + 0.05, y + 0.55, w, 0.85, txt, size=11, color=COLOR_TEXT, bold=True)
        y += 1.5

    # キーメッセージ
    add_rect(s, 0.5, 5.7, 12.3, 1.1, COLOR_BG_DANGER, line_color=COLOR_DANGER)
    add_body(s, 0.8, 5.85, 11.7, 0.4,
             "🚨 6F が最大の課題：必要度Ⅰ・Ⅱ の両方が応需係数加算後でも未達",
             size=15, color=COLOR_DANGER, bold=True)
    add_body(s, 0.8, 6.3, 11.7, 0.5,
             "5F は両指標とも達成圏内。6F のⅠ -1.46pt（=月17.5患者日不足）、Ⅱ -3.38pt（=月40.7患者日不足）を、"
             "犯人探しではなく病棟全体の運用で取り戻す。",
             size=12, color=COLOR_TEXT)
    add_footer(s, "Slide 4 / 15  •  応需係数 = 救急279件 ÷ 94床 × 0.005 = +1.48pt")
    add_notes(s,
        "（90秒）これが当院の現状です。緑が達成見込み、黄色が要注意、赤が未達。\n"
        "ご覧の通り、5F病棟は両指標ともほぼ達成圏内です。あと一歩です。"
        "問題は6F病棟です。看護必要度Ⅰ・Ⅱの両方が、救急患者応需係数を加算した後でも"
        "新基準を超えられていません。\n"
        "ここで強調したいのは、これは『誰か特定の医師や看護師の問題』ではないということです。"
        "6F は内科とペイン科の混合病棟で、診療科の特性上どうしても必要度が低めに出やすい。"
        "病棟全体の運用ルールを変えることで、必要度を取り戻していきます。"
    )

    # =====================================================================
    # SLIDE 5: 看護必要度 12ヶ月平均
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_title(s, "❶ 看護必要度：6F が新基準 19% / 18% から最も離れている")

    add_body(s, 0.5, 1.3, 12.3, 0.5,
             "12ヶ月平均（2025-04 〜 2026-03、応需係数 +1.48pt 加算後）",
             size=13, color=COLOR_SUBTEXT)

    # 5F カード
    add_rect(s, 0.5, 1.95, 6.0, 2.4, COLOR_BG_SUCCESS, line_color=COLOR_SUCCESS)
    add_body(s, 0.7, 2.05, 5.6, 0.45,
             "5F 病棟（外科・整形・脳外）", size=16, color=COLOR_PRIMARY, bold=True)
    add_body(s, 0.7, 2.55, 5.6, 0.4,
             "→ あと一歩で達成", size=13, color=COLOR_SUCCESS, bold=True)
    add_body(s, 0.7, 3.05, 5.6, 0.5,
             "Ⅰ: 18.21% → +応需 19.70%   ✅ 19%達成",
             size=13, color=COLOR_TEXT, bold=True)
    add_body(s, 0.7, 3.55, 5.6, 0.5,
             "Ⅱ: 16.36% → +応需 17.85%   ⚠️ 18%まで -0.15pt",
             size=13, color=COLOR_TEXT, bold=True)

    # 6F カード
    add_rect(s, 6.8, 1.95, 6.0, 2.4, COLOR_BG_DANGER, line_color=COLOR_DANGER)
    add_body(s, 7.0, 2.05, 5.6, 0.45,
             "6F 病棟（内科・ペイン科）", size=16, color=COLOR_PRIMARY, bold=True)
    add_body(s, 7.0, 2.55, 5.6, 0.4,
             "→ Ⅰ・Ⅱ 両方で未達", size=13, color=COLOR_DANGER, bold=True)
    add_body(s, 7.0, 3.05, 5.6, 0.5,
             "Ⅰ: 16.06% → +応需 17.54%   🔴 19%まで -1.46pt",
             size=13, color=COLOR_TEXT, bold=True)
    add_body(s, 7.0, 3.55, 5.6, 0.5,
             "Ⅱ: 13.13% → +応需 14.62%   🔴 18%まで -3.38pt",
             size=13, color=COLOR_TEXT, bold=True)

    # 救急患者応需係数の説明
    add_rect(s, 0.5, 4.6, 12.3, 1.0, COLOR_BG_INFO, line_color=COLOR_ACCENT)
    add_body(s, 0.7, 4.7, 12.0, 0.4,
             "💡 救急患者応需係数とは？", size=14, color=COLOR_ACCENT, bold=True)
    add_body(s, 0.7, 5.1, 12.0, 0.5,
             "令和8改定で新設された加算: (年間救急搬送 ÷ 病床数) × 0.005、上限10%。"
             "当院 = 279件 ÷ 94床 × 0.005 = +1.48pt。"
             "つまり、自力で残り 17.5% / 16.5% を稼ぐ必要がある。",
             size=12, color=COLOR_TEXT)

    # 取り戻す目安
    add_body(s, 0.5, 5.85, 12.3, 0.4,
             "📐 6F が達成するために取り戻す目安（月平均）",
             size=14, color=COLOR_PRIMARY, bold=True)
    add_body(s, 0.7, 6.3, 12.0, 0.7,
             "Ⅰ：あと月 17.5 患者日（=1日0.6人 該当判定が増えれば達成）\n"
             "Ⅱ：あと月 40.7 患者日（=1日1.4人 該当判定が増えれば達成）",
             size=13, color=COLOR_TEXT, bold=True)
    add_footer(s, "Slide 5 / 15  •  出典: data/nursing_necessity_2025fy.csv")
    add_notes(s,
        "（90秒）看護必要度の現状を、もう少し細かく見ます。\n"
        "5Fは緑、6Fは赤です。6Fの不足を『月あたり何人を該当に持っていけば達成か』に翻訳すると、"
        "Ⅰ で月17.5患者日、つまり1日あたり0.6人。Ⅱ で月40.7患者日、つまり1日あたり1.4人。\n"
        "これは『新しい処置を増やす』のではなく、"
        "『すでに実施している医学的に必要なケアを正しく数字に残す』ことで十分達成できる範囲です。\n"
        "数字だけ見ると重く感じますが、後ほどの協力依頼の中で具体的にどのケアを記録すればよいかを示します。"
    )

    # =====================================================================
    # SLIDE 6: 直近3ヶ月の悪化
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_title(s, "❶続：直近3ヶ月で 6F が崖に近づいている", color=COLOR_DANGER)

    add_body(s, 0.5, 1.3, 12.3, 0.5,
             "2026-01〜2026-03 の3ヶ月平均（rolling 3ヶ月で 6/1 以降は判定される）",
             size=13, color=COLOR_SUBTEXT)

    # 比較テーブル
    add_rect(s, 0.5, 1.9, 12.3, 0.5, COLOR_PRIMARY)
    headers = ["指標", "12ヶ月平均", "直近3ヶ月", "推移", "新基準ギャップ（直近3ヶ月）"]
    widths = [2.5, 2.0, 2.0, 1.5, 4.3]
    x = 0.5
    for h, w in zip(headers, widths):
        add_body(s, x + 0.1, 1.95, w - 0.2, 0.4, h, size=13, color=COLOR_WHITE, bold=True)
        x += w

    rows = [
        ("5F 必要度Ⅰ", "19.70%", "16.92%", "↓ 2.78pt", "−2.08pt 未達", COLOR_DANGER),
        ("5F 必要度Ⅱ", "17.85%", "14.99%", "↓ 2.86pt", "−3.01pt 未達", COLOR_DANGER),
        ("6F 必要度Ⅰ", "17.54%", "12.57%", "↓ 4.97pt", "−6.43pt 未達", COLOR_DANGER),
        ("6F 必要度Ⅱ", "14.62%", "9.11%",  "↓ 5.51pt", "−8.89pt 未達 🚨", COLOR_DANGER),
    ]
    y = 2.5
    for label, m12, m3, delta, gap, gap_color in rows:
        add_rect(s, 0.5, y, 12.3, 0.55, COLOR_BG_LIGHT)
        add_body(s, 0.6, y + 0.1, 2.4, 0.4, label, size=12, color=COLOR_TEXT, bold=True)
        add_body(s, 2.6, y + 0.1, 1.8, 0.4, m12, size=12, color=COLOR_TEXT)
        add_body(s, 4.6, y + 0.1, 1.8, 0.4, m3, size=12, color=COLOR_TEXT, bold=True)
        add_body(s, 6.6, y + 0.1, 1.4, 0.4, delta, size=12, color=COLOR_DANGER, bold=True)
        add_body(s, 8.1, y + 0.1, 4.5, 0.4, gap, size=12, color=gap_color, bold=True)
        y += 0.6

    # 何が起きているか
    add_rect(s, 0.5, 5.05, 12.3, 1.5, COLOR_BG_WARNING, line_color=COLOR_WARNING)
    add_body(s, 0.8, 5.2, 11.7, 0.4,
             "⚠️ 直近3ヶ月で 1〜5pt の悪化、6F は 5pt 級。原因は1つではなく複合的：",
             size=14, color=COLOR_WARNING, bold=True)
    add_body(s, 0.8, 5.65, 11.7, 0.85,
             "・冬季の入院構成変化（軽症の予定外入院が増加）\n"
             "・在院日数延長による分母の膨らみ（同じ該当人数でも%は下がる）\n"
             "・実施しているケアが評価項目として記録されていない可能性（要確認）",
             size=12, color=COLOR_TEXT)
    add_footer(s, "Slide 6 / 15")
    add_notes(s,
        "（90秒）こちらが、直近3ヶ月のデータです。\n"
        "12ヶ月平均と比べて、5Fも6Fも1〜5pt 悪化しています。特に6Fの必要度Ⅱは"
        "直近3ヶ月で9.11%まで落ちており、新基準18%まで実に8.89pt 不足。これは"
        "rolling 3ヶ月評価が導入される6月1日に、このまま当てはめれば未達確定の数字です。\n"
        "原因は1つではなく複合的です。冬季の入院構成、在院日数の延長による分母の膨らみ、"
        "そして — 『実施しているケアが評価項目として正しく記録されていない可能性』。\n"
        "後者については、現場でやっている処置・ケアを評価項目の言葉に変換する余地があります。"
        "後の協力依頼スライドでお願いする予定です。"
    )

    # =====================================================================
    # SLIDE 7: 救急搬送後15%
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_title(s, "❷ 救急搬送後15%：救急車搬送＋下り搬送ルートで分子を稼ぐ")

    add_body(s, 0.5, 1.3, 12.3, 0.5,
             "2026-06-01以降、rolling 3ヶ月（病棟別）、短手3患者を分母に含めて判定",
             size=13, color=COLOR_SUBTEXT)

    # 分子の中身
    add_rect(s, 0.5, 1.9, 12.3, 0.5, COLOR_ACCENT)
    add_body(s, 0.7, 1.95, 12.0, 0.4,
             "📐 分子に含まれる入院ルート（救急15% を稼ぐ患者）",
             size=14, color=COLOR_WHITE, bold=True)

    # 2列レイアウト
    add_rect(s, 0.5, 2.5, 6.0, 1.7, COLOR_BG_LIGHT, line_color=COLOR_DANGER)
    add_body(s, 0.7, 2.6, 5.6, 0.45,
             "🚑 救急車搬送", size=15, color=COLOR_DANGER, bold=True)
    add_body(s, 0.7, 3.1, 5.6, 1.05,
             "・自院救急受入（直接救急車から入院）\n"
             "・2025FY実績: 年間 279 件 / 1,823 件中 = 15.3%\n"
             "・月別変動大（2.4% 〜 27.4%）",
             size=12, color=COLOR_TEXT)

    add_rect(s, 6.8, 2.5, 6.0, 1.7, COLOR_BG_LIGHT, line_color=COLOR_WARNING)
    add_body(s, 7.0, 2.6, 5.6, 0.45,
             "🏥 他院からの下り搬送", size=15, color=COLOR_WARNING, bold=True)
    add_body(s, 7.0, 3.1, 5.6, 1.05,
             "・他院で「救急患者連携搬送料」を算定して当院へ搬送\n"
             "・地域の急性期病院からの転院受入が分子に含まれる\n"
             "・連携室・救急室との情報共有が鍵",
             size=12, color=COLOR_TEXT)

    # 2025FY データの限界
    add_rect(s, 0.5, 4.4, 12.3, 1.5, COLOR_BG_WARNING, line_color=COLOR_WARNING)
    add_body(s, 0.7, 4.55, 12.0, 0.45,
             "⚠️ 2025FYデータの解釈について（正直にお伝えします）",
             size=14, color=COLOR_WARNING, bold=True)
    add_body(s, 0.7, 5.0, 12.0, 0.85,
             "2025FYデータの「予定外入院割合」は、制度上の「救急搬送後患者割合」と完全には一致しません。"
             "外来紹介・連携室・ウォークインなど、救急搬送ではない予定外入院も混在している可能性があります。"
             "→ 2025FYデータは『制度判定の厳密値』ではなく、『リスク把握・傾向確認』として位置付けます。",
             size=12, color=COLOR_TEXT)

    # 移行プラン
    add_rect(s, 0.5, 6.05, 12.3, 0.85, COLOR_BG_INFO, line_color=COLOR_ACCENT)
    add_body(s, 0.7, 6.15, 12.0, 0.4,
             "📅 厳密判定への段階的移行（既に運用中）",
             size=14, color=COLOR_ACCENT, bold=True)
    add_body(s, 0.7, 6.55, 12.0, 0.4,
             "2026-04 以降: 詳細5区分の入院経路 + 手術有無を記録 → 5月までは手動シードで補完 → "
             "2026-07 以降: 純実データで rolling 3ヶ月計算へ移行",
             size=11, color=COLOR_TEXT)
    add_footer(s, "Slide 7 / 15")
    add_notes(s,
        "（90秒）次は救急搬送後15%です。\n"
        "ここで重要なのは『分子に含まれるのは何か』を全員で共有することです。"
        "分子は『救急車で当院に直接入院した患者』と『他院で救急患者連携搬送料を算定して当院に転院してきた患者』の合計。\n"
        "この『下り搬送』のルートは、地域の急性期病院との連携が鍵です。"
        "連携室・救急室との情報共有を密にして、救急15%基準の達成に貢献していきます。\n"
        "なお、2025年度のデータについては正直にお伝えしますが、"
        "事務システムの『予定外入院』カテゴリと制度上の『救急搬送後』が完全には一致しません。"
        "今年4月から詳細記録を開始しているので、7月以降には純実データで rolling 3ヶ月の判定ができるようになります。"
    )

    # =====================================================================
    # SLIDE 8: 平均在院日数 + 短手3
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_title(s, "❸ 平均在院日数20日 + 短手3患者の階段関数：Day 5/Day 6 が境界")

    add_body(s, 0.5, 1.3, 12.3, 0.5,
             "2026-06-01以降、rolling 90日 で判定。短手3患者の特殊扱いに注意。",
             size=13, color=COLOR_SUBTEXT)

    # 階段関数の説明
    add_body(s, 0.5, 1.85, 12.3, 0.4,
             "📐 短手3患者（短期滞在手術等基本料3）の特殊ルール",
             size=15, color=COLOR_PRIMARY, bold=True)

    # Day1-5 のカード（含めない）
    add_rect(s, 0.5, 2.4, 6.0, 2.0, COLOR_BG_SUCCESS, line_color=COLOR_SUCCESS)
    add_body(s, 0.7, 2.5, 5.6, 0.45,
             "Day 1〜Day 5 まで滞在（予定通り）", size=14, color=COLOR_SUCCESS, bold=True)
    add_body(s, 0.7, 2.95, 5.6, 1.4,
             "→ 平均在院日数の分母に「含めない」\n"
             "→ 病棟全体の LOS を押し下げる効果\n"
             "→ つまり、計画通り 5 日以内に退院すれば\n"
             "　 LOS 基準達成に味方をしてくれる",
             size=12, color=COLOR_TEXT)

    # Day6+ のカード（遡って計上）
    add_rect(s, 6.8, 2.4, 6.0, 2.0, COLOR_BG_DANGER, line_color=COLOR_DANGER)
    add_body(s, 7.0, 2.5, 5.6, 0.45,
             "Day 6 以降に滞在が延びた瞬間 🚨", size=14, color=COLOR_DANGER, bold=True)
    add_body(s, 7.0, 2.95, 5.6, 1.4,
             "→ 入院初日まで遡って「全日数」を分母計上\n"
             "→ Day 5 と Day 6 の境界で LOS 分母が +6日 ジャンプ\n"
             "→ 1人の Day 6 突破で病棟 LOS が悪化\n"
             "→ Day 5 アラートで全力退院支援が必須",
             size=12, color=COLOR_TEXT)

    # 視覚化（簡易棒）
    add_body(s, 0.5, 4.5, 12.3, 0.4,
             "イメージ：分母への寄与（短手3 1名）",
             size=12, color=COLOR_SUBTEXT)
    # Day 1-5 = 0、Day 6+ = 6+
    add_body(s, 0.5, 4.85, 1.5, 0.3, "Day 1-5", size=11, color=COLOR_TEXT)
    add_rect(s, 2.0, 4.9, 0.2, 0.25, COLOR_SUCCESS)
    add_body(s, 2.4, 4.85, 4.0, 0.3, "0 日（分母に含めない）", size=11, color=COLOR_TEXT)
    add_body(s, 0.5, 5.25, 1.5, 0.3, "Day 6", size=11, color=COLOR_TEXT)
    add_rect(s, 2.0, 5.30, 3.0, 0.25, COLOR_DANGER)
    add_body(s, 5.2, 5.25, 4.0, 0.3, "+6 日 ジャンプ（入院初日まで遡及計上）",
             size=11, color=COLOR_DANGER, bold=True)

    # 現場で何ができるか
    add_rect(s, 0.5, 5.85, 12.3, 1.1, COLOR_BG_INFO, line_color=COLOR_ACCENT)
    add_body(s, 0.7, 5.95, 12.0, 0.4,
             "✅ 現場でできること：Day 5 アラートを「全員で見る」",
             size=14, color=COLOR_ACCENT, bold=True)
    add_body(s, 0.7, 6.4, 12.0, 0.5,
             "ベッドコントロールアプリで Day 5 到達の短手3患者を自動表示済み。"
             "回診時・カンファでこのアラートを確認し、Day 6 突入を防ぐ。",
             size=12, color=COLOR_TEXT)
    add_footer(s, "Slide 8 / 15")
    add_notes(s,
        "（90秒）平均在院日数20日基準で見落としがちなのが、短手3患者の特殊ルールです。\n"
        "短手3とは、白内障やポリペク等の予定された短期入院患者のこと。"
        "Day 5 までに退院すれば、その入院は分母に含めないルールがあります。"
        "つまり病棟全体のLOSを押し下げてくれる味方なんです。\n"
        "でも、Day 6 以降に滞在が延びた瞬間、ルールが変わります。"
        "入院初日まで遡って全日数が分母計上される。"
        "つまり1人の Day 6 突破で、病棟LOSが+6日ジャンプします。\n"
        "これを防ぐため、アプリには Day 5 到達の短手3患者を自動アラート表示する機能があります。"
        "回診時、カンファ時に必ずこのアラートを全員で確認してください。"
    )

    # =====================================================================
    # SLIDE 9: 退院集中
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_title(s, "❹ 退院集中：金+土に退院が偏り、月曜入院ピークと重なる")

    add_body(s, 0.5, 1.3, 12.3, 0.5,
             "過去1年データ（n=1,793件）— 退院曜日と土日入院補充の現実",
             size=13, color=COLOR_SUBTEXT)

    # 4 KPI カード
    add_metric_card(s, 0.5, 1.9, 3.0,
                    "退院集中（5名超/日）の発生回数",
                    "49 回 / 年",
                    note="5F: 22日 + 6F: 27日 = 約2週に1回",
                    accent=COLOR_DANGER, value_size=22)
    add_metric_card(s, 3.7, 1.9, 3.0,
                    "金+土 退院シェア",
                    "5F 36% / 6F 35%",
                    note="退院全体の 1/3 以上が金土に集中",
                    accent=COLOR_WARNING, value_size=18)
    add_metric_card(s, 6.9, 1.9, 3.0,
                    "土曜入院補充",
                    "2.2 件/日",
                    note="平日入院 7.4件/日 の 30% しかない",
                    accent=COLOR_WARNING, value_size=22)
    add_metric_card(s, 10.1, 1.9, 2.7,
                    "日曜入院補充",
                    "0.3 件/日",
                    note="ほぼゼロ（病棟は静止）",
                    accent=COLOR_DANGER, value_size=22)

    # 因果図
    add_rect(s, 0.5, 3.6, 12.3, 1.7, COLOR_BG_DANGER, line_color=COLOR_DANGER)
    add_body(s, 0.7, 3.7, 12.0, 0.45,
             "🔄 因果のサイクル（毎週繰り返し起きている）",
             size=14, color=COLOR_DANGER, bold=True)
    add_body(s, 0.7, 4.15, 12.0, 1.1,
             "金+土に退院が集中  →  土日に新規入院ほぼなし → 週末空床滞留 → 月曜にまとめて入院ピーク → "
             "月曜の現場負荷急増 → 看護師オーバーワーク → 必要度評価の精度低下リスク → "
             "稼働率は不安定なまま、施設基準にも負の影響",
             size=12, color=COLOR_TEXT)

    # 機会損失（経営）
    add_rect(s, 0.5, 5.5, 6.0, 1.4, COLOR_BG_WARNING, line_color=COLOR_WARNING)
    add_body(s, 0.7, 5.6, 5.6, 0.4,
             "💴 機会損失（経営目線）", size=13, color=COLOR_WARNING, bold=True)
    add_body(s, 0.7, 6.0, 5.6, 0.85,
             "・年間 812 〜 1,215 万円の機会損失\n"
             "・退院集中対策で 50〜70% 実現すれば年 400〜850 万円の改善",
             size=11, color=COLOR_TEXT)

    # 現場負荷（看護目線）
    add_rect(s, 6.8, 5.5, 6.0, 1.4, COLOR_BG_INFO, line_color=COLOR_ACCENT)
    add_body(s, 7.0, 5.6, 5.6, 0.4,
             "🩺 現場負荷（看護目線）", size=13, color=COLOR_ACCENT, bold=True)
    add_body(s, 7.0, 6.0, 5.6, 0.85,
             "・月曜入院ピーク日の看護師負荷が突出\n"
             "・退院処理 + 新規入院対応 + 必要度記録 が同日集中",
             size=11, color=COLOR_TEXT)
    add_footer(s, "Slide 9 / 15  •  出典: data/past_admissions_2025fy.csv")
    add_notes(s,
        "（90秒）退院集中の話です。\n"
        "過去1年で、5F・6F合わせて49回も『1日5名以上の退院』が発生していました。"
        "約2週に1回のペース。これは偶然ではなく構造的な問題です。\n"
        "金土に退院が偏っているのに、土日は新規入院がほぼ来ない。"
        "結果、月曜にまとめて入院が来て、現場負荷が爆発する。\n"
        "経営目線では年間812〜1,215万円の機会損失。"
        "でも今日強調したいのは、これは経営の問題ではなく『現場負荷の問題』だということ。"
        "看護師の月曜のオーバーワークが、必要度評価の精度低下にも繋がる悪循環を生んでいます。"
    )

    # =====================================================================
    # SLIDE 10: このまま 6/1 を迎えると
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_title(s, "💀 このまま 6/1 を迎えると、病棟運営の継続が難しくなる",
              color=COLOR_DANGER)

    add_body(s, 0.5, 1.3, 12.3, 0.5,
             "経過措置（困難月の除外）は 5/31 で終了。例外なし、本則完全適用。",
             size=13, color=COLOR_SUBTEXT)

    # シナリオブロック
    add_rect(s, 0.5, 1.9, 12.3, 1.4, COLOR_BG_DANGER, line_color=COLOR_DANGER)
    add_body(s, 0.7, 2.0, 12.0, 0.45,
             "🚨 ステップ1: 看護必要度未達 → 入院基本料 上位施設基準から脱落",
             size=14, color=COLOR_DANGER, bold=True)
    add_body(s, 0.7, 2.45, 12.0, 0.85,
             "現状の6Fは Ⅰ 1.46pt / Ⅱ 3.38pt 未達（直近3ヶ月では更に悪化）。"
             "rolling 3ヶ月で6月以降 1度でも未達月が連続すれば、地域包括医療病棟入院料の届出取下げのリスク。",
             size=12, color=COLOR_TEXT)

    add_rect(s, 0.5, 3.45, 12.3, 1.4, COLOR_BG_DANGER, line_color=COLOR_DANGER)
    add_body(s, 0.7, 3.55, 12.0, 0.45,
             "🚨 ステップ2: 救急15% / LOS 未達 → 病棟運営継続そのものが困難",
             size=14, color=COLOR_DANGER, bold=True)
    add_body(s, 0.7, 4.0, 12.0, 0.85,
             "施設基準未達が複数指標に及ぶと、月単位で 1,500〜2,500 万円規模の入院料減収。"
             "地域の救急受け皿としての役割を縮小せざるを得ない可能性。",
             size=12, color=COLOR_TEXT)

    # 患者・地域への波及
    add_rect(s, 0.5, 5.0, 12.3, 1.9, COLOR_BG_WARNING, line_color=COLOR_WARNING)
    add_body(s, 0.7, 5.1, 12.0, 0.45,
             "⚠️ 患者さん・地域への波及",
             size=14, color=COLOR_WARNING, bold=True)
    add_body(s, 0.7, 5.55, 12.0, 1.3,
             "・近隣の急性期病院からの転院受入（下り搬送）が縮小 → 急性期病院の在院日数長期化\n"
             "・救急受入縮小 → 地域救急の受け皿が当院から失われる\n"
             "・職員の処遇への影響（賞与・昇給）と、新人教育の余力縮小\n"
             "・最悪、入院料の格下げで病棟機能の見直しを迫られる",
             size=12, color=COLOR_TEXT)
    add_footer(s, "Slide 10 / 15  •  ※ 警告ではなく、6/1以降に実際に起こりうる現実シナリオ")
    add_notes(s,
        "（90秒）ここからは少し重い話です。\n"
        "このまま6月1日を迎えるとどうなるか。\n"
        "まずステップ1: 6Fの看護必要度が未達のままだと、地域包括医療病棟入院料の届出取下げのリスク。\n"
        "ステップ2: 救急15%や LOS が未達になると、月単位で1,500〜2,500万円規模の減収。\n"
        "そして患者さん、地域への波及。当院は地域の急性期病院からの下り搬送の受け皿でもあります。"
        "ここが弱まると、急性期病院の在院日数も延びて、地域全体の医療資源が逼迫します。\n"
        "脅したいのではありません。これは6月1日以降に実際に起こりうる現実シナリオです。"
        "だからこそ、残り35日で運用を変えていきたい。次のスライドから、具体的にお願いしたいことを話します。"
    )

    # =====================================================================
    # SLIDE 11: 看護必要度 — 医師・看護師アクション
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_title(s, "🤝 協力依頼①：看護必要度 — 実施しているケアを正しく数字に残す")

    add_body(s, 0.5, 1.3, 12.3, 0.5,
             "新しい処置を増やすのではない。すでにやっているケアを評価項目の言葉に変換する。",
             size=13, color=COLOR_SUBTEXT)

    # 医師カード
    add_rect(s, 0.5, 1.95, 6.0, 4.5, COLOR_BG_LIGHT, line_color=COLOR_ACCENT)
    add_body(s, 0.7, 2.05, 5.6, 0.5,
             "👨‍⚕️ 医師にお願いしたい3つ", size=16, color=COLOR_ACCENT, bold=True)
    actions_doc = [
        ("①", "Day 3 で退院見通し共有",
         "入院 Day 3 のカンファで、想定退院日と退院時の状態を病棟と共有。"
         "Day 7 以降の退院支援を加速できる。"),
        ("②", "処置・モニタリングの記録",
         "実施したA項目（呼吸ケア・点滴管理など）の処置情報を、看護記録と連動させる。"
         "「やったけど書いていない」をゼロに。"),
        ("③", "C項目（手術・侵襲的処置）の事前共有",
         "予定処置の前日までに病棟へ通知し、評価対象として漏れなく拾えるようにする。"
         "C1以上の処置が分子に直接効く。"),
    ]
    y = 2.55
    for num, head, desc in actions_doc:
        add_body(s, 0.7, y, 0.4, 0.4, num, size=14, color=COLOR_DANGER, bold=True)
        add_body(s, 1.1, y, 5.2, 0.4, head, size=13, color=COLOR_PRIMARY, bold=True)
        add_body(s, 0.9, y + 0.4, 5.4, 0.85, desc, size=11, color=COLOR_TEXT)
        y += 1.3

    # 看護師カード
    add_rect(s, 6.8, 1.95, 6.0, 4.5, COLOR_BG_LIGHT, line_color=COLOR_SUCCESS)
    add_body(s, 7.0, 2.05, 5.6, 0.5,
             "👩‍⚕️ 看護師にお願いしたい3つ", size=16, color=COLOR_SUCCESS, bold=True)
    actions_nrs = [
        ("①", "毎日の必要度記録漏れチェック",
         "勤務終了前の5分で、当日の必要度評価を見直す。"
         "実施したケアが項目に反映されているか確認。"),
        ("②", "A2 / C1 トリガーの即時記録",
         "「呼吸ケア」「点滴ライン3本以上」「術後管理」など、トリガー該当の判断が出た瞬間に記録。"
         "後回しにしない。"),
        ("③", "評価困難な患者を朝礼で共有",
         "「該当か微妙」「医師に確認したい」を朝礼でオープンにし、リーダーが医師と詰める。"
         "個別判断を属人化しない。"),
    ]
    y = 2.55
    for num, head, desc in actions_nrs:
        add_body(s, 7.0, y, 0.4, 0.4, num, size=14, color=COLOR_SUCCESS, bold=True)
        add_body(s, 7.4, y, 5.2, 0.4, head, size=13, color=COLOR_PRIMARY, bold=True)
        add_body(s, 7.2, y + 0.4, 5.4, 0.85, desc, size=11, color=COLOR_TEXT)
        y += 1.3

    add_body(s, 0.5, 6.7, 12.3, 0.4,
             "💡 これらは「初期案」です。現場から修正提案をお願いします。",
             size=13, color=COLOR_ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s, "Slide 11 / 15")
    add_notes(s,
        "（90秒）ここから具体的な協力依頼です。\n"
        "看護必要度については、医師に3つ、看護師に3つお願いがあります。\n"
        "医師には、Day 3 での退院見通し共有、A項目処置の記録連動、C項目処置の事前共有。\n"
        "看護師には、毎日5分の記録漏れチェック、A2/C1トリガー即時記録、判断に迷う患者の朝礼共有。\n"
        "繰り返しますが、これらは『初期案』です。現場で『これは難しい』『こうしたほうが現実的』があれば、"
        "むしろそれを聞かせてください。一緒にアクションを練り直します。"
    )

    # =====================================================================
    # SLIDE 12: 救急15% + LOS 協力依頼
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_title(s, "🤝 協力依頼②：救急15% と LOS — 受入判断と退院支援の協力")

    # 上半分: 救急15%
    add_rect(s, 0.5, 1.4, 12.3, 0.4, COLOR_DANGER)
    add_body(s, 0.7, 1.45, 12.0, 0.35,
             "❷ 救急搬送後 15% を満たすために",
             size=14, color=COLOR_WHITE, bold=True)

    add_body(s, 0.5, 1.95, 6.0, 0.4,
             "👨‍⚕️ 医師", size=14, color=COLOR_ACCENT, bold=True)
    add_body(s, 0.7, 2.4, 6.0, 1.5,
             "・救急車搬送と他院からの下り搬送の受入判断を、病棟稼働だけでなく\n"
             "　月単位の救急15%状況と連動して判断\n"
             "・連携室・救急室と「今月の必要救急枠」を共有",
             size=11, color=COLOR_TEXT)

    add_body(s, 6.8, 1.95, 6.0, 0.4,
             "👩‍⚕️ 看護師", size=14, color=COLOR_SUCCESS, bold=True)
    add_body(s, 7.0, 2.4, 6.0, 1.5,
             "・救急患者の受入準備の標準化（夜間・週末も同質に）\n"
             "・転院搬送患者の情報（救急患者連携搬送料の有無）を医事と共有\n"
             "・「救急車有り」「下り搬送」のフラグ漏れを防ぐ",
             size=11, color=COLOR_TEXT)

    # 下半分: LOS / 短手3
    add_rect(s, 0.5, 3.95, 12.3, 0.4, COLOR_WARNING)
    add_body(s, 0.7, 4.0, 12.0, 0.35,
             "❸ 平均在院日数 20日 を維持するために（短手3 階段関数対策）",
             size=14, color=COLOR_WHITE, bold=True)

    add_body(s, 0.5, 4.5, 6.0, 0.4,
             "👨‍⚕️ 医師", size=14, color=COLOR_ACCENT, bold=True)
    add_body(s, 0.7, 4.95, 6.0, 1.7,
             "・短手3患者の Day 5 到達アラートを毎日確認、Day 6 突破を全力で防ぐ\n"
             "・Day 7 以降の退院支援を加速（薬剤調整・家族面談前倒し）\n"
             "・退院サマリーは退院日に間に合わせる（病棟回転を止めない）",
             size=11, color=COLOR_TEXT)

    add_body(s, 6.8, 4.5, 6.0, 0.4,
             "👩‍⚕️ 看護師", size=14, color=COLOR_SUCCESS, bold=True)
    add_body(s, 7.0, 4.95, 6.0, 1.7,
             "・短手3患者の Day 5 を朝礼で必ず共有、医師に直接確認\n"
             "・退院困難因子（家族・住居・ADL）の早期発見と退院支援NSへの連絡\n"
             "・カンファで退院見通しを能動的に提案する",
             size=11, color=COLOR_TEXT)

    add_body(s, 0.5, 6.8, 12.3, 0.4,
             "💡 アプリの「Day 5 到達アラート」「rolling 90日 LOS」を毎朝の確認ルーチンに",
             size=13, color=COLOR_ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s, "Slide 12 / 15")
    add_notes(s,
        "（90秒）救急15%と平均在院日数の協力依頼です。\n"
        "救急15%については、医師は受入判断を月単位の救急15%状況と連動させてほしい。"
        "看護師は救急患者の受入準備を夜間・週末も同質に保ち、フラグ漏れを防いでほしい。\n"
        "LOSについては、医師には短手3 Day 5 アラートの毎日確認と Day 7以降の退院支援加速を。"
        "看護師には短手3 Day 5 の朝礼共有と退院困難因子の早期発見をお願いします。\n"
        "アプリにはこれら全部を支援する機能があります。Day 5 アラートも rolling 90日 LOS も自動表示されています。"
        "毎朝5分の確認ルーチンに組み込んでいただけると助かります。"
    )

    # =====================================================================
    # SLIDE 13: 退院集中 — 月曜以降への振替
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_title(s, "🤝 協力依頼③：退院集中 — 金+土退院を月曜以降へ振り替える")

    add_body(s, 0.5, 1.3, 12.3, 0.5,
             "現状: 金+土退院 = 35〜36% / 5名超の集中日 = 年49回 / 月曜入院ピーク 7.4件/日",
             size=13, color=COLOR_SUBTEXT)

    # 提案ロジック
    add_rect(s, 0.5, 1.95, 12.3, 1.4, COLOR_BG_INFO, line_color=COLOR_ACCENT)
    add_body(s, 0.7, 2.05, 12.0, 0.4,
             "📐 退院日振替の根拠（過去1年実データ）", size=14, color=COLOR_ACCENT, bold=True)
    add_body(s, 0.7, 2.5, 12.0, 0.85,
             "・金曜退院 1件 → 月曜以降に振替で 土日 2日分 の空床を防止\n"
             "・土曜退院 1件 → 月曜以降に振替で 日曜 1日分 の空床を防止\n"
             "・日曜・祝日退院も 1病棟 2人/日 まで補助枠として活用可（無理なく）",
             size=12, color=COLOR_TEXT)

    # 医師
    add_rect(s, 0.5, 3.5, 6.0, 2.7, COLOR_BG_LIGHT, line_color=COLOR_ACCENT)
    add_body(s, 0.7, 3.6, 5.6, 0.45,
             "👨‍⚕️ 医師にお願い（病棟担当医・主治医）",
             size=14, color=COLOR_ACCENT, bold=True)
    add_body(s, 0.7, 4.05, 5.6, 2.1,
             "・カンファで「金土退院予定」の患者をリストアップし、\n"
             "　月曜以降に延ばせる方を能動的に検討\n"
             "・「家族都合・処置スケジュール」が金土集中の主因なら、\n"
             "　処置・退院前面談を平日後半に組み直す\n"
             "・処置スケジュールの前倒し（金曜の処置を木曜以前に）も視野",
             size=11, color=COLOR_TEXT)

    # 看護師
    add_rect(s, 6.8, 3.5, 6.0, 2.7, COLOR_BG_LIGHT, line_color=COLOR_SUCCESS)
    add_body(s, 7.0, 3.6, 5.6, 0.45,
             "👩‍⚕️ 看護師にお願い（リーダー・退院支援NS）",
             size=14, color=COLOR_SUCCESS, bold=True)
    add_body(s, 7.0, 4.05, 5.6, 2.1,
             "・退院カレンダーで「金土退院集中日」を事前に検出 → 主治医アラート\n"
             "・退院候補リスト（A群→B群への移行）の早期更新\n"
             "・日曜・祝日の退院候補（補助枠 2人/日）を金曜カンファで合意\n"
             "・退院困難因子の早期共有で退院日後ろ倒しの選択肢を増やす",
             size=11, color=COLOR_TEXT)

    add_body(s, 0.5, 6.4, 12.3, 0.4,
             "✅ 効果試算: 50% 実現で 年 400〜570 万円改善 + 月曜の現場負荷ピークが平準化",
             size=13, color=COLOR_SUCCESS, bold=True, align=PP_ALIGN.CENTER)
    add_footer(s, "Slide 13 / 15  •  既存ベッドコントロールアプリで「金+土退院率」可視化済")
    add_notes(s,
        "（90秒）退院集中対策です。\n"
        "ベッドコントロールアプリには『金+土退院率』を可視化する機能があり、Hint 2/3 で What-If シミュレーションも可能です。\n"
        "医師には、カンファで金土退院予定をリストアップし、月曜以降に延ばせる方を能動的に検討してほしい。"
        "処置スケジュールが金集中の主因なら、処置を木曜以前に前倒すという選択肢もあります。\n"
        "看護師には、退院カレンダーで金土集中日を事前に検出し主治医にアラート、"
        "そして日曜・祝日の補助枠2人/日を金曜カンファで合意してほしいです。\n"
        "50%実現できれば年間400〜570万円の改善ですが、それ以上に月曜の現場負荷ピークが平準化される効果が大きい。"
    )

    # =====================================================================
    # SLIDE 14: 達成シナリオ
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_title(s, "🌅 一緒なら達成できる — 達成シナリオと月次振り返り")

    # 段階目標
    add_body(s, 0.5, 1.3, 12.3, 0.4,
             "🎯 1ヶ月目（5月） — 運用変更の定着フェーズ",
             size=15, color=COLOR_PRIMARY, bold=True)
    add_rect(s, 0.5, 1.7, 12.3, 1.0, COLOR_BG_INFO)
    add_body(s, 0.7, 1.8, 12.0, 0.85,
             "・看護必要度: 6F の必要度Ⅰを月次で +1pt 改善（17.5患者日/月の取り戻し）\n"
             "・退院集中: 金+土退院率を 30% 未満に抑える（カンファで月曜振替を試行）\n"
             "・短手3 Day 5 アラート確認を全員のルーチンに",
             size=12, color=COLOR_TEXT)

    add_body(s, 0.5, 2.85, 12.3, 0.4,
             "🎯 3ヶ月目（7月） — rolling 3ヶ月で本則達成",
             size=15, color=COLOR_PRIMARY, bold=True)
    add_rect(s, 0.5, 3.25, 12.3, 1.0, COLOR_BG_SUCCESS)
    add_body(s, 0.7, 3.35, 12.0, 0.85,
             "・看護必要度Ⅰ・Ⅱ: 6F も応需係数加算後で新基準達成（rolling 3ヶ月）\n"
             "・救急15%: 純実データで rolling 3ヶ月計算へ移行（5月までの手動シードを卒業）\n"
             "・退院集中: 月曜入院ピーク日の看護師オーバーワーク日数を半減",
             size=12, color=COLOR_TEXT)

    # 月次振り返り
    add_rect(s, 0.5, 4.4, 12.3, 1.4, COLOR_BG_INFO, line_color=COLOR_ACCENT)
    add_body(s, 0.7, 4.5, 12.0, 0.4,
             "🔄 毎月最終週: 進捗を病棟会議で共有", size=14, color=COLOR_ACCENT, bold=True)
    add_body(s, 0.7, 4.95, 12.0, 0.85,
             "ベッドコントロールアプリで自動算出 →"
             " 看護必要度の月次推移、救急15%の rolling 3ヶ月、退院集中の発生回数 を共有。\n"
             "数字だけでなく、「どの運用変更が効いたか／まだ効いていないか」を全員で議論。",
             size=12, color=COLOR_TEXT)

    # 修正ループ
    add_rect(s, 0.5, 5.95, 12.3, 1.0, COLOR_BG_WARNING, line_color=COLOR_WARNING)
    add_body(s, 0.7, 6.05, 12.0, 0.4,
             "🛠 現場修正提案を歓迎します（責める運用にしない）",
             size=14, color=COLOR_WARNING, bold=True)
    add_body(s, 0.7, 6.5, 12.0, 0.45,
             "本日の提案は『初期案』です。「これは無理」「こうしたほうが現実的」を月次振り返りで遠慮なく出してください。"
             "アクション内容は毎月見直します。",
             size=12, color=COLOR_TEXT)
    add_footer(s, "Slide 14 / 15")
    add_notes(s,
        "（60秒）達成シナリオです。\n"
        "1ヶ月目（5月）は運用変更の定着フェーズ。看護必要度6Fを+1pt 改善、"
        "金+土退院率を30%未満に、短手3 Day 5アラート確認をルーチン化。\n"
        "3ヶ月目（7月）には rolling 3ヶ月で本則達成、救急15%も純実データで判定可能に。\n"
        "毎月最終週に病棟会議で進捗を共有します。アプリで自動算出された数字を見ながら、"
        "『どの運用変更が効いたか』を全員で議論する。\n"
        "そして繰り返しますが、責める運用にしません。修正提案を月次で受け付けます。"
    )

    # =====================================================================
    # SLIDE 15: 一致団結メッセージ
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_rect(s, 0, 0, 13.333, 7.5, COLOR_PRIMARY)

    add_body(s, 0.5, 1.0, 12.3, 1.0,
             "🤝 病棟を守る共同作戦",
             size=44, color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_body(s, 0.5, 2.3, 12.3, 0.7,
             "誰かを責めるための資料ではありません。",
             size=22, color=RGBColor(0xD1, 0xD5, 0xDB), align=PP_ALIGN.CENTER)
    add_body(s, 0.5, 3.0, 12.3, 0.7,
             "地域の高齢救急の受け皿である当院の病棟を、",
             size=22, color=RGBColor(0xD1, 0xD5, 0xDB), align=PP_ALIGN.CENTER)
    add_body(s, 0.5, 3.7, 12.3, 0.7,
             "医師・看護師全員で続けていくための共同作戦です。",
             size=22, color=RGBColor(0xD1, 0xD5, 0xDB), align=PP_ALIGN.CENTER)

    add_rect(s, 1.5, 4.7, 10.3, 1.3, COLOR_ACCENT)
    add_body(s, 1.5, 4.85, 10.3, 0.4,
             "次のアクション", size=16, color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_body(s, 1.5, 5.3, 10.3, 0.6,
             "明日から：協力依頼の3項目を、自分の現場で1つ試す",
             size=18, color=COLOR_WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_body(s, 0.5, 6.2, 12.3, 0.4,
             "毎月最終週に病棟会議で進捗共有 — 月次データで成果を見える化します",
             size=14, color=RGBColor(0xE5, 0xE7, 0xEB), align=PP_ALIGN.CENTER)
    add_body(s, 0.5, 6.7, 12.3, 0.4,
             "質問・修正提案は遠慮なく → 副院長・病棟師長まで",
             size=14, color=RGBColor(0xE5, 0xE7, 0xEB), align=PP_ALIGN.CENTER)

    add_notes(s,
        "（30秒）最後にもう一度。\n"
        "今日のお話は誰かを責めるための資料ではありません。"
        "地域の高齢救急の受け皿である当院の病棟を、医師・看護師全員で続けていくための共同作戦です。\n"
        "明日から、協力依頼の3項目のうち、自分の現場で1つ試してみてください。"
        "完璧でなくていいです。1つから始めましょう。\n"
        "ご質問、修正提案は遠慮なく。これから10分、質疑の時間にしたいと思います。ありがとうございました。"
    )

    return prs


def main() -> int:
    prs = build_presentation()
    out = Path("docs/admin/2026新基準_医師看護師向けブリーフィング_2026-04-26.pptx")
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"✅ 生成完了: {out}")
    print(f"   スライド数: {len(prs.slides)} 枚")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
