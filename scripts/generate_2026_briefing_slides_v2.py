"""2026新基準 医師看護師ブリーフィング v2（ビジュアル中心、20枚）.

副院長フィードバック (2026-04-26): v1は文字過多・小フォント・グラフ不在で
行動変容に繋がらない構成だった。v2では:
- 1枚 1メッセージ徹底
- 主要メッセージは 40-72pt の超大文字
- グラフは matplotlib で生成済み PNG を全画面サイズで埋め込み
- 看護必要度の前提知識から教える（医師は忙しく知らない前提）
- 危機ページは全画面赤、希望ページは全画面緑

使い方:
    .venv/bin/python scripts/generate_2026_briefing_slides_v2.py
    → docs/admin/2026新基準_医師看護師向けブリーフィング_v2_2026-04-26.pptx
"""
from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# テーマカラー
# ---------------------------------------------------------------------------

PRIMARY = RGBColor(0x1F, 0x29, 0x37)      # ダーク
ACCENT = RGBColor(0x25, 0x63, 0xEB)       # 青
SUCCESS = RGBColor(0x10, 0xB9, 0x81)      # 緑
WARNING = RGBColor(0xF5, 0x9E, 0x0B)      # オレンジ
DANGER = RGBColor(0xDC, 0x26, 0x26)       # 赤
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_RED = RGBColor(0x99, 0x1B, 0x1B)
DARK_GREEN = RGBColor(0x06, 0x59, 0x3D)
DARK_BLUE = RGBColor(0x1E, 0x40, 0xAF)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
TEXT_DARK = RGBColor(0x11, 0x18, 0x27)

ASSETS = Path('docs/admin/briefing_2026_assets')
TODAY = date(2026, 4, 26)
TRANSITION_END = date(2026, 5, 31)
DAYS_LEFT = (TRANSITION_END - TODAY).days


# ---------------------------------------------------------------------------
# ヘルパー
# ---------------------------------------------------------------------------

def add_full_bg(slide, color):
    """スライド全面を指定色で塗りつぶす."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0), Inches(0),
                                   Inches(13.333), Inches(7.5))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False


def add_text(slide, left, top, width, height, text,
             size, color, bold=True, align=PP_ALIGN.CENTER):
    """テキストボックス追加（デフォルトは中央揃え・太字）."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color


def add_image_full(slide, image_path, left=0.5, top=1.0, width=12.3, height=6.0):
    """画像を埋め込み（デフォルトはほぼ全画面サイズ）."""
    slide.shapes.add_picture(str(image_path),
                              Inches(left), Inches(top),
                              width=Inches(width), height=Inches(height))


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_page_no(slide, text):
    add_text(slide, 11.5, 7.05, 1.5, 0.3, text, 9,
             RGBColor(0x9C, 0xA3, 0xAF), bold=False, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# プレゼン本体（20枚）
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # =====================================================================
    # SLIDE 1: タイトル
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_full_bg(s, PRIMARY)

    add_text(s, 0.5, 1.0, 12.3, 1.5,
             "あと " + str(DAYS_LEFT) + " 日",
             size=120, color=DANGER, bold=True)
    add_text(s, 0.5, 2.5, 12.3, 0.7,
             "2026年6月1日 新基準 完全適用まで",
             size=24, color=RGBColor(0xD1, 0xD5, 0xDB))

    add_text(s, 0.5, 3.8, 12.3, 1.2,
             "病棟を、守る。",
             size=72, color=WHITE, bold=True)
    add_text(s, 0.5, 5.0, 12.3, 0.8,
             "医師・看護師合同ブリーフィング",
             size=28, color=ACCENT)
    add_text(s, 0.5, 6.5, 12.3, 0.5,
             "おもろまちメディカルセンター  副院長  久保田 徹  /  2026-04-26",
             size=14, color=RGBColor(0xE5, 0xE7, 0xEB), bold=False)
    add_notes(s,
        "（30秒）\n"
        "今日、皆さんに本当に伝えたいことを、20分でお話しします。\n"
        "5月31日まで、あと35日。6月1日からは、当院の地域包括医療病棟は新しい基準で運営されます。"
        "経過措置はもうありません。今までのように『困難月だったから例外』は通じない。\n"
        "結論を先に言います。このままだと、当院の病棟は基準を満たせず、運営継続が難しくなります。"
        "でも、私たち全員で運用を変えれば、達成できます。今日のお話はその共同作戦です。"
    )

    # =====================================================================
    # SLIDE 2: 結論一行
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_full_bg(s, WHITE)

    add_text(s, 0.5, 1.5, 12.3, 1.0,
             "今日の結論",
             size=36, color=PRIMARY)

    add_text(s, 0.5, 2.8, 12.3, 1.5,
             "このままでは、",
             size=44, color=TEXT_DARK)
    add_text(s, 0.5, 4.0, 12.3, 1.2,
             "6月1日に病棟基準が満たせない。",
             size=52, color=DANGER, bold=True)
    add_text(s, 0.5, 5.5, 12.3, 1.2,
             "でも、全員で運用を変えれば達成できる。",
             size=40, color=SUCCESS, bold=True)
    add_page_no(s, "2 / 21")
    add_notes(s,
        "（45秒）\n"
        "結論は1行です。このままでは6月1日に病棟基準が満たせない。"
        "でも、全員で運用を変えれば達成できる。\n"
        "今日のお話は、医師の先生方、看護師の皆さんに『こうしてほしい』をお願いするものです。"
        "後で具体的な3つずつのアクションをお見せします。\n"
        "そして大事なのは、これは初期案であって、現場からの修正提案を歓迎します。"
        "1人で決めて押し付ける運用は、続かないので。"
    )

    # =====================================================================
    # SLIDE 3: 4つのハードル
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_full_bg(s, WHITE)

    add_text(s, 0.5, 0.5, 12.3, 0.8,
             "6月1日から、4つのハードルを同時にクリアし続ける",
             size=28, color=PRIMARY)

    # 4つの大きな円をピル型で
    items = [
        ("看護必要度", "Ⅰ19% / Ⅱ18%", DANGER, 0.5),
        ("救急搬送後", "15%", WARNING, 3.7),
        ("平均在院日数", "20日", WARNING, 6.9),
        ("退院集中", "現場負荷", ACCENT, 10.1),
    ]
    for label, val, color, x in items:
        # 円形カード
        rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x), Inches(2.0),
                                   Inches(2.7), Inches(3.5))
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.fill.background()
        rect.shadow.inherit = False
        # ラベル
        add_text(s, x, 2.4, 2.7, 0.7, label, size=22, color=WHITE)
        # 値
        add_text(s, x, 3.4, 2.7, 1.5, val, size=36, color=WHITE, bold=True)

    add_text(s, 0.5, 5.9, 12.3, 0.8,
             "1つでも未達 → 病棟入院料の届出取下げリスク",
             size=28, color=DANGER, bold=True)
    add_text(s, 0.5, 6.6, 12.3, 0.5,
             "（4つは互いに引っ張り合う。単一指標で頑張ると別が悪化する）",
             size=16, color=RGBColor(0x6B, 0x72, 0x80), bold=False)
    add_page_no(s, "3 / 21")
    add_notes(s,
        "（45秒）\n"
        "6月1日から、私たちが守らないといけないハードルは4つです。\n"
        "看護必要度Ⅰ19%とⅡ18%、救急搬送後の患者割合15%、平均在院日数20日以下、"
        "そして退院集中による現場負荷の問題。\n"
        "1つでも未達になると、地域包括医療病棟入院料の届出取下げリスクがあります。\n"
        "そして厄介なのが、4つは互いに引っ張り合うこと。在院日数を短くすれば退院増で集中、"
        "稼働率上げるために延ばせば必要度が下がる。だから単一指標で頑張る運用じゃなく、"
        "4つを同時に成り立たせる運用が必要です。"
    )

    # =====================================================================
    # SLIDE 4: 看護必要度 A項目（正確な令和8改定版）
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_full_bg(s, WHITE)
    add_text(s, 0.5, 0.4, 12.3, 0.6,
             "まず、看護必要度 A項目 を正確に共有します",
             size=24, color=PRIMARY)
    add_image_full(s, ASSETS / 'g7_what_is_necessity.png',
                   left=0.3, top=1.0, width=12.7, height=6.3)
    add_page_no(s, "4 / 21")
    add_notes(s,
        "（120秒）\n"
        "医師の先生方の中には看護必要度をご存知ない方もいらっしゃるので、令和8改定版を正確に共有します。\n"
        "A項目は7つ。A1 創傷処置（褥瘡除く）— 蜂窩織炎・糖尿病性足病変・術後縫合創。"
        "A2 呼吸ケア — ネブライザー・体位ドレナージ。喀痰吸引のみは除外。\n"
        "A3 注射薬剤3種類以上の管理 — 令和8改定で旧『点滴ライン3本以上』から名称・概念が変わった項目です。"
        "最大7日間、3種類以上の注射薬剤を投与していれば該当。\n"
        "A4 シリンジポンプ、A5 輸血、A7 救急搬送後 — これらは比較的明確。\n"
        "そして A6 専門的な治療・処置の中には3点項目があり、内科で意識すれば取れます。次のスライドで詳しく。\n"
        "判定は『A2点以上 OR C1点以上』。令和8改定で『AND』から『OR』に変わったので、"
        "A項目だけで2点取れれば該当判定です。"
    )

    # =====================================================================
    # SLIDE 4.5 (新): 内科で意識すれば取れる 3点項目
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_full_bg(s, WHITE)
    add_text(s, 0.5, 0.4, 12.3, 0.6,
             "内科で「意識すれば即取れる」3 点項目（A6 のうち4つ）",
             size=24, color=DANGER)
    add_image_full(s, ASSETS / 'g9_a6_3pt.png',
                   left=0.3, top=1.0, width=12.7, height=6.3)
    add_page_no(s, "5 / 21")
    add_notes(s,
        "（120秒）\n"
        "ここが今日のキースライドです。A6 の中で 3点（最高得点）の項目が6つあり、"
        "内科系で取りやすいのは下記4つ。\n"
        "A6③ 麻薬注射 — 急性膵炎、胆石症、がん性疼痛で経口移行を急ぎすぎず、注射が継続するうちはしっかり該当判定。\n"
        "A6⑦ 昇圧剤注射 — 敗血症、心不全急性増悪。短時間でも該当します。記録漏れが多い項目。\n"
        "A6⑧ 抗不整脈剤注射 — アミオダロン、リドカイン点滴。ペイン科のPHN治療も対象になりうる。\n"
        "A6⑨ 抗血栓塞栓薬持続点滴 — ヘパリン持続。『持続』がキーで、単発静注は対象外。\n"
        "医師が回診時に『この患者、これに該当する？』を意識するだけで、該当割合は変わります。"
        "後ほどの協力依頼で『どう運用に組み込むか』をお話しします。"
    )

    # =====================================================================
    # SLIDE 5: 看護必要度Ⅰ — 6F が未達
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_full_bg(s, WHITE)
    add_text(s, 0.5, 0.4, 12.3, 0.7,
             "看護必要度Ⅰ：6F病棟は、新基準に届かない",
             size=28, color=DANGER)
    add_image_full(s, ASSETS / 'g1_necessity_I.png',
                   left=0.5, top=1.2, width=12.3, height=5.7)
    add_page_no(s, "6 / 21")
    add_notes(s,
        "（60秒）\n"
        "これが当院の看護必要度Ⅰの現状です。グラフは応需係数（救急患者を多く受け入れている病院に与えられる加算）を加算した後の値。\n"
        "5Fは19.7%で達成見込み。でも6Fは17.5%で、新基準19%まで1.5pt 不足です。"
        "これが2025年度1年分の平均。直近の状況はもっと厳しいことを次のスライドでお見せします。"
    )

    # =====================================================================
    # SLIDE 6: 看護必要度Ⅱ — もっと深刻
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_full_bg(s, WHITE)
    add_text(s, 0.5, 0.4, 12.3, 0.7,
             "看護必要度Ⅱ：6Fはさらに深刻、5Fもギリギリ",
             size=28, color=DANGER)
    add_image_full(s, ASSETS / 'g2_necessity_II.png',
                   left=0.5, top=1.2, width=12.3, height=5.7)
    add_page_no(s, "7 / 21")
    add_notes(s,
        "（60秒）\n"
        "必要度Ⅱはもっと深刻です。"
        "5Fは17.85%で新基準18%まで0.15pt と紙一重。"
        "6Fは14.62% で、新基準まで3.38pt も不足しています。\n"
        "つまり、5Fは『あと一歩』、6Fは『現状の延長線では到底届かない』。"
        "そして、これは1年平均。直近を見るとさらに悪化しています。"
    )

    # =====================================================================
    # SLIDE 7: 直近3ヶ月の崖
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_full_bg(s, WHITE)
    add_text(s, 0.5, 0.4, 12.3, 0.7,
             "6F病棟は、直近3ヶ月で崖のように悪化",
             size=28, color=DANGER)
    add_image_full(s, ASSETS / 'g3_6f_trend_cliff.png',
                   left=0.5, top=1.2, width=12.3, height=5.7)
    add_page_no(s, "8 / 21")
    add_notes(s,
        "（90秒）\n"
        "これが6F病棟の月次推移です。直近3ヶ月、つまり1月から3月にかけて、"
        "必要度Ⅰ・Ⅱの両方が崖のように落ちています。\n"
        "Ⅱは9.11% まで落ちて、新基準18% まで8.89pt 不足。"
        "rolling 3ヶ月評価が始まる6月1日に、このまま当てはめれば未達確定の数字です。\n"
        "原因は複合的です。冬季の入院構成変化、在院日数の延長による分母の膨らみ、"
        "そして、実施したケアが評価項目として正しく記録されていない可能性があります。\n"
        "後者は、私たちが運用で取り戻せる部分です。"
    )

    # =====================================================================
    # SLIDE 8: 救急15% って何？
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_full_bg(s, WHITE)
    add_text(s, 0.5, 0.4, 12.3, 0.7,
             "次に、救急搬送後15% を 30秒で説明します",
             size=28, color=PRIMARY)
    add_image_full(s, ASSETS / 'g8_what_is_emergency15.png',
                   left=0.5, top=1.2, width=12.3, height=5.7)
    add_page_no(s, "9 / 21")
    add_notes(s,
        "（60秒）\n"
        "救急搬送後15%は、分子と分母で考えます。\n"
        "分子は『救急車で当院に直接入院した患者』と『他院で救急患者連携搬送料を算定して当院に転院してきた下り搬送患者』の合計。"
        "分母は入院患者全員（短手3も含む）。これが15%以上を rolling 3ヶ月で病棟別に維持。\n"
        "つまり下り搬送は独立した別基準ではなく、救急15%の分子の一部として地域急性期病院との連携で稼ぐルートです。"
    )

    # =====================================================================
    # SLIDE 9: 短手3 階段関数
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_full_bg(s, WHITE)
    add_text(s, 0.5, 0.4, 12.3, 0.7,
             "在院日数：短手3 患者が Day 6 を超えた瞬間、LOSが+6日ジャンプ",
             size=24, color=DANGER)
    add_image_full(s, ASSETS / 'g5_short3_stair.png',
                   left=0.5, top=1.2, width=12.3, height=5.7)
    add_page_no(s, "10 / 21")
    add_notes(s,
        "（90秒）\n"
        "もう1つ、見落としがちな短手3 患者の特殊ルールです。\n"
        "短手3とは白内障やポリペクなど予定された短期入院。"
        "Day 5までに退院すれば分母に含めない。つまり病棟LOSを押し下げる味方です。\n"
        "ところがDay 6 以降に滞在が延びた瞬間、入院初日まで遡って全日数を分母に計上。"
        "1人のDay 6 突破で病棟LOSが+6日ジャンプします。\n"
        "アプリにはDay 5 到達アラート機能があります。回診・カンファで必ずこれを全員で確認してください。"
    )

    # =====================================================================
    # SLIDE 10: 退院集中 — 49回
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_full_bg(s, WHITE)
    add_text(s, 0.5, 0.4, 12.3, 0.7,
             "退院集中：年間49回（約2週に1回）も発生している",
             size=28, color=DANGER)
    add_image_full(s, ASSETS / 'g6_discharge_overflow.png',
                   left=0.5, top=1.2, width=12.3, height=5.7)
    add_page_no(s, "11 / 21")
    add_notes(s,
        "（45秒）\n"
        "退院集中の話に移ります。\n"
        "1日5名以上の退院が発生した日は、過去1年で5F 22日、6F 27日、合計49日。"
        "約2週に1回のペースで起きています。これは偶然ではなく構造的な問題です。"
    )

    # =====================================================================
    # SLIDE 11: 退院は金土に偏る
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_full_bg(s, WHITE)
    add_text(s, 0.5, 0.4, 12.3, 0.7,
             "退院は金土に偏り、土日の入院補充はほぼゼロ",
             size=26, color=DANGER)
    add_image_full(s, ASSETS / 'g4_discharge_admission.png',
                   left=0.5, top=1.2, width=12.3, height=5.7)
    add_page_no(s, "12 / 21")
    add_notes(s,
        "（90秒）\n"
        "なぜ退院集中が起きるか。グラフが因果を示しています。\n"
        "退院は金土に集中。一方、入院は土曜2.2件/日、日曜0.3件/日と、補充がほぼ来ません。"
        "結果、土日に空床が滞留し、月曜にまとめて入院ピークが来る。\n"
        "つまり、毎週『金土退院 → 土日空床 → 月曜入院ピーク』のサイクル。"
        "月曜の現場負荷が突出し、看護師オーバーワーク、必要度評価精度の低下にも繋がります。"
    )

    # =====================================================================
    # SLIDE 12: このままだとこうなる（全画面赤）
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_full_bg(s, DARK_RED)

    add_text(s, 0.5, 0.8, 12.3, 1.2,
             "このまま 6月1日 を迎えると",
             size=44, color=WHITE)

    add_text(s, 0.5, 2.5, 12.3, 1.2,
             "病棟入院料の", size=36, color=RGBColor(0xFE, 0xCA, 0xCA))
    add_text(s, 0.5, 3.4, 12.3, 1.5,
             "届出取下げ",
             size=88, color=WHITE, bold=True)

    add_text(s, 0.5, 5.1, 12.3, 0.8,
             "→ 月 1,500〜2,500 万円規模の減収",
             size=30, color=WHITE)
    add_text(s, 0.5, 5.9, 12.3, 0.7,
             "→ 地域救急の受け皿が、当院から失われる",
             size=26, color=RGBColor(0xFE, 0xCA, 0xCA))
    add_text(s, 0.5, 6.7, 12.3, 0.5,
             "脅しではありません。本則完全適用の現実です。",
             size=18, color=RGBColor(0xFE, 0xCA, 0xCA), bold=False)
    add_notes(s,
        "（60秒）\n"
        "ここからは少し重い話です。\n"
        "このまま6月1日を迎えると、地域包括医療病棟入院料の届出取下げのリスクがあります。"
        "月単位で1,500〜2,500万円規模の減収。\n"
        "でもそれよりも大きいのは、当院は地域の高齢救急の受け皿だということ。"
        "ここが弱まると、急性期病院の在院日数も延びて、地域全体の医療資源が逼迫します。\n"
        "脅しではありません。これは本則完全適用の現実シナリオ。だからこそ、残り35日で運用を変えていきたい。"
    )

    # =====================================================================
    # SLIDE 13: でも達成できる（全画面緑への転換）
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_full_bg(s, DARK_GREEN)

    add_text(s, 0.5, 1.8, 12.3, 1.5,
             "でも、",
             size=72, color=RGBColor(0xA7, 0xF3, 0xD0))
    add_text(s, 0.5, 3.0, 12.3, 1.7,
             "全員でやれば、",
             size=64, color=WHITE, bold=True)
    add_text(s, 0.5, 4.6, 12.3, 1.7,
             "達成できる。",
             size=84, color=WHITE, bold=True)
    add_page_no(s, "14 / 21")
    add_notes(s,
        "（30秒）\n"
        "でも、全員でやれば達成できます。\n"
        "次のスライドから、医師の先生方、看護師の皆さんにお願いしたいことを具体的にお見せします。\n"
        "完璧でなくていい。まず1つ、明日から自分の現場で試してください。"
    )

    # =====================================================================
    # SLIDE 14: 医師にお願い 3つ（看護必要度）
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_full_bg(s, WHITE)
    add_text(s, 0.5, 0.5, 12.3, 0.7,
             "医師の先生方にお願いしたい 3 つ",
             size=28, color=ACCENT)

    actions = [
        ("①", "Day 3 で退院見通しを共有",
         "入院 Day 3 のカンファで想定退院日を病棟と共有。Day 7 以降の退院支援を前倒しできる。"),
        ("②", "実施したA項目を看護記録に連動",
         "呼吸ケア・点滴管理など、やったことを記録に残す。「やったけど書いてない」をゼロに。"),
        ("③", "C項目（処置・手術）の事前共有",
         "予定処置を前日までに病棟へ通知。C1以上の処置が必要度の分子に直接効く。"),
    ]
    y = 1.7
    for num, head, desc in actions:
        # 円い番号
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(0.7), Inches(y),
                                     Inches(1.0), Inches(1.0))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT
        circle.line.fill.background()
        circle.shadow.inherit = False
        add_text(s, 0.7, y + 0.18, 1.0, 0.7, num, size=44, color=WHITE)
        # ヘッドライン
        add_text(s, 1.95, y, 10.5, 0.6, head, size=26, color=PRIMARY,
                 align=PP_ALIGN.LEFT)
        # 説明
        add_text(s, 1.95, y + 0.65, 10.5, 0.7, desc, size=15,
                 color=RGBColor(0x4B, 0x55, 0x63), bold=False, align=PP_ALIGN.LEFT)
        y += 1.7

    add_page_no(s, "15 / 21")
    add_notes(s,
        "（90秒）医師の先生方にお願いしたい3つ。\n"
        "①Day 3で退院見通し共有 — 入院3日目のカンファで想定退院日を病棟と共有してください。"
        "Day 7以降の退院支援を前倒しできます。\n"
        "②A項目を看護記録に連動 — 呼吸ケア、点滴管理など、実施したことを記録に残してください。"
        "『やったけど書いてない』をゼロに。\n"
        "③C項目の事前共有 — 予定処置を前日までに病棟へ通知。C1以上の処置は必要度の分子に直接効きます。"
    )

    # =====================================================================
    # SLIDE 15: 看護師にお願い 3つ
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_full_bg(s, WHITE)
    add_text(s, 0.5, 0.5, 12.3, 0.7,
             "看護師の皆さんにお願いしたい 3 つ",
             size=28, color=SUCCESS)

    actions = [
        ("①", "勤務終了前の 5 分で記録漏れチェック",
         "当日の必要度評価を見直す。実施したケアが評価項目に反映されているか確認する。"),
        ("②", "A2 / C1 トリガーを即時記録",
         "「呼吸ケア」「点滴ライン3本以上」「術後管理」など、該当判断が出た瞬間に記録。後回しにしない。"),
        ("③", "判断に迷う患者を朝礼で共有",
         "「該当か微妙」「医師に確認したい」を朝礼でオープンに。リーダーが医師と詰める。属人化を防ぐ。"),
    ]
    y = 1.7
    for num, head, desc in actions:
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(0.7), Inches(y),
                                     Inches(1.0), Inches(1.0))
        circle.fill.solid()
        circle.fill.fore_color.rgb = SUCCESS
        circle.line.fill.background()
        circle.shadow.inherit = False
        add_text(s, 0.7, y + 0.18, 1.0, 0.7, num, size=44, color=WHITE)
        add_text(s, 1.95, y, 10.5, 0.6, head, size=26, color=PRIMARY,
                 align=PP_ALIGN.LEFT)
        add_text(s, 1.95, y + 0.65, 10.5, 0.7, desc, size=15,
                 color=RGBColor(0x4B, 0x55, 0x63), bold=False, align=PP_ALIGN.LEFT)
        y += 1.7

    add_page_no(s, "16 / 21")
    add_notes(s,
        "（90秒）看護師の皆さんにお願いしたい3つ。\n"
        "①勤務終了前の5分で記録漏れチェック — 当日の必要度評価を見直してください。"
        "実施したケアが評価項目に反映されているか確認。\n"
        "②A2 / C1 トリガーの即時記録 — 該当判断が出た瞬間に記録。後回しにしない。\n"
        "③判断に迷う患者を朝礼で共有 — 『該当か微妙』をオープンにして、リーダーが医師と詰める。属人化を防ぐ。"
    )

    # =====================================================================
    # SLIDE 16: 救急15% + LOS 共通アクション
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_full_bg(s, WHITE)
    add_text(s, 0.5, 0.4, 12.3, 0.7,
             "救急15% と 在院日数20日 を守るための合言葉",
             size=26, color=PRIMARY)

    # 2 ブロック
    add_text(s, 0.5, 1.5, 12.3, 0.7,
             "救急15%", size=28, color=WARNING)

    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.5), Inches(2.3),
                               Inches(12.3), Inches(1.6))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    rect.line.color.rgb = WARNING
    rect.line.width = Pt(2)
    rect.shadow.inherit = False

    add_text(s, 0.7, 2.5, 11.9, 0.7,
             "「救急車も下り搬送も、月の目標から逆算して受ける」",
             size=24, color=PRIMARY, align=PP_ALIGN.LEFT)
    add_text(s, 0.7, 3.2, 11.9, 0.6,
             "連携室・救急室と「今月あと何件必要か」を朝礼で共有",
             size=16, color=RGBColor(0x4B, 0x55, 0x63),
             bold=False, align=PP_ALIGN.LEFT)

    add_text(s, 0.5, 4.2, 12.3, 0.7,
             "在院日数20日（短手3 階段関数）", size=28, color=WARNING)

    rect2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.5), Inches(5.0),
                                Inches(12.3), Inches(1.6))
    rect2.fill.solid()
    rect2.fill.fore_color.rgb = RGBColor(0xFE, 0xF3, 0xC7)
    rect2.line.color.rgb = WARNING
    rect2.line.width = Pt(2)
    rect2.shadow.inherit = False

    add_text(s, 0.7, 5.2, 11.9, 0.7,
             "「短手3 患者の Day 5 を、毎朝 全員で確認する」",
             size=24, color=PRIMARY, align=PP_ALIGN.LEFT)
    add_text(s, 0.7, 5.9, 11.9, 0.6,
             "アプリの Day 5 アラートを朝礼でチェック → Day 6 突破ゼロを目指す",
             size=16, color=RGBColor(0x4B, 0x55, 0x63),
             bold=False, align=PP_ALIGN.LEFT)
    add_page_no(s, "17 / 21")
    add_notes(s,
        "（60秒）救急15%とLOSは、それぞれ1つの合言葉でまとめます。\n"
        "救急15%は『救急車も下り搬送も、月の目標から逆算して受ける』。"
        "連携室・救急室と『今月あと何件必要か』を朝礼で共有することが鍵です。\n"
        "LOS は『短手3 患者の Day 5 を、毎朝 全員で確認する』。"
        "アプリにDay 5 アラート機能があるので、毎朝のチェックをルーチンに。"
    )

    # =====================================================================
    # SLIDE 17: 退院集中対策
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_full_bg(s, WHITE)
    add_text(s, 0.5, 0.4, 12.3, 0.7,
             "退院集中：金+土退院を、月曜以降に振り替える",
             size=26, color=PRIMARY)

    # 大きな提案
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.5), Inches(1.4),
                               Inches(12.3), Inches(2.2))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xDB, 0xEA, 0xFE)
    rect.line.color.rgb = ACCENT
    rect.line.width = Pt(3)
    rect.shadow.inherit = False
    add_text(s, 0.7, 1.6, 11.9, 0.7,
             "「金土に退院予定の患者、月曜まで延ばせない？」",
             size=28, color=ACCENT)
    add_text(s, 0.7, 2.4, 11.9, 1.1,
             "これを毎週の退院調整カンファで全員で問い直す。\n"
             "日曜・祝日も 1 病棟あたり 2 人/日 まで補助枠として活用可。",
             size=18, color=PRIMARY, bold=False)

    # 効果
    add_text(s, 0.5, 4.0, 12.3, 0.6,
             "効果", size=22, color=SUCCESS)

    rect2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.5), Inches(4.7),
                                Inches(12.3), Inches(2.2))
    rect2.fill.solid()
    rect2.fill.fore_color.rgb = RGBColor(0xD1, 0xFA, 0xE5)
    rect2.line.color.rgb = SUCCESS
    rect2.line.width = Pt(3)
    rect2.shadow.inherit = False
    add_text(s, 0.7, 4.9, 11.9, 0.7,
             "月曜入院ピーク日の現場負荷を平準化",
             size=24, color=DARK_GREEN, align=PP_ALIGN.LEFT)
    add_text(s, 0.7, 5.7, 11.9, 1.1,
             "・看護師の月曜オーバーワークを抑える → 必要度評価の精度向上\n"
             "・年間 400〜570 万円の機会損失を回収（50% 実現で）",
             size=16, color=PRIMARY, bold=False, align=PP_ALIGN.LEFT)
    add_page_no(s, "18 / 21")
    add_notes(s,
        "（60秒）退院集中対策です。\n"
        "提案は1つ。『金土に退院予定の患者、月曜まで延ばせない？』を、毎週の退院調整カンファで全員で問い直す。\n"
        "日曜祝日も1病棟2人/日までは補助枠として活用できます。\n"
        "効果は2つ。月曜の現場負荷が平準化されて、看護師オーバーワークが減る。"
        "結果、必要度評価の精度も上がります。経営的には年間400〜570万円の機会損失回収。"
    )

    # =====================================================================
    # SLIDE 18: 達成シナリオ
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_full_bg(s, WHITE)
    add_text(s, 0.5, 0.5, 12.3, 0.7,
             "全員で取り組めば、3ヶ月でこうなる",
             size=28, color=SUCCESS)

    # 2段階のロードマップ
    # 1ヶ月目
    rect = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.5), Inches(1.7),
                               Inches(12.3), Inches(2.3))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xDB, 0xEA, 0xFE)
    rect.line.color.rgb = ACCENT
    rect.line.width = Pt(2)
    rect.shadow.inherit = False
    add_text(s, 0.7, 1.9, 11.9, 0.6,
             "5月（1ヶ月目）", size=24, color=ACCENT, align=PP_ALIGN.LEFT)
    add_text(s, 0.7, 2.6, 11.9, 1.3,
             "・必要度: 6F の必要度Ⅰを 月 +1pt 改善\n"
             "・退院集中: 金+土退院率を 30% 未満に抑える\n"
             "・短手3 Day 5 アラート確認をルーチンに",
             size=18, color=PRIMARY, bold=False, align=PP_ALIGN.LEFT)

    # 3ヶ月目
    rect2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.5), Inches(4.3),
                                Inches(12.3), Inches(2.3))
    rect2.fill.solid()
    rect2.fill.fore_color.rgb = RGBColor(0xD1, 0xFA, 0xE5)
    rect2.line.color.rgb = SUCCESS
    rect2.line.width = Pt(2)
    rect2.shadow.inherit = False
    add_text(s, 0.7, 4.5, 11.9, 0.6,
             "7月（3ヶ月目）— rolling 3ヶ月で本則達成",
             size=24, color=DARK_GREEN, align=PP_ALIGN.LEFT)
    add_text(s, 0.7, 5.2, 11.9, 1.3,
             "・必要度Ⅰ・Ⅱ: 6F も応需係数加算後で新基準達成\n"
             "・救急15%: 純実データで rolling 3ヶ月計算へ移行\n"
             "・退院集中: 月曜入院ピーク日の看護師オーバーワーク日数を半減",
             size=18, color=PRIMARY, bold=False, align=PP_ALIGN.LEFT)

    add_text(s, 0.5, 6.8, 12.3, 0.4,
             "毎月最終週に病棟会議で進捗共有 — 数字を見ながら、運用を見直す",
             size=14, color=RGBColor(0x6B, 0x72, 0x80), bold=False)
    add_page_no(s, "19 / 21")
    add_notes(s,
        "（45秒）達成シナリオです。\n"
        "5月は運用変更の定着フェーズ。必要度6Fを+1pt 改善、金+土退院率を30%未満に、Day 5 アラート確認をルーチン化。\n"
        "7月にはrolling 3ヶ月で本則達成。救急15%も純実データで判定可能になる。\n"
        "毎月最終週に病棟会議で進捗を共有して、運用を見直していきます。"
    )

    # =====================================================================
    # SLIDE 19: 修正提案歓迎
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_full_bg(s, WHITE)

    add_text(s, 0.5, 1.5, 12.3, 1.0,
             "今日の提案は",
             size=44, color=PRIMARY)
    add_text(s, 0.5, 2.7, 12.3, 1.5,
             "「初期案」",
             size=88, color=ACCENT, bold=True)
    add_text(s, 0.5, 4.5, 12.3, 1.0,
             "「これは無理」「こうしたほうが現実的」",
             size=28, color=PRIMARY)
    add_text(s, 0.5, 5.5, 12.3, 1.0,
             "現場の修正提案を歓迎します。",
             size=32, color=SUCCESS, bold=True)
    add_text(s, 0.5, 6.7, 12.3, 0.5,
             "毎月の振り返りで、運用ルールは更新していきます",
             size=16, color=RGBColor(0x6B, 0x72, 0x80), bold=False)
    add_page_no(s, "20 / 21")
    add_notes(s,
        "（45秒）大事な点をもう一度。\n"
        "今日の提案は『初期案』です。『これは無理』『こうしたほうが現実的』があれば、"
        "むしろそれを聞かせてください。\n"
        "毎月の振り返り会議で、運用ルールを更新していきます。"
        "全員が当事者として作戦を作り直していく前提です。"
    )

    # =====================================================================
    # SLIDE 20: クロージング
    # =====================================================================
    s = prs.slides.add_slide(blank)
    add_full_bg(s, PRIMARY)

    add_text(s, 0.5, 1.0, 12.3, 1.5,
             "病棟を、",
             size=72, color=WHITE)
    add_text(s, 0.5, 2.5, 12.3, 1.5,
             "守ろう。",
             size=120, color=SUCCESS, bold=True)

    add_text(s, 0.5, 4.5, 12.3, 0.7,
             "明日から、自分の現場で 1 つだけ試してください。",
             size=24, color=WHITE)
    add_text(s, 0.5, 5.4, 12.3, 0.7,
             "完璧でなくていい。1つから始めましょう。",
             size=22, color=RGBColor(0xD1, 0xD5, 0xDB))

    add_text(s, 0.5, 6.7, 12.3, 0.5,
             "ご質問・修正提案は、副院長・病棟師長まで",
             size=14, color=RGBColor(0x9C, 0xA3, 0xAF), bold=False)
    add_notes(s,
        "（30秒）最後に。\n"
        "病棟を、守りましょう。\n"
        "明日から、自分の現場で1つだけ試してください。完璧でなくていい。1つから始めましょう。\n"
        "ご質問、修正提案、いつでも遠慮なく副院長・病棟師長まで。"
        "これから10分間、質疑応答の時間にします。ありがとうございました。"
    )

    return prs


def main():
    prs = build()
    out = Path('docs/admin/2026新基準_医師看護師向けブリーフィング_v2_2026-04-26.pptx')
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"✅ 生成完了: {out}")
    print(f"   スライド数: {len(prs.slides)}")
    print(f"   サイズ: {out.stat().st_size / 1024:.1f} KB")
    return 0


if __name__ == '__main__':
    raise SystemExit(main())
