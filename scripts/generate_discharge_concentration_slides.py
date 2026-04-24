"""退院集中対策スライド（経営会議・師長会議向け、10 枚）生成スクリプト.

副院長の「ベッドコントローラー（病棟主任・入退院調整NS・退院支援SW）に
退院集中の重要性を納得させたい」要望に応えるための、稼働率と診療報酬の
観点からストーリー立てた 10 枚のスライドを python-pptx で生成する。

使い方:
    python3 scripts/generate_discharge_concentration_slides.py
    → docs/admin/退院集中対策_2026-04-24.pptx を出力

データソース（本スクリプト実行時点の数字を使用）:
    - past_admissions_2025fy.csv (1823 件、5F/6F 2025-04〜2026-04)
    - slot_policy_impact_analysis.py の計算結果（ハードコード）
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# テーマカラー（design_tokens と揃える）
# ---------------------------------------------------------------------------

COLOR_PRIMARY = RGBColor(0x37, 0x41, 0x51)    # ダークグレー
COLOR_ACCENT = RGBColor(0x25, 0x63, 0xEB)     # 青
COLOR_SUCCESS = RGBColor(0x10, 0xB9, 0x81)    # 緑
COLOR_WARNING = RGBColor(0xF5, 0x9E, 0x0B)    # オレンジ
COLOR_DANGER = RGBColor(0xDC, 0x26, 0x26)     # 赤
COLOR_TEXT = RGBColor(0x1F, 0x29, 0x37)       # 本文テキスト
COLOR_SUBTEXT = RGBColor(0x6B, 0x72, 0x80)    # サブテキスト
COLOR_BG_LIGHT = RGBColor(0xF9, 0xFA, 0xFB)   # 薄いグレー背景


# ---------------------------------------------------------------------------
# ヘルパー関数
# ---------------------------------------------------------------------------

def add_title(slide, text: str, color=COLOR_PRIMARY, size: int = 32) -> None:
    """スライドにタイトルを追加（上部にテキストボックス）."""
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(12.3)
    height = Inches(0.9)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = color


def add_body(slide, left_in, top_in, width_in, height_in, text: str,
             size: int = 16, color=COLOR_TEXT, bold: bool = False,
             align=PP_ALIGN.LEFT) -> None:
    """本文テキストボックスを追加."""
    tb = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.alignment = align
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color


def add_rect(slide, left_in, top_in, width_in, height_in,
             fill_color, line_color=None) -> None:
    """背景装飾の矩形."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False


def add_metric_card(slide, left_in, top_in, width_in, label: str, value: str,
                    note: str = "", accent_color=COLOR_ACCENT) -> None:
    """KPI カード（数値 + ラベル + 注記）."""
    height_in = 1.5
    # 背景
    add_rect(slide, left_in, top_in, width_in, height_in,
             COLOR_BG_LIGHT, line_color=accent_color)
    # ラベル
    add_body(slide, left_in + 0.15, top_in + 0.08,
             width_in - 0.3, 0.35,
             label, size=12, color=COLOR_SUBTEXT)
    # 数値（大きく）
    add_body(slide, left_in + 0.15, top_in + 0.45,
             width_in - 0.3, 0.7,
             value, size=28, color=accent_color, bold=True)
    # 注記
    if note:
        add_body(slide, left_in + 0.15, top_in + 1.1,
                 width_in - 0.3, 0.35,
                 note, size=10, color=COLOR_SUBTEXT)


def add_footer(slide, text: str) -> None:
    """各スライド下部の footer."""
    add_body(slide, 0.5, 7.0, 12.3, 0.3,
             text, size=10, color=COLOR_SUBTEXT)


# ---------------------------------------------------------------------------
# スライド生成
# ---------------------------------------------------------------------------

def build_presentation() -> Presentation:
    prs = Presentation()
    # 16:9 wide screen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    # ---- Slide 1: タイトル ----
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.333, 7.5, COLOR_PRIMARY)
    add_rect(slide, 0, 5.5, 13.333, 2.0, COLOR_ACCENT)
    add_body(slide, 0.5, 2.2, 12.3, 1.5,
             "退院集中という隠れた損失",
             size=54, color=RGBColor(0xFF, 0xFF, 0xFF),
             bold=True, align=PP_ALIGN.CENTER)
    add_body(slide, 0.5, 3.8, 12.3, 0.9,
             "— 実データが示す構造的問題と、枠管理による対策 —",
             size=24, color=RGBColor(0xD1, 0xD5, 0xDB),
             align=PP_ALIGN.CENTER)
    add_body(slide, 0.5, 5.8, 12.3, 0.5,
             "おもろまちメディカルセンター  副院長  久保田 徹",
             size=18, color=RGBColor(0xFF, 0xFF, 0xFF),
             align=PP_ALIGN.CENTER)
    add_body(slide, 0.5, 6.5, 12.3, 0.5,
             "2026-04-24  病棟師長会議 / ベッドコントロール説明資料",
             size=14, color=RGBColor(0xE5, 0xE7, 0xEB),
             align=PP_ALIGN.CENTER)

    # ---- Slide 2: 問題提起 ----
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "🚨 今日のカンファで起きたこと")
    add_body(slide, 0.5, 1.4, 12.3, 0.6,
             "2026-04-24（今日）— 1 日で退院 8 名が重なった",
             size=20, color=COLOR_DANGER, bold=True)
    add_rect(slide, 0.5, 2.1, 12.3, 2.3, COLOR_BG_LIGHT)
    add_body(slide, 0.8, 2.3, 11.7, 0.5,
             "病棟の声:", size=16, color=COLOR_SUBTEXT, bold=True)
    add_body(slide, 1.0, 2.8, 11.5, 1.5,
             "・「業務がパンクしている」病棟看護師\n"
             "・「稼働率が一気に下がる」経営的リスク\n"
             "・「コントローラー不在で、度々こういうことが起きる」",
             size=16, color=COLOR_TEXT)
    add_body(slide, 0.5, 4.7, 12.3, 0.6,
             "これは偶然の一度きり？それとも構造的な問題？",
             size=22, color=COLOR_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_body(slide, 0.5, 5.6, 12.3, 0.8,
             "→ 過去 1 年の実データ（1,793 件）を分析して検証しました",
             size=18, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    add_footer(slide, "Slide 2 / 10")

    # ---- Slide 3: データで見る ----
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "📊 実データが示す「退院集中」の構造")
    add_body(slide, 0.5, 1.3, 12.3, 0.5,
             "過去 1 年（2025-04〜2026-04、5F+6F 計 1,793 件）の曜日別 枠超過日数",
             size=14, color=COLOR_SUBTEXT)

    # 5F
    add_body(slide, 0.5, 2.0, 6.0, 0.4,
             "5F 病棟（22 日 / 382 日 = 5.8%）",
             size=16, color=COLOR_PRIMARY, bold=True)
    dow_5f = [("月", 2), ("火", 4), ("水", 4), ("木", 2),
              ("金", 2), ("土", 8), ("日", 0)]
    y = 2.5
    for dow, n in dow_5f:
        add_body(slide, 0.6, y, 0.4, 0.3, dow,
                 size=14, color=COLOR_TEXT)
        bar_w = max(0.05, n * 0.5)
        color = COLOR_DANGER if n >= 5 else (COLOR_WARNING if n >= 3 else COLOR_SUCCESS)
        add_rect(slide, 1.1, y + 0.03, bar_w, 0.25, color)
        add_body(slide, 1.2 + bar_w, y, 1.0, 0.3, f"{n} 日",
                 size=14, color=COLOR_TEXT, bold=True)
        y += 0.38

    # 6F
    add_body(slide, 7.0, 2.0, 6.0, 0.4,
             "6F 病棟（27 日 / 387 日 = 7.0%）",
             size=16, color=COLOR_PRIMARY, bold=True)
    dow_6f = [("月", 1), ("火", 8), ("水", 4), ("木", 4),
              ("金", 8), ("土", 2), ("日", 0)]
    y = 2.5
    for dow, n in dow_6f:
        add_body(slide, 7.1, y, 0.4, 0.3, dow,
                 size=14, color=COLOR_TEXT)
        bar_w = max(0.05, n * 0.5)
        color = COLOR_DANGER if n >= 5 else (COLOR_WARNING if n >= 3 else COLOR_SUCCESS)
        add_rect(slide, 7.6, y + 0.03, bar_w, 0.25, color)
        add_body(slide, 7.7 + bar_w, y, 1.0, 0.3, f"{n} 日",
                 size=14, color=COLOR_TEXT, bold=True)
        y += 0.38

    add_rect(slide, 0.5, 5.5, 12.3, 1.5, COLOR_BG_LIGHT)
    add_body(slide, 0.8, 5.6, 11.7, 0.5,
             "見えてきた病棟の個性（偶然ではなく構造的）",
             size=16, color=COLOR_PRIMARY, bold=True)
    add_body(slide, 0.8, 6.1, 11.7, 0.8,
             "・5F（外科・整形中心）: 土曜に退院が集中（8 日 = 36%）\n"
             "・6F（内科中心）: 火曜・金曜に集中（各 8 日 = 合計 59%）",
             size=14, color=COLOR_TEXT)
    add_footer(slide, "Slide 3 / 10")

    # ---- Slide 4: 頻度 ----
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "📅 「2 週間に 1 回」起きている現象")

    # 大きな数字
    add_metric_card(slide, 0.5, 1.5, 3.8,
                    "年間枠超過日数", "49 日",
                    "5F 22 日 + 6F 27 日", accent_color=COLOR_DANGER)
    add_metric_card(slide, 4.6, 1.5, 3.8,
                    "平均発生頻度", "2 週に 1 回",
                    "5F 約 17 日に 1 回 / 6F 約 14 日に 1 回",
                    accent_color=COLOR_WARNING)
    add_metric_card(slide, 8.7, 1.5, 3.8,
                    "最大退院数 / 日", "9 名",
                    "6F 2025-12-30 (火) 記録", accent_color=COLOR_DANGER)

    add_body(slide, 0.5, 3.5, 12.3, 0.5,
             "直近の 8 名退院事件は「最大級」ではなく「よくある重なり」",
             size=18, color=COLOR_PRIMARY, bold=True, align=PP_ALIGN.CENTER)

    add_rect(slide, 0.5, 4.3, 12.3, 2.2, COLOR_BG_LIGHT)
    add_body(slide, 0.8, 4.4, 11.7, 0.5,
             "超過が大きい日 TOP 5（年間）",
             size=14, color=COLOR_SUBTEXT, bold=True)
    lines = [
        "🔴 2025-12-30 (火) 6F   退院 9 名 / 枠 5 → +4 超過",
        "🔴 2026-01-27 (火) 6F   退院 9 名 / 枠 5 → +4 超過",
        "🟠 2025-05-29 (木) 5F   退院 8 名 / 枠 5 → +3 超過",
        "🟠 2025-07-19 (土) 5F   退院 8 名 / 枠 5 → +3 超過",
        "🟠 2026-02-17 (火) 6F   退院 8 名 / 枠 5 → +3 超過",
    ]
    add_body(slide, 0.8, 4.9, 11.7, 1.5,
             "\n".join(lines), size=13, color=COLOR_TEXT)
    add_footer(slide, "Slide 4 / 10")

    # ---- Slide 5: 稼働率への影響 ----
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "📉 退院集中が稼働率に与える影響")

    add_body(slide, 0.5, 1.3, 12.3, 0.6,
             "退院が一日に集中する → 翌日の在院数ガクッと減少 → 数日かけて元に戻る",
             size=16, color=COLOR_TEXT)

    # 視覚化
    add_rect(slide, 0.5, 2.1, 12.3, 3.3, COLOR_BG_LIGHT)
    add_body(slide, 0.8, 2.3, 11.7, 0.4,
             "【パターン A】枠制限なし（実績）",
             size=14, color=COLOR_DANGER, bold=True)
    add_body(slide, 0.8, 2.7, 11.7, 0.6,
             "月  : 稼働率 92% ▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬ (通常)\n"
             "火  : 稼働率 80% ▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬ (8 名退院で急落)\n"
             "水  : 稼働率 82% ▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬\n"
             "木  : 稼働率 84% ▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬\n"
             "金  : 稼働率 87% ▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬ (徐々に回復)",
             size=12, color=COLOR_TEXT)

    add_body(slide, 0.8, 4.2, 11.7, 0.4,
             "【パターン B】枠制限あり（現仕様）",
             size=14, color=COLOR_SUCCESS, bold=True)
    add_body(slide, 0.8, 4.6, 11.7, 0.6,
             "月  : 稼働率 92% ▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬\n"
             "火  : 稼働率 85% ▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬ (5 名で抑制、3 名翌日へ)\n"
             "水  : 稼働率 85% ▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬ (横ばい、急落回避)\n"
             "木  : 稼働率 89% ▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬\n"
             "金  : 稼働率 91% ▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬▬ (順調に回復)",
             size=12, color=COLOR_TEXT)

    add_body(slide, 0.5, 5.9, 12.3, 0.5,
             "「累積は同じ」でも「平均稼働率」は改善",
             size=18, color=COLOR_ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_body(slide, 0.5, 6.5, 12.3, 0.4,
             "目標レンジ（90-95%）から外れる日数が減る → 月次平均が押し上がる",
             size=14, color=COLOR_SUBTEXT, align=PP_ALIGN.CENTER)
    add_footer(slide, "Slide 5 / 10")

    # ---- Slide 6: 診療報酬への換算 ----
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "💰 診療報酬換算：年間 812 万円の機会損失")

    add_body(slide, 0.5, 1.3, 12.3, 0.5,
             "シミュレーションで算出した「退院集中による機会損失」",
             size=14, color=COLOR_SUBTEXT)

    # 計算式
    add_rect(slide, 0.5, 2.0, 12.3, 1.8, COLOR_BG_LIGHT)
    add_body(slide, 0.8, 2.1, 11.7, 0.4,
             "計算の前提",
             size=14, color=COLOR_SUBTEXT, bold=True)
    add_body(slide, 0.8, 2.5, 11.7, 1.2,
             "・地域包括医療病棟 入院料 (イ, 緊急×無手術) = 3,367 点/日 = 33,670 円/日\n"
             "・平均稼働率改善幅: 5F +0.51 pt, 6F +0.90 pt（過去 1 年実データベース）\n"
             "・年間延べ床日増 = 改善幅 × 病棟床数 47 × 365 日",
             size=12, color=COLOR_TEXT)

    # 結果
    add_metric_card(slide, 0.5, 4.2, 3.8,
                    "5F 改善効果", "+ 293 万円/年",
                    "稼働率 +0.51 pt × 床日換算",
                    accent_color=COLOR_SUCCESS)
    add_metric_card(slide, 4.6, 4.2, 3.8,
                    "6F 改善効果", "+ 519 万円/年",
                    "稼働率 +0.90 pt × 床日換算",
                    accent_color=COLOR_SUCCESS)
    add_metric_card(slide, 8.7, 4.2, 3.8,
                    "合計経営改善効果", "+ 812 万円/年",
                    "常勤医師 1 名の手取りに相当",
                    accent_color=COLOR_ACCENT)

    add_body(slide, 0.5, 6.2, 12.3, 0.5,
             "💡 現実の運用では 50-70% 実現できれば 400-570 万円/年",
             size=16, color=COLOR_TEXT, align=PP_ALIGN.CENTER)
    add_footer(slide, "Slide 6 / 10")

    # ---- Slide 7: 現仕様の効果 ----
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "✅ 枠制限 UI の導入効果（シミュレーション結果）")

    add_body(slide, 0.5, 1.3, 12.3, 0.5,
             "実データに「月〜土 5 枠 / 日祝 2 枠 / 超過時は翌日 -N」のルールを当てはめると",
             size=14, color=COLOR_SUBTEXT)

    # 3 カラム: Before / After / 理論上限
    for i, (label, before_val, after_val, note, accent) in enumerate([
        ("5F 枠超過日数", "22 日", "0 日", "介入率 100%", COLOR_SUCCESS),
        ("6F 枠超過日数", "27 日", "0 日", "介入率 100%", COLOR_SUCCESS),
        ("総退院数（保存）", "変化なし", "同じ", "患者の追い出しではない", COLOR_ACCENT),
    ]):
        x = 0.5 + i * 4.3
        add_rect(slide, x, 2.1, 4.0, 1.8, COLOR_BG_LIGHT, line_color=accent)
        add_body(slide, x + 0.2, 2.2, 3.6, 0.3,
                 label, size=12, color=COLOR_SUBTEXT, bold=True)
        add_body(slide, x + 0.2, 2.5, 3.6, 0.4,
                 f"Before: {before_val}", size=14, color=COLOR_DANGER)
        add_body(slide, x + 0.2, 2.9, 3.6, 0.4,
                 f"After:  {after_val}", size=14, color=COLOR_SUCCESS, bold=True)
        add_body(slide, x + 0.2, 3.4, 3.6, 0.4,
                 note, size=11, color=COLOR_SUBTEXT)

    add_body(slide, 0.5, 4.4, 12.3, 0.5,
             "枠制限 UI は「過去 1 年の枠超過日 49 回すべてを未然に防げた」水準の効果",
             size=16, color=COLOR_PRIMARY, bold=True, align=PP_ALIGN.CENTER)

    add_rect(slide, 0.5, 5.2, 12.3, 1.7, COLOR_BG_LIGHT)
    add_body(slide, 0.8, 5.3, 11.7, 0.4,
             "現仕様で実装済みの機能",
             size=14, color=COLOR_SUBTEXT, bold=True)
    add_body(slide, 0.8, 5.7, 11.7, 1.1,
             "✓ 月〜土 5 枠 / 日祝 2 枠の基本枠ルール（単一ソースで管理）\n"
             "✓ 超過時に翌営業日の枠を -N 自動縮小（連続累積なし）\n"
             "✓ 日曜退院を推奨表示（⭐ マーカー、枠 2 の短いバー）\n"
             "✓ 動的枠調整（稼働率 95% 以上で +1 枠）\n"
             "✓ 主治医別突発退院の頻度集計（運用改善向け）",
             size=12, color=COLOR_TEXT)
    add_footer(slide, "Slide 7 / 10")

    # ---- Slide 8: 現場の負荷軽減 ----
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "🏥 現場の負荷軽減：稼働率以外のメリット")

    benefits = [
        ("病棟業務の平準化", "退院処理 8 名 → 5 名で退院手続き・書類・ベッドメイク・送迎が分散",
         COLOR_ACCENT),
        ("病棟師長の裁量確保", "スタッフ配置の急遽変更が不要、計画通りのシフト運用",
         COLOR_SUCCESS),
        ("入退院調整 NS の業務リズム", "週の後半に負荷が集中しない、カンファの判断材料が見える化",
         COLOR_ACCENT),
        ("経営部の予見性", "稼働率の急落を事前予測 → 投資判断・人員配置に反映",
         COLOR_SUCCESS),
        ("副院長（ベッドコントローラー）の不在対応",
         "アプリが判断材料を提示 → 不在時も事故なく運営", COLOR_WARNING),
    ]

    y = 1.4
    for title, desc, color in benefits:
        add_rect(slide, 0.5, y, 12.3, 0.95, COLOR_BG_LIGHT, line_color=color)
        add_body(slide, 0.7, y + 0.1, 3.5, 0.4,
                 "● " + title, size=14, color=color, bold=True)
        add_body(slide, 4.3, y + 0.1, 8.3, 0.8,
                 desc, size=12, color=COLOR_TEXT)
        y += 1.05
    add_footer(slide, "Slide 8 / 10")

    # ---- Slide 9: 運用フロー提案 ----
    slide = prs.slides.add_slide(blank_layout)
    add_title(slide, "🎯 ベッドコントローラー向け運用フロー（提案）")

    add_body(slide, 0.5, 1.3, 12.3, 0.5,
             "現仕様＋追加策を組み合わせた、1 週間の運用サイクル",
             size=14, color=COLOR_SUBTEXT)

    # フロー 5 ステップ
    steps = [
        ("水曜夜",
         "カンファ前「混雑警告」サマリー自動生成",
         "→ 翌週の超過リスク日 TOP 3 が副院長にプッシュ",
         COLOR_ACCENT),
        ("木曜カンファ",
         "病棟ごとに退院予定を確定",
         "→ 枠超過時は即座に別日候補を検討、動かせない理由もタグで管理",
         COLOR_PRIMARY),
        ("木曜午後〜金曜",
         "カレンダー上で確認・微調整",
         "→ 日曜推奨枠を活用、超過予想日の早期分散",
         COLOR_SUCCESS),
        ("退院当日",
         "突発退院（主治医独断）の監視",
         "→ ⚡ マーカーで集計、月次レポートで主治医別頻度を共有",
         COLOR_WARNING),
        ("月初",
         "月次レポート自動生成",
         "→ 経営会議で枠制限の効果と稼働率改善を報告",
         COLOR_ACCENT),
    ]

    y = 1.9
    for i, (when, what, result, color) in enumerate(steps):
        add_rect(slide, 0.5, y, 1.3, 0.8, color)
        add_body(slide, 0.6, y + 0.1, 1.1, 0.6,
                 when, size=13, color=RGBColor(0xFF, 0xFF, 0xFF),
                 bold=True, align=PP_ALIGN.CENTER)
        add_rect(slide, 2.0, y, 10.8, 0.8, COLOR_BG_LIGHT)
        add_body(slide, 2.2, y + 0.05, 10.4, 0.4,
                 what, size=13, color=COLOR_PRIMARY, bold=True)
        add_body(slide, 2.2, y + 0.45, 10.4, 0.4,
                 result, size=11, color=COLOR_TEXT)
        y += 0.95
    add_footer(slide, "Slide 9 / 10")

    # ---- Slide 10: まとめ ----
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.333, 7.5, COLOR_BG_LIGHT)
    add_title(slide, "💎 まとめ: 退院の「日」を選ぶだけで、病院は変わる")

    # 3 つのキーメッセージ
    msgs = [
        ("見えない損失を可視化",
         "退院集中は 2 週間に 1 回、\n年間 49 回起きている構造的問題",
         COLOR_DANGER),
        ("数字で意思決定",
         "稼働率 +0.5〜0.9 pt、\n年間 812 万円の改善余地",
         COLOR_SUCCESS),
        ("全員で支える",
         "病棟・調整 NS・副院長が\n同じ画面で判断できる仕組み",
         COLOR_ACCENT),
    ]
    for i, (title, body, color) in enumerate(msgs):
        x = 0.5 + i * 4.3
        add_rect(slide, x, 1.5, 4.0, 2.6, RGBColor(0xFF, 0xFF, 0xFF),
                 line_color=color)
        add_body(slide, x + 0.2, 1.7, 3.6, 0.5,
                 title, size=16, color=color, bold=True)
        add_body(slide, x + 0.2, 2.4, 3.6, 1.5,
                 body, size=14, color=COLOR_TEXT)

    # 最後のメッセージ
    add_rect(slide, 0.5, 4.5, 12.3, 1.8, COLOR_PRIMARY)
    add_body(slide, 0.5, 4.8, 12.3, 0.6,
             "ベッドコントローラーは、特別な誰かではない。",
             size=22, color=RGBColor(0xFF, 0xFF, 0xFF),
             bold=True, align=PP_ALIGN.CENTER)
    add_body(slide, 0.5, 5.5, 12.3, 0.6,
             "全員がカレンダーを見て、数字で判断する文化があれば、",
             size=18, color=RGBColor(0xD1, 0xD5, 0xDB),
             align=PP_ALIGN.CENTER)
    add_body(slide, 0.5, 6.0, 12.3, 0.6,
             "不在のコントローラーを超える運営が、アプリと一緒に実現できる。",
             size=18, color=RGBColor(0xD1, 0xD5, 0xDB),
             align=PP_ALIGN.CENTER)

    add_body(slide, 0.5, 6.8, 12.3, 0.4,
             "2026-04-24  おもろまちメディカルセンター  副院長 久保田 徹",
             size=11, color=COLOR_SUBTEXT, align=PP_ALIGN.CENTER)

    return prs


def main() -> None:
    out_dir = Path(__file__).resolve().parent.parent / "docs" / "admin"
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / "退院集中対策_2026-04-24.pptx"

    prs = build_presentation()
    prs.save(str(out_path))

    print(f"✅ 生成完了: {out_path}")
    print(f"   スライド数: {len(prs.slides)}")
    print(f"   サイズ: {out_path.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    main()
