"""医師別 病棟運営パターン分析 — 理事会・運営会議 プレゼン PPTX 生成（4 層分析 抜本改訂版）.

副院長指示 (2026-05-04 / Rev. 4):
  - Codex 批判的レビューを受け、ランキング型 → 4 層分析へ全面再構成
  - スライド 2 を「本日の結論」から「本資料の見方（4 層 framework）」へ置換
  - Layer 2/3/4 のスライド群を新規追加
"""
from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from generate_doctor_kpi_charts import (
    OUT_DIR as CHART_DIR,
    PROFIT_A, PROFIT_B, PROFIT_C,
    PROFIT_A_LOW, PROFIT_A_HIGH,
    PROFIT_B_LOW, PROFIT_B_HIGH,
    PROFIT_C_LOW, PROFIT_C_HIGH,
    REV_A, REV_B, REV_C,
    COST_A, COST_B, COST_C,
    COST_UNCERTAINTY,
    EMERGENCY_RATIO_THRESHOLD,
    HOME_DISCHARGE_THRESHOLD,
    TYPE_ORDER,
    TYPE_DESCRIPTIONS,
    load_and_compute,
    make_anonymous_map,
)

REPO_ROOT = Path(__file__).resolve().parent.parent
JP_FONT = "Hiragino Mincho ProN"

C_TITLE = RGBColor(0x1F, 0x29, 0x37)
C_ACCENT = RGBColor(0x25, 0x63, 0xEB)
C_DANGER = RGBColor(0xDC, 0x26, 0x26)
C_OK = RGBColor(0x10, 0xB9, 0x81)
C_GOLD = RGBColor(0xD9, 0x77, 0x06)
C_MUTED = RGBColor(0x6B, 0x72, 0x80)
C_FAINT = RGBColor(0x9C, 0xA3, 0xAF)
C_BG_BLUE = RGBColor(0xEF, 0xF6, 0xFF)
C_BG_RED = RGBColor(0xFE, 0xF2, 0xF2)
C_BG_GREEN = RGBColor(0xEC, 0xFD, 0xF5)
C_BG_GOLD = RGBColor(0xFE, 0xF3, 0xC7)


def _set_font(run, size_pt=18, bold=False, color=None):
    run.font.name = JP_FONT
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, bold=False,
                 color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_font(run, size_pt=size, bold=bold, color=color)
    return tb


def add_multi_para_textbox(slide, left, top, width, height, lines, default_size=16,
                           color=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            text, size, bold = line + (False,) * (3 - len(line))
        else:
            text, size, bold = line, default_size, False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        _set_font(run, size_pt=size, bold=bold, color=color or C_TITLE)
    return tb


def add_bg_rect(slide, left, top, width, height, fill_color):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_image_to_slide(slide, image_path, left, top, width, height=None):
    if height:
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    else:
        slide.shapes.add_picture(str(image_path), left, top, width=width)


def add_section_header(slide, title, bg=None, fg=None, size=28):
    """ヘッダーバー + タイトル."""
    add_bg_rect(slide, 0, 0, Inches(13.333), Inches(1.0), bg or C_BG_BLUE)
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7),
                title, size=size, bold=True, color=fg or C_TITLE)


def build(mode: str):
    df, agg, ward_total_days, total_hospital_days = load_and_compute()
    anon_map = make_anonymous_map(agg)
    label_map = (lambda x: x) if mode == 'named' else (lambda x: anon_map.get(x, x))
    title_audience = "理事会" if mode == 'named' else "運営会議"

    okuk_label = label_map('OKUK')
    sub_5f = agg[agg['5F_寄与率'] > 0].sort_values('5F_寄与率', ascending=False)
    sub_6f = agg[agg['6F_寄与率'] > 0].sort_values('6F_寄与率', ascending=False)
    overall_surgery = agg['手術あり件数'].sum() / agg['入院件数'].sum() * 100
    overall_emerg = agg['救急搬送後件数'].sum() / agg['入院件数'].sum() * 100
    overall_home = agg['在宅復帰件数'].sum() / agg['入院件数'].sum() * 100
    overall_unp = agg['予定外件数'].sum() / agg['入院件数'].sum() * 100
    type_counts = {t: int((agg['タイプ分類'] == t).sum()) for t in TYPE_ORDER}

    # 分母情報（CSV 全体 vs 分析対象）
    total_admissions = len(df)
    analysis_admissions = int(agg['入院件数'].sum())
    excluded_admissions = total_admissions - analysis_admissions

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank_layout = prs.slide_layouts[6]

    # ========================================================
    # Slide 1: 表紙
    # ========================================================
    s = prs.slides.add_slide(blank_layout)
    add_bg_rect(s, 0, 0, SW, SH, RGBColor(0xF9, 0xFA, 0xFB))
    add_textbox(s, Inches(1), Inches(1.5), Inches(11.5), Inches(1.2),
                "医師別 病棟運営パターン分析",
                size=42, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(2.7), Inches(11.5), Inches(0.7),
                f"{title_audience}説明資料",
                size=26, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.5),
                "病床利用 + 出来高上乗せ可能性 + 制度適合 + 経営上の解釈 の 4 層分析",
                size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.5),
                "（医師個人の評価ではなく、地域包括医療病棟の運営構造を理解するための資料）",
                size=14, color=C_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(5.5), Inches(11.5), Inches(0.6),
                "おもろまちメディカルセンター",
                size=18, color=C_TITLE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(6.0), Inches(11.5), Inches(0.5),
                "副院長 久保田 透",
                size=14, color=C_TITLE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1), Inches(6.5), Inches(11.5), Inches(0.5),
                "2026 年 5 月 4 日 / Rev. 4（4 層分析 抜本改訂版）",
                size=12, color=C_FAINT, align=PP_ALIGN.CENTER)
    if mode == 'anonymous':
        add_textbox(s, Inches(1), Inches(6.95), Inches(11.5), Inches(0.4),
                    "※ 個人名はすべて A 医師, B 医師 ... と匿名化",
                    size=11, color=C_DANGER, align=PP_ALIGN.CENTER)

    # ========================================================
    # Slide 2: 本資料の見方（4 層 framework）— 旧「結論スライド」を置換
    # ========================================================
    s = prs.slides.add_slide(blank_layout)
    add_section_header(s, "本資料の見方 — 4 層分析の枠組み", bg=C_BG_GOLD, size=28)

    body = [
        ("従来「病床日数」と「ベッド粗利単価」だけで医師別貢献を語ると、", 16, False),
        ("手術・リハ・処置等の出来高貢献や、施設基準への適合貢献が抜け落ちます。", 16, False),
        ("そのため本資料は次の 4 層に分けて読み解きます。", 16, True),
        ("", 8, False),
        ("Layer 1: 病床利用貢献", 22, True),
        ("    稼働率寄与率・病床日数シェア・ベッド粗利単価仮定（管理用スコア）", 14, False),
        ("", 6, False),
        ("Layer 2: 出来高上乗せ可能性", 22, True),
        ("    手術あり症例比率（入口指標）／手術料・麻酔料・リハ単位は本資料未反映", 14, False),
        ("", 6, False),
        ("Layer 3: 制度適合貢献", 22, True),
        ("    救急車入院比率（参考 proxy）・在宅復帰率・予定外比率（施設基準と対比）", 14, False),
        ("", 6, False),
        ("Layer 4: 経営上の解釈", 22, True),
        ("    5+1 類型による医師運営パターン分類（順位ではなくタイプ）", 14, False),
        ("", 10, False),
        ("📌 大前提: 医師個人の優劣判定ではなく、病棟運営の構造把握が目的です", 16, True),
    ]
    add_multi_para_textbox(s, Inches(0.7), Inches(1.2), Inches(12.0), Inches(6.0), body)

    # ========================================================
    # Slide 3: 目的・大前提
    # ========================================================
    s = prs.slides.add_slide(blank_layout)
    add_section_header(s, "なぜ「医師別 病棟運営パターン」を見える化するのか")

    body = [
        ("● 病棟稼働率は施設基準（地域包括医療病棟は 90%）達成のための最重要指標", 18, True),
        ("", 8, False),
        ("● 実態は各主治医の入院・在院・退院判断の総和が病棟稼働率を作る", 18, False),
        ("", 16, False),
        ("本資料の問い:", 20, True),
        ("    1. 各医師が病棟稼働率をどれだけ支えているか？（Layer 1）", 17, False),
        ("    2. 出来高上乗せの可能性をどれだけ持つ症例構成か？（Layer 2）", 17, False),
        ("    3. 施設基準維持・在宅復帰にどう貢献しているか？（Layer 3）", 17, False),
        ("    4. 病棟をどのような形で支えているのか？（Layer 4）", 17, False),
    ]
    add_multi_para_textbox(s, Inches(0.7), Inches(1.4), Inches(12.0), Inches(5.0), body)
    add_bg_rect(s, Inches(0.7), Inches(6.4), Inches(12.0), Inches(0.8), C_BG_RED)
    add_textbox(s, Inches(0.9), Inches(6.5), Inches(11.6), Inches(0.6),
                "📌 大前提: 個別医師の評価ではなく、病棟運営の構造把握を目的とします",
                size=14, bold=True, color=C_DANGER)

    # ========================================================
    # Slide 4: データ概要・限界
    # ========================================================
    s = prs.slides.add_slide(blank_layout)
    add_section_header(s, "データ概要・本資料の限界（必読）")

    box_y = Inches(1.3); box_h = Inches(1.8); box_w = Inches(2.9); box_gap = Inches(0.2)
    box_x_start = Inches(0.5)
    boxes = [
        ("対象期間", "13 ヶ月", "2025-04 〜 2026-04"),
        ("CSV 全体", f"{total_admissions:,} 件", "事務提供（医師コード化）"),
        ("分析対象",
         f"{len(agg)} 医師 / {analysis_admissions:,} 件",
         f"件数 10 件以上（{excluded_admissions} 件除外）"),
        ("総延日数", f"{int(agg['総延日数'].sum()):,}", "床日（5F + 6F）"),
    ]
    for i, (label, value, sub) in enumerate(boxes):
        x = box_x_start + (box_w + box_gap) * i
        add_bg_rect(s, x, box_y, box_w, box_h, C_BG_BLUE)
        add_textbox(s, x, box_y + Inches(0.1), box_w, Inches(0.4), label,
                    size=12, color=C_MUTED, align=PP_ALIGN.CENTER)
        add_textbox(s, x, box_y + Inches(0.5), box_w, Inches(0.9), value,
                    size=28, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(s, x, box_y + Inches(1.4), box_w, Inches(0.4), sub,
                    size=11, color=C_MUTED, align=PP_ALIGN.CENTER)

    # 限界の集約表示
    add_bg_rect(s, Inches(0.5), Inches(3.4), Inches(12.3), Inches(3.7), C_BG_RED)
    add_multi_para_textbox(s, Inches(0.7), Inches(3.5), Inches(12.0), Inches(3.5), [
        ("⚠️ 本資料の限界 — 以下は本資料未反映", 16, True),
        ("", 6, False),
        ("【未反映 ① コスト精緻化】 ベッド粗利単価のコストは業界目安 ±20%、当院 DPC 検証未実施", 13, False),
        ("【未反映 ② 出来高点数】 手術料・麻酔料・リハ単位・処置加算・検査加算・薬剤料は CSV 未収録", 13, False),
        ("【未反映 ③ 看護必要度】 病棟単位では別管理、医師別集計は本資料に含まず", 13, False),
        ("", 8, False),
        ("→ Layer 1（管理用スコア）は確定的会計利益ではなく、Layer 2/3 と併読すべき", 14, True),
        ("→ Stage 2 で当院 DPC 等を活用した社内分析を段階的に整備予定", 14, False),
    ], color=C_DANGER)

    # ========================================================
    # Slide 5: Layer 1 概念（粗利フェーズ）
    # ========================================================
    s = prs.slides.add_slide(blank_layout)
    add_section_header(s, "Layer 1: 病床利用貢献 — 粗利仮定の考え方", size=24)
    add_textbox(s, Inches(0.7), Inches(1.2), Inches(12.0), Inches(0.5),
                "アプリ実装の 2026 年度プリセット（入院料 1）に基づく管理用スコア（コスト ±20%）:",
                size=14, color=C_MUTED)

    phase_y = Inches(1.9); phase_h = Inches(3.5); phase_w = Inches(3.9); phase_gap = Inches(0.3)
    phase_x_start = Inches(0.5)
    phases = [
        ("A 群（急性期）", "1 〜 5 日",
         f"報酬 {REV_A:,}", f"変動費 {COST_A:,}", f"粗利仮定 {PROFIT_A:,}",
         RGBColor(0xFE, 0xCA, 0xCA)),
        ("B 群（回復期）", "6 〜 14 日",
         f"報酬 {REV_B:,}", f"変動費 {COST_B:,}", f"粗利仮定 {PROFIT_B:,}",
         RGBColor(0xBB, 0xF7, 0xD0)),
        ("C 群（退院準備期）", "15 日以降",
         f"報酬 {REV_C:,}", f"変動費 {COST_C:,}", f"粗利仮定 {PROFIT_C:,}",
         RGBColor(0xFD, 0xE0, 0x8A)),
    ]
    for i, (label, days, rev, cost, profit, color) in enumerate(phases):
        x = phase_x_start + (phase_w + phase_gap) * i
        add_bg_rect(s, x, phase_y, phase_w, phase_h, color)
        add_textbox(s, x, phase_y + Inches(0.15), phase_w, Inches(0.5), label,
                    size=20, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER)
        add_textbox(s, x, phase_y + Inches(0.7), phase_w, Inches(0.4), days,
                    size=14, color=C_MUTED, align=PP_ALIGN.CENTER)
        add_textbox(s, x, phase_y + Inches(1.2), phase_w, Inches(0.5), rev,
                    size=16, color=C_TITLE, align=PP_ALIGN.CENTER)
        add_textbox(s, x, phase_y + Inches(1.7), phase_w, Inches(0.5), "− " + cost,
                    size=16, color=C_DANGER, align=PP_ALIGN.CENTER)
        add_textbox(s, x, phase_y + Inches(2.3), phase_w, Inches(0.6), "= " + profit,
                    size=20, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    add_bg_rect(s, Inches(0.7), Inches(5.7), Inches(12.0), Inches(1.5), C_BG_GOLD)
    add_multi_para_textbox(s, Inches(0.9), Inches(5.85), Inches(11.6), Inches(1.3), [
        ("⚠️ この粗利仮定は管理用スコア — 会計上の確定利益ではない", 14, True),
        ("    Layer 1 だけでは出来高（手術等）と制度適合は捕捉できない", 13, False),
        ("    Layer 2/3/4 を併読して初めて貢献の姿が立体的に見える", 13, False),
    ], color=C_TITLE)

    # ========================================================
    # Slide 6: Layer 1 ポジショニング（散布図）
    # ========================================================
    s = prs.slides.add_slide(blank_layout)
    add_section_header(s, "Layer 1: 医師別 ポジショニング（散布図）")
    add_image_to_slide(s, CHART_DIR / f"02_scatter_{mode}.png",
                       Inches(0.7), Inches(1.2), Inches(11.9))

    # ========================================================
    # Slide 7: Layer 1 病棟別ランキング
    # ========================================================
    s = prs.slides.add_slide(blank_layout)
    add_section_header(s, "Layer 1: 病棟別 稼働率寄与率")
    add_image_to_slide(s, CHART_DIR / f"03_ward_ranking_{mode}.png",
                       Inches(0.5), Inches(1.2), Inches(12.3))

    # ========================================================
    # Slide 8: Layer 1 ツリーマップ（病床日数シェア）
    # ========================================================
    s = prs.slides.add_slide(blank_layout)
    total_oku = agg['総粗利'].sum() / 1e8
    add_section_header(s,
        "Layer 1: 病床日数シェア — 各医師の占有割合",
        size=24)
    add_image_to_slide(s, CHART_DIR / f"04_treemap_{mode}.png",
                       Inches(0.7), Inches(1.2), Inches(11.9))
    add_textbox(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                f"※ 参考：管理用スコア合計 {total_oku:.2f} 億円相当（コスト ±20% の業界目安、確定利益ではない）",
                size=11, color=C_MUTED, align=PP_ALIGN.CENTER)

    # ========================================================
    # Slide 9: Layer 1 相関と注意（r=0.96 は循環論法）
    # ========================================================
    s = prs.slides.add_slide(blank_layout)
    add_section_header(s, "Layer 1: 病床日数シェアと管理用スコアの関係")
    add_image_to_slide(s, CHART_DIR / f"05_correlation_{mode}.png",
                       Inches(0.5), Inches(1.1), Inches(12.3))
    add_bg_rect(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.85), C_BG_RED)
    add_multi_para_textbox(s, Inches(0.7), Inches(6.6), Inches(12.0), Inches(0.7), [
        ("⚠️ r=0.96 は寄与率を使ってスコアを計算しているため統計的に独立ではない", 13, True),
        ("→「単価差より病床日数シェアがスコアを規定する」という説明として読む（独立検証ではない）", 12, False),
    ], color=C_DANGER)

    # ========================================================
    # Slide 10: Layer 2 概要（出来高未反映）
    # ========================================================
    s = prs.slides.add_slide(blank_layout)
    add_section_header(s, "Layer 2: 出来高上乗せ収益の可能性 — 概要", bg=C_BG_GREEN, size=24)
    body = [
        ("地域包括医療病棟入院料の包括範囲外として、収益に影響する要素:", 18, True),
        ("    手術・麻酔・リハビリテーション・1,000 点以上の処置・特定の検査", 16, False),
        ("    （出典: 厚生労働省「令和 6 年度診療報酬改定の概要 入院 I」）", 13, False),
        ("", 12, False),
        ("Layer 1 の病床日数だけでは、これらの出来高貢献を捕捉できない", 18, True),
        ("", 12, False),
        ("📌 本資料の Layer 2 は入口指標のみ", 18, True),
        ("    手術あり症例比率を提示。手術料・麻酔料・リハ単位の点数は本資料未反映", 14, False),
        ("    「単価が低い = 経営貢献が小さい」と判断する前に必ず併読する", 14, False),
        ("", 12, False),
        ("Stage 2 で点数データを取得し、絶対金額化を段階的に実施", 16, True),
    ]
    add_multi_para_textbox(s, Inches(0.7), Inches(1.3), Inches(12.0), Inches(5.8), body)

    # ========================================================
    # Slide 11: Layer 2 手術あり比率
    # ========================================================
    s = prs.slides.add_slide(blank_layout)
    add_section_header(s,
        f"Layer 2: 医師別 手術あり症例比率（分析対象 {len(agg)} 医師 {analysis_admissions:,} 件中の平均 {overall_surgery:.1f}%）",
        bg=C_BG_GREEN, size=18)
    add_image_to_slide(s, CHART_DIR / f"08_layer2_surgery_{mode}.png",
                       Inches(0.5), Inches(1.1), Inches(12.3))
    add_bg_rect(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.85), C_BG_GOLD)
    add_textbox(s, Inches(0.7), Inches(6.75), Inches(12.0), Inches(0.6),
                "💡 手術あり比率 ≥ 40% の医師は「手術・処置型」へ分類。Layer 1 では見えない出来高上乗せの可能性",
                size=14, bold=True, color=C_GOLD)

    # ========================================================
    # Slide 12: Layer 3 制度適合 3 指標
    # ========================================================
    s = prs.slides.add_slide(blank_layout)
    add_section_header(s, "Layer 3: 医師別 制度適合貢献（3 指標）", bg=C_BG_RED, size=24)
    add_image_to_slide(s, CHART_DIR / f"09_layer3_compliance_{mode}.png",
                       Inches(0.3), Inches(1.1), Inches(12.7))
    add_bg_rect(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.85), C_BG_GREEN)
    add_textbox(s, Inches(0.7), Inches(6.65), Inches(12.0), Inches(0.6),
                "💡 高単価でも施設基準悪化症例 と 単価中等でも制度維持に効く症例 を区別する",
                size=14, bold=True, color=C_OK)

    # ========================================================
    # Slide 13: Layer 3 表（基準対比）
    # ========================================================
    s = prs.slides.add_slide(blank_layout)
    add_section_header(s, "Layer 3: 当院全体平均と施設基準の対比", bg=C_BG_RED, size=24)
    body = [
        ("● 救急車入院比率（参考 proxy、医師別）:", 18, True),
        (f"    分析対象平均 {overall_emerg:.1f}%  ／  病棟別 rolling 制度線 {EMERGENCY_RATIO_THRESHOLD:.0f}%（参考）", 16, False),
        ("    ※ 制度上の「救急搬送後 15%」判定は病棟別 rolling 3 ヶ月。医師別比較とは別物", 13, False),
        ("", 8, False),
        ("● 在宅復帰率（自宅+居住系）:", 18, True),
        (f"    分析対象平均 {overall_home:.1f}%  ／  施設基準 {HOME_DISCHARGE_THRESHOLD:.0f}% 以上", 16, False),
        ("", 8, False),
        ("● 予定外入院比率（緊急入院全般）:", 18, True),
        (f"    分析対象平均 {overall_unp:.1f}%  ／  参考値（救急受入機能の厚みを示す）", 16, False),
        ("", 12, False),
        ("📌 病棟単位 rolling 3 ヶ月での制度判定は既存アプリの「制度管理」セクションで実施中", 14, True),
    ]
    add_multi_para_textbox(s, Inches(0.7), Inches(1.3), Inches(12.0), Inches(5.8), body)

    # ========================================================
    # Slide 14: Layer 4 5+1 類型定義
    # ========================================================
    s = prs.slides.add_slide(blank_layout)
    add_section_header(s, "Layer 4: 経営上の解釈 — 5+1 類型の定義", bg=C_BG_GOLD, size=24)
    body = [
        ("Layer 1〜3 を総合して、各医師を「病棟をどのような形で支えているか」で類型化:", 16, True),
        ("", 10, False),
        ("① 手術・処置型: 手術あり比率 ≥ 40%", 18, True),
        ("    → 在院は短〜中等でも、手術関連の出来高上乗せが大きい可能性", 13, False),
        ("② 救急・予定外受入基盤型: 救急車入院比率 ≥ 25% かつ 手術 < 30%", 18, True),
        ("    → 救急車入院・予定外入院を多く担当、地域包括病棟の受入機能を支える", 13, False),
        ("③ 長期安定型: 平均在院日数 ≥ 14 日", 18, True),
        ("    → 長めの在院日数で病床日数を支え、稼働率維持に効く", 13, False),
        ("④ リハ・在宅復帰型: 在宅復帰率 ≥ 95% かつ 平均在院 8〜13 日", 18, True),
        ("    → 適切な在院後に高い在宅復帰率、施設基準・退院支援に効く", 13, False),
        ("⑤ 回転供給型: 平均在院日数 ≤ 8 日", 18, True),
        ("    → 短期入院を多く回し、入退院の流れを作る", 13, False),
        ("⑥ 混合型: 上記いずれにも該当せず（複数の役割）", 18, True),
    ]
    add_multi_para_textbox(s, Inches(0.7), Inches(1.2), Inches(12.0), Inches(6.0), body)

    # ========================================================
    # Slide 15: Layer 4 散布図
    # ========================================================
    s = prs.slides.add_slide(blank_layout)
    add_section_header(s, "Layer 4: 当院の医師別 タイプ分類（散布図）", bg=C_BG_GOLD, size=24)
    add_image_to_slide(s, CHART_DIR / f"10_layer4_typing_{mode}.png",
                       Inches(0.5), Inches(1.1), Inches(12.3))

    # ========================================================
    # Slide 16: Layer 4 構成内訳
    # ========================================================
    s = prs.slides.add_slide(blank_layout)
    add_section_header(s, "Layer 4: 当院の医師構成内訳", bg=C_BG_GOLD, size=24)
    lines = [("当院 14 名の運営パターン構成:", 20, True), ("", 8, False)]
    for t in TYPE_ORDER:
        n = type_counts.get(t, 0)
        if n == 0:
            continue
        lines.append((f"● {t}: {n} 名", 20, True))
        lines.append((f"    {TYPE_DESCRIPTIONS[t]}", 13, False))
        lines.append(("", 6, False))
    lines.extend([
        ("", 10, False),
        ("💡 観察: 当院は手術・処置型と救急・予定外受入基盤型が中心、長期安定型がそれを支える構成", 16, True),
        ("        Layer 1 のスコアだけでは見えなかった貢献構造が浮かび上がる", 14, False),
    ])
    add_multi_para_textbox(s, Inches(0.7), Inches(1.2), Inches(12.0), Inches(6.0), lines)

    # ========================================================
    # Slide 17: 主な観察（問いかけ型）
    # ========================================================
    s = prs.slides.add_slide(blank_layout)
    add_section_header(s, "主な観察 — 経営会議で問うべきこと")
    body = [
        ("本資料は順位を見るためのものではありません。次の問いに使います。", 14, False),
        ("", 8, False),
        ("Q1. 当院は、どの類型の医師群で稼働を支えているか？（Layer 4）", 16, True),
        ("Q2. Layer 1 でスコア低い医師は、本当に経営貢献が小さいのか？", 16, True),
        ("    → Layer 2/3 を併読して、別軸の貢献を確認する", 13, False),
        ("Q3. 5F の単一医師集中（23.3%）は、脆弱性か強みか？", 16, True),
        ("    → 不在時の代替体制をどう組むかが構造課題", 13, False),
        ("Q4. 施設基準を支える症例構成の厚みはどこにあるか？（Layer 3）", 16, True),
        ("Q5. 手術・処置型 5 名の出来高貢献を、現資料は反映できているか？", 16, True),
        ("    → できていない（Stage 2 で点数データ取得して精緻化）", 13, False),
        ("Q6. 回転供給型・リハ・在宅復帰型が 0 名なのは何を意味するか？", 16, True),
        ("    → 「専門化が進んでいない」と断定はできない。閾値依存の heuristic で、", 13, False),
        ("       リハ単位データが CSV 未収録のため判断不能。追加データで再評価が必要", 13, False),
    ]
    add_multi_para_textbox(s, Inches(0.7), Inches(1.2), Inches(12.0), Inches(6.0), body)

    # ========================================================
    # Slide 18: 提言 1 — 5F 構造リスク
    # ========================================================
    s = prs.slides.add_slide(blank_layout)
    add_section_header(s, "提言 1: 病棟運営構造の継続監視（5F 構造）", bg=C_BG_RED, size=24)
    body = [
        (f"課題: 5F は寄与率トップの医師（{okuk_label}、救急・予定外受入基盤型）に病床日数の 23.3% が集中", 18, True),
        ("", 12, False),
        ("→ 個人攻撃ではなく構造リスクとして次の対応を提言:", 16, True),
        ("    ● 5F に主たる病棟を持つ医師の中期的な増員を検討", 16, False),
        (f"    ● {okuk_label} の患者バックアップ体制（指示・引継ぎ）の整備", 16, False),
        ("    ● 5F 主治医の月次寄与率を可視化し、依存度を継続監視", 16, False),
        ("", 12, False),
        ("● 6F 病棟は寄与率上位 3 名で分散しており、依存リスクは相対的に小さい", 14, False),
    ]
    add_multi_para_textbox(s, Inches(0.7), Inches(1.5), Inches(12.0), Inches(5.5), body)

    # ========================================================
    # Slide 19: 提言 2 — 制度適合維持
    # ========================================================
    s = prs.slides.add_slide(blank_layout)
    add_section_header(s, "提言 2: 制度適合の維持（Layer 3 の月次監視）", bg=C_BG_RED, size=24)
    body = [
        ("背景: 2026-06-01 から救急搬送後 15% は本則完全適用（rolling 3 ヶ月）", 18, True),
        ("       在宅復帰 70% は地域包括医療病棟の継続要件", 16, False),
        ("", 12, False),
        ("対応:", 18, True),
        ("    ● 病棟別 Layer 3 指標の月次ダッシュボード化（既存アプリ組込み済）", 16, False),
        ("    ● 施設基準下振れ徴候があれば、救急受入比率と在宅復帰調整の介入を検討", 16, False),
        ("    ● ALOS の階段関数（短手 3 Day 5/6 境界）の遵守をベッドコントロールで継続", 16, False),
        ("", 12, False),
        ("既存アプリ「📊 今日の運営」「🛡️ 制度管理」セクションで毎日監視可能", 14, True),
    ]
    add_multi_para_textbox(s, Inches(0.7), Inches(1.5), Inches(12.0), Inches(5.5), body)

    # ========================================================
    # Slide 20: 提言 3 — Stage 2 データ整備
    # ========================================================
    s = prs.slides.add_slide(blank_layout)
    add_section_header(s, "提言 3: 必要データの整備（Stage 2）", bg=C_BG_BLUE, size=24)
    body = [
        ("本資料は CSV から計算可能な範囲。事務は患者別コスト計算ができない前提のため、", 16, False),
        ("**事務に重い負担を強いる順序ではなく、請求収入側の集計可能性確認を先行**する。", 16, True),
        ("", 8, False),
        ("【優先度 1】請求収入側の月次集計可能性確認", 18, True),
        ("    ● 手術料・麻酔料の患者別 月次集計が可能か確認", 14, False),
        ("    ● リハビリ単位数（PT / OT / ST）と リハビリ料の月次集計可能性", 14, False),
        ("    ● 処置加算（中心静脈、人工呼吸器、ドレーン管理）と検査加算", 14, False),
        ("    → 既存の請求データ・レセプトの集計範囲で見える出来高情報を整理", 13, False),
        ("", 8, False),
        ("【優先度 2】コスト見積もりの精緻化（社内分析）", 18, True),
        ("    ● 優先度 1 が見えてから、当院 DPC データを活用したコスト精緻化を検討", 14, False),
        ("", 8, False),
        ("📌 事務に対する精密コスト計算依頼ではなく、社内分析として段階的に整備", 14, True),
    ]
    add_multi_para_textbox(s, Inches(0.7), Inches(1.3), Inches(12.0), Inches(5.8), body)

    # ========================================================
    # Slide 21: 提言 4 — KPI 月次運用（4 層併読）
    # ========================================================
    s = prs.slides.add_slide(blank_layout)
    add_section_header(s, "提言 4: KPI 月次運用 — 4 層併読の文化", bg=C_BG_BLUE, size=24)
    body = [
        ("「稼働率寄与率（Layer 1）」だけ見ると、Layer 2/3 の貢献が抜け落ちる", 18, True),
        ("", 12, False),
        ("運用提言:", 18, True),
        ("    ● 月次会議で Layer 1〜3 の指標を併記して提示", 16, False),
        ("    ● 順位ではなく Layer 4 散布図で「タイプ構成」として議論", 16, False),
        ("    ● 既存ベッドコントロールアプリに月次自動更新機能を組込み（次フェーズ）", 16, False),
        ("    ● 四半期に 1 回、Stage 2 データ取込み状況をレビュー", 16, False),
        ("", 12, False),
        ("⚠️ 倫理的配慮:", 16, True),
        ("    患者選別・適応外入院・在院延長の誘因にしてはならない", 13, False),
        ("    「適応のある患者を、適切な期間で治療し、自宅復帰を支援する」が大前提", 13, False),
    ]
    add_multi_para_textbox(s, Inches(0.7), Inches(1.4), Inches(12.0), Inches(5.7), body)

    # ========================================================
    # Slide 22: ご質問
    # ========================================================
    s = prs.slides.add_slide(blank_layout)
    add_bg_rect(s, 0, 0, SW, SH, RGBColor(0xF9, 0xFA, 0xFB))
    add_textbox(s, Inches(0.5), Inches(3.0), Inches(12.3), Inches(1.5),
                "ご質問・ご討議",
                size=60, bold=True, color=C_TITLE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.7),
                "おもろまちメディカルセンター",
                size=18, color=C_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.7),
                "副院長 久保田 透",
                size=14, color=C_MUTED, align=PP_ALIGN.CENTER)

    suffix = '理事会プレゼン' if mode == 'named' else '運営会議プレゼン'
    out_path = REPO_ROOT / 'docs' / 'admin' / f'医師別KPI解析_{suffix}_2026-05-04.pptx'
    prs.save(out_path)
    print(f"✅ Generated: {out_path} ({len(prs.slides)} slides)")
    return out_path


if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('--mode', choices=['named', 'anonymous'], default='named')
    args = parser.parse_args()
    build(args.mode)
