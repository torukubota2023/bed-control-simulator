"""ブリーフィング PPTX に「6 階病棟の医療提供の工夫」スライドを挿入.

副院長要望 (2026-04-26):
    新基準（2026-06-01 適用）では 6F は「漏れなく拾う」だけでは届かない。
    医学的適応の範囲内で、複数の選択肢の中から A6/C 項目に貢献する選び方を
    取ることで、看護必要度の構造に医療提供を添わせる。これを独立スライド
    として既存 PPTX に追加する（γ: 行動主導）。

    倫理ガードレール:
        - 適応外の処置は絶対に行わない
        - 医学的判断は変わらない
        - 複数の選択肢がある中での「選び方」の問題として捉える

実行:
    .venv/bin/python scripts/insert_achievability_slide.py

入力:
    docs/admin/2026新基準_医師看護師向けブリーフィング_2026-04-26.pptx (既存、副院長編集済 13 スライド)

出力:
    同ファイルを上書き、末尾に 1 スライド追加 (合計 14 スライド)
    その後 PowerPoint で 「100 患者日まとめ」(現スライド 13) の直前に
    手動でドラッグ移動してください (推奨位置: 新スライド 13、まとめが 14 へ)。
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PPTX_PATH = Path(__file__).resolve().parent.parent / "docs" / "admin" / "2026新基準_医師看護師向けブリーフィング_2026-04-26.pptx"

# 色トークン（既存スライドと統一）
COLOR_TEXT = RGBColor(0x1F, 0x29, 0x37)         # ダークグレー
COLOR_TITLE = RGBColor(0x37, 0x41, 0x51)        # アクセント
COLOR_SUCCESS = RGBColor(0x10, 0xB9, 0x81)      # 緑（行動の方向）
COLOR_CAPTION = RGBColor(0x6B, 0x72, 0x80)      # キャプション
COLOR_TABLE_HEAD = RGBColor(0x37, 0x41, 0x51)
COLOR_TABLE_HEAD_BG = RGBColor(0xF3, 0xF4, 0xF6)


def _add_text(text_frame, text: str, size: int, bold: bool = False,
              color: RGBColor = COLOR_TEXT, alignment: int = PP_ALIGN.LEFT,
              new_para: bool = True) -> None:
    if new_para:
        para = text_frame.add_paragraph()
    else:
        para = text_frame.paragraphs[0]
    para.alignment = alignment
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Hiragino Kaku Gothic ProN"


def _build_achievability_slide(prs: Presentation) -> None:
    """6F 医療提供の工夫スライドを構築（γ: 行動主導版）."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # タイトル
    title_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.9)
    )
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    _add_text(
        title_tf,
        "新基準が求める内科型の地域包括医療病棟である 6 階病棟の医療提供の工夫",
        size=22, bold=True, color=COLOR_TITLE, new_para=False,
    )

    # キーメッセージ（γ: 行動主導 ― 「記録する」だけでなく「選ぶ」）
    key_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.7)
    )
    key_tf = key_box.text_frame
    key_tf.word_wrap = True
    _add_text(
        key_tf,
        "「記録する」だけでなく「選ぶ」 ― 医学的適応の中で A6 / C 項目に貢献する選択肢を取る",
        size=20, bold=True, color=COLOR_SUCCESS,
        alignment=PP_ALIGN.CENTER, new_para=False,
    )

    # 倫理ガードレール（最上段に明示）
    guard_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(1.85), Inches(12.5), Inches(0.4)
    )
    guard_tf = guard_box.text_frame
    guard_tf.word_wrap = True
    _add_text(
        guard_tf,
        "前提: 適応外の処置は絶対に行わない / 医学的判断は変わらない / "
        "複数の選択肢がある中での「選び方」の問題として捉える",
        size=11, color=COLOR_CAPTION,
        alignment=PP_ALIGN.CENTER, new_para=False,
    )

    # テーブル：7 行（項目 × 現状 × 工夫の方向 × A/C 寄与）
    rows, cols = 8, 4  # ヘッダー + 7 データ
    table_left = Inches(0.4)
    table_top = Inches(2.35)
    table_width = Inches(12.5)
    table_height = Inches(4.3)
    table = slide.shapes.add_table(rows, cols, table_left, table_top,
                                    table_width, table_height).table

    # 列幅
    table.columns[0].width = Inches(2.2)
    table.columns[1].width = Inches(3.6)
    table.columns[2].width = Inches(4.4)
    table.columns[3].width = Inches(2.3)

    headers = [
        "項目",
        "現状の医療提供",
        "工夫の方向（適応の範囲内）",
        "A / C 項目寄与",
    ]
    data = [
        ("A6③ 麻薬注射",
         "がん性疼痛で経口移行を急ぐ",
         "適応継続中は持続点滴を維持(疼痛コントロール優先)",
         "A 3 点 / 3-5 日"),
        ("A6⑦ 昇圧剤",
         "敗血症で短時間使用、早期離脱",
         "必要十分な期間まで維持(早すぎる中止を避ける)",
         "A 3 点 / 1-3 日"),
        ("A6⑧ 抗不整脈剤",
         "Af で経口管理を優先",
         "急性期は IV アミオダロン / ペイン PHN は IV リドカイン",
         "A 3 点 / 1-7 日"),
        ("A6⑨ 抗血栓持続",
         "PE で DOAC 移行を急ぐ",
         "高リスク・腎不全例は UFH 持続を維持",
         "A 3 点 / 3-7 日"),
        ("A4 + A3 ペア",
         "β-ラクタム間欠投与",
         "重症感染症で PIPC/TAZ・MEPM 持続点滴を選択",
         "A 2 点 / 5-7 日"),
        ("C21③ 内視鏡治療",
         "適応症例で見送り傾向",
         "ESD・ERCP・内視鏡止血を積極化",
         "C 項目 / 各 4 日"),
        ("C23 PEG・CV 挿入",
         "末梢困難でも見送り",
         "適応症例で閾値を下げる",
         "C 項目 / PEG 5 日・CV 4 日"),
    ]

    # ヘッダー行
    for col_i, header in enumerate(headers):
        cell = table.cell(0, col_i)
        cell.text = ""
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = header
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = COLOR_TABLE_HEAD
        run.font.name = "Hiragino Kaku Gothic ProN"
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_TABLE_HEAD_BG

    # データ行
    for row_i, row_data in enumerate(data, start=1):
        for col_i, value in enumerate(row_data):
            cell = table.cell(row_i, col_i)
            cell.text = ""
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if col_i in (1, 2) else PP_ALIGN.CENTER
            run = para.add_run()
            run.text = value
            run.font.size = Pt(11)
            run.font.color.rgb = COLOR_TEXT
            run.font.name = "Hiragino Kaku Gothic ProN"
            # 工夫の方向 列は緑強調（行動の方向）
            if col_i == 2:
                run.font.color.rgb = COLOR_SUCCESS
            # A/C 寄与列は太字でアクセント色
            if col_i == 3:
                run.font.bold = True
                run.font.color.rgb = COLOR_TITLE

    # 行動メッセージ（フッター、γ 思想を凝縮）
    action_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.6)
    )
    action_tf = action_box.text_frame
    action_tf.word_wrap = True
    _add_text(
        action_tf,
        "📌 達成は「漏れを拾う」ではなく「選択肢の取り方を変える」 ― "
        "医学的適応の範囲で 6F の医療提供を看護必要度の構造に添わせる",
        size=13, bold=True, color=COLOR_SUCCESS,
        alignment=PP_ALIGN.CENTER, new_para=False,
    )

    # スピーカーノート: 安全にスキップ（python-pptx でノート挿入が
    # 既存 PPTX のレイアウト次第で不安定。PowerPoint で副院長が手動追加可能）


def _remove_existing_achievability_slide(prs: Presentation) -> int:
    """既存の挿入済みスライドがあれば削除 + 孤立 slide rel を掃除.

    旧版（"100 患者日の達成は…"）と新版（"…6 階病棟の医療提供の工夫"）の
    両方を検出する。

    過去の不完全な remove 実装が残した「sldIdLst には載っていないが
    presentation.xml.rels には残っている」孤立 rel もここで一掃する。
    これがないと ZIP に slide{N}.xml が二重で書き込まれ、PowerPoint が
    古い方を読み続ける（python-pptx は新しい方を読むため検出できない）。

    Returns:
        削除した枚数（sldIdLst から除いた slide 数）
    """
    from pptx.oxml.ns import qn

    target_keywords = (
        "6 階病棟の医療提供の工夫",  # 新版（γ）
        "100 患者日の達成は",          # 旧版（reality check）
    )
    xml_slides = prs.slides._sldIdLst
    sldId_elements = list(xml_slides)
    indices_to_remove: list[int] = []
    for idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            if any(kw in text for kw in target_keywords):
                indices_to_remove.append(idx)
                break

    # 後ろから削除（インデックスずれ防止）
    for idx in reversed(indices_to_remove):
        sld_id_elem = sldId_elements[idx]
        rId = sld_id_elem.get(qn("r:id"))
        prs.part.drop_rel(rId)
        xml_slides.remove(sld_id_elem)

    # 孤立 rel 掃除: sldIdLst に載っていない slide part への rel を全削除
    valid_rIds = {e.get(qn("r:id")) for e in xml_slides}
    pres_rels = prs.part.rels
    orphan_rIds: list[str] = []
    for rId, rel in list(pres_rels.items()):
        if rel.is_external:
            continue
        partname = str(rel.target_part.partname)
        if partname.startswith("/ppt/slides/slide") and rId not in valid_rIds:
            orphan_rIds.append(rId)
    for rId in orphan_rIds:
        prs.part.drop_rel(rId)

    return len(indices_to_remove)


def main() -> int:
    if not PPTX_PATH.exists():
        print(f"❌ PPTX が見つかりません: {PPTX_PATH}")
        return 1

    prs = Presentation(str(PPTX_PATH))
    n_before_total = len(prs.slides)
    n_removed = _remove_existing_achievability_slide(prs)
    n_before = len(prs.slides)
    if n_removed > 0:
        print(f"♻️ 既存の挿入済みスライド {n_removed} 枚を削除（重複防止）")
    _build_achievability_slide(prs)
    prs.save(str(PPTX_PATH))
    n_after = len(prs.slides)

    print(f"✅ スライド追加完了: {n_before} → {n_after} スライド")
    print(f"   追加位置: 末尾（PowerPoint スライド {n_after}）")
    print(f"   ファイル: {PPTX_PATH}")
    print()
    print("📋 次の手順（PowerPoint で副院長が手動操作）:")
    print(f"   1. PowerPoint で {PPTX_PATH.name} を開く")
    print(f"   2. 左ペインのスライド一覧で、新規追加されたスライド {n_after} を")
    print(f"      『100 患者日まとめ』（現スライド 13）の直前にドラッグ移動")
    print(f"   3. 結果: 6 階病棟の医療提供の工夫 → 100 患者日まとめ の流れになる")
    print(f"   4. 保存して完了")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
